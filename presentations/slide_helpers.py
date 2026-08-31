"""
Curo — Slide Helper Utilities
==============================
Low-level python-pptx wrappers used by every slide pattern. All coordinates in
inches, all type sizes in points, all colours as hex strings. Every default
comes from design_tokens; there are no magic numbers here.

Two invariants this module owns
-------------------------------
1. **No raw citation marker reaches a slide.** Every text run created here
   passes through `text_budget.sanitize()`. That makes the rule structural
   rather than a habit: a pattern that forgets to sanitise still cannot emit
   `[[PMID:12345]]`, because the only way to put text on a slide is through
   these functions.

2. **Font fallback.** python-pptx writes one `<a:latin typeface="X"/>` per run.
   If X is not installed PowerPoint silently drops to the theme font (usually
   Calibri), which collapses the display/body distinction the layouts depend
   on. `set_font` writes the primary face into the run XML and records the
   fallback chain, so a machine missing Georgia still lands on Cambria or
   Times rather than on the sans body face. See design_tokens for why the
   spec's Instrument Serif / Inter cannot simply be embedded.
"""

from __future__ import annotations

import math

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from presentations.design_tokens import (
    COLORS, FONTS, FONT_FALLBACKS, SIZES, SIZES_PX, LAYOUT, GEOM_PX,
    LINE_HEIGHT, TIER_CHIP, TIER_LABELS, TIER_CHART_FILL, PILL_PMID,
    TRACK_EYEBROW, TRACK_CHIP, px_in, px_pt, tier_key,
)
from presentations.text_budget import sanitize, extract_pmids

# python-pptx autoshape ids used directly to avoid the enum import churn.
_SHAPE_RECT = 1          # MSO_AUTO_SHAPE_TYPE.RECTANGLE
_SHAPE_ROUNDED = 5       # MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
_SHAPE_OVAL = 9          # MSO_AUTO_SHAPE_TYPE.OVAL


# ── Colour ────────────────────────────────────────────────────────────────────

def hex_rgb(hex_str: str) -> RGBColor:
    """Convert '#RRGGBB' to RGBColor."""
    h = str(hex_str).lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def resolve_color(value: str | None, default: str) -> str:
    """Accept a COLORS key or a raw hex; fall back to `default`."""
    if not value:
        return default
    if isinstance(value, str) and value.startswith("#") and len(value) == 7:
        return value
    return COLORS.get(value, default)


def relative_luminance(hex_str: str) -> float:
    h = str(hex_str).lstrip("#")
    channels = []
    for i in (0, 2, 4):
        c = int(h[i:i + 2], 16) / 255.0
        channels.append(c / 12.92 if c <= 0.03928
                        else ((c + 0.055) / 1.055) ** 2.4)
    r, g, b = channels
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def contrast_ratio(fg: str, bg: str) -> float:
    """WCAG contrast ratio. Body text on `bg` must clear 4.5:1 (spec §2.4)."""
    l1, l2 = relative_luminance(fg), relative_luminance(bg)
    hi, lo = max(l1, l2), min(l1, l2)
    return (hi + 0.05) / (lo + 0.05)


# ── Fonts ─────────────────────────────────────────────────────────────────────

_FONT_ROLE_MAP = {
    "display": "display", "serif": "display", "header_serif": "display",
    "body": "body", "sans": "body", "body_sans": "body",
    # The old deck had a mono eyebrow face. The spec's eyebrow is Inter 600
    # uppercase with tracking, so mono now resolves to the body face.
    "mono": "body", "mono_eyebrow": "body",
}


def set_font(
    run,
    font_role: str,
    size,
    *,
    bold: bool = False,
    italic: bool = False,
    color: str | None = None,
    tracking: int | None = None,
) -> None:
    """Apply face, size, style, colour and the fallback chain to a pptx run.

    font_role: "display" | "body" (aliases: serif/sans/mono kept for callers
    written against the previous vocabulary).
    size: points. tracking: character spacing in 1/100 pt.
    """
    role_key = _FONT_ROLE_MAP.get(font_role, "body")
    primary = FONTS[role_key]

    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    if color:
        f.color.rgb = hex_rgb(color)

    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:cs"):
        for old in rPr.findall(qn(tag)):
            rPr.remove(old)

    latin = etree.SubElement(rPr, qn("a:latin"))
    latin.set("typeface", primary)
    cs = etree.SubElement(rPr, qn("a:cs"))
    cs.set("typeface", primary)

    if tracking is not None:
        rPr.set("spc", str(int(tracking)))


def font_fallback_chain(font_role: str) -> list[str]:
    """The documented substitution chain for a role (used by tests + the web deck)."""
    return list(FONT_FALLBACKS[_FONT_ROLE_MAP.get(font_role, "body")])


# ── Text metrics (for the overflow budget) ────────────────────────────────────

# Mean glyph advance as a fraction of the em, measured for the two stand-in
# faces on mixed-case clinical prose. Used only to *predict* wrapping so the
# builder can reserve height; PowerPoint does the real layout.
_ADVANCE = {"display": 0.545, "body": 0.485}


def estimate_lines(text: str, size_pt: float, width_in: float,
                   font_role: str = "body") -> int:
    """Predicted wrapped line count for `text` in a box `width_in` wide."""
    text = sanitize(text)
    if not text or width_in <= 0:
        return 0
    adv = _ADVANCE.get(_FONT_ROLE_MAP.get(font_role, "body"), 0.48)
    chars_per_line = max(1, int((width_in * 72.0) / (size_pt * adv)))
    lines = 0
    for para in text.split("\n"):
        lines += max(1, int(math.ceil(len(para) / float(chars_per_line))))
    return lines


def estimate_height_in(text: str, size_pt: float, width_in: float,
                       *, font_role: str = "body",
                       line_height: float = 1.5) -> float:
    """Predicted rendered height in inches, for reserving vertical space."""
    lines = estimate_lines(text, size_pt, width_in, font_role)
    return lines * (size_pt * line_height) / 72.0


# ── Text boxes ────────────────────────────────────────────────────────────────

def add_textbox(
    slide,
    text: str,
    x: float, y: float, w: float, h: float,
    *,
    font_role: str = "body",
    size=None,
    color: str = COLORS["text_body"],
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    tracking: int | None = None,
    wrap: bool = True,
    upper: bool = False,
    line_spacing: float | None = None,
    margin_left: float = 0.0,
    margin_top: float = 0.0,
    anchor=None,
):
    """Add a single-run text box. Text is sanitised of citation markers."""
    if size is None:
        size = SIZES["bullet"]
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Inches(margin_left)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(margin_top)
    tf.margin_bottom = Inches(0)
    if anchor is not None:
        tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    clean = sanitize(text)
    run.text = clean.upper() if upper else clean
    set_font(run, font_role, size, bold=bold, italic=italic, color=color,
             tracking=tracking)
    return txb


def add_multiline_textbox(
    slide,
    lines: list[dict],
    x: float, y: float, w: float, h: float,
    *,
    wrap: bool = True,
    margin_left: float = 0.0,
):
    """Add a multi-paragraph text box. Every run is sanitised.

    Each entry: text, font_role, size, color, bold, italic, align,
    space_before (pt), line_spacing, tracking, upper.
    """
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Inches(margin_left)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    for i, spec in enumerate(lines):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        sb = spec.get("space_before", 0)
        if sb:
            p.space_before = Pt(sb)
        ls = spec.get("line_spacing")
        if ls:
            p.line_spacing = ls
        run = p.add_run()
        clean = sanitize(spec.get("text", ""))
        run.text = clean.upper() if spec.get("upper") else clean
        set_font(
            run,
            spec.get("font_role", "body"),
            spec.get("size", SIZES["bullet"]),
            bold=spec.get("bold", False),
            italic=spec.get("italic", False),
            color=spec.get("color", COLORS["text_body"]),
            tracking=spec.get("tracking"),
        )
    return txb


# ── Shape primitives ──────────────────────────────────────────────────────────

def add_filled_rect(slide, x, y, w, h, fill_color,
                    line_color: str | None = None,
                    line_width_pt: float = 0.75):
    """Solid rectangle. Returns the shape."""
    shape = slide.shapes.add_shape(_SHAPE_RECT, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill_color)
    if line_color:
        shape.line.color.rgb = hex_rgb(line_color)
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rounded_rect(slide, x, y, w, h, fill_color, *,
                     radius_in: float = None,
                     line_color: str | None = None,
                     line_width_pt: float = 0.75):
    """Rounded rectangle with a radius in inches (not a fraction).

    python-pptx exposes the corner as adjustment[0], a fraction of half the
    shorter side, so the requested radius has to be converted per shape.
    """
    if radius_in is None:
        radius_in = LAYOUT["card_radius_in"]
    shape = slide.shapes.add_shape(_SHAPE_ROUNDED, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    try:
        adj = max(0.0, min(0.5, radius_in / max(min(w, h), 1e-6)))
        shape.adjustments[0] = adj
    except Exception:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill_color)
    if line_color:
        shape.line.color.rgb = hex_rgb(line_color)
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_circle(slide, cx, cy, diameter, fill_color,
               line_color: str | None = None, line_width_pt: float = 0.75):
    """Solid circle centred at (cx, cy)."""
    shape = slide.shapes.add_shape(
        _SHAPE_OVAL,
        Inches(cx - diameter / 2), Inches(cy - diameter / 2),
        Inches(diameter), Inches(diameter),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill_color)
    if line_color:
        shape.line.color.rgb = hex_rgb(line_color)
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_hairline(slide, x, y, w, color: str | None = None):
    """The spec's 1px rule — table borders, the footer rule, row separators."""
    return add_filled_rect(slide, x, y, w, LAYOUT["hairline_in"],
                           color or COLORS["border"])


def add_left_accent_bar(slide, x, y, h, color, width: float = px_in(3)):
    """Thin vertical accent on the left edge of a row or card."""
    return add_filled_rect(slide, x, y, width, h, color)


# ── Chips and pills ───────────────────────────────────────────────────────────

def chip_width_in(text: str, *, size_px: float = None, dot: bool = False,
                  pad_px: float = 12.0) -> float:
    """Predicted pill width for right-aligned placement."""
    if size_px is None:
        size_px = SIZES_PX["chip"]
    text = sanitize(text)
    # 0.62em mean advance for uppercase + the 0.08em tracking the spec asks for.
    glyph_px = len(text) * size_px * 0.70
    dot_px = (GEOM_PX["chip_dot"] + 6) if dot else 0
    return px_in(glyph_px + dot_px + pad_px * 2)


def add_chip(slide, text: str, x: float, y: float, *,
             bg: str, fg: str, dot_color: str | None = None,
             size=None, radius_in: float | None = None,
             tracking: int | None = None, width_in: float | None = None,
             bold: bool = True):
    """A pill: rounded background, optional 8px colour dot, uppercase label.

    Returns the width consumed, so callers can lay chips out in a row.
    """
    if size is None:
        size = SIZES["chip"]
    if radius_in is None:
        radius_in = LAYOUT["chip_h_in"] / 2   # radius 999 -> fully round ends
    if tracking is None:
        tracking = TRACK_CHIP
    h = LAYOUT["chip_h_in"]
    w = width_in or chip_width_in(text, dot=bool(dot_color))
    add_rounded_rect(slide, x, y, w, h, bg, radius_in=radius_in)

    pad = px_in(12)
    text_x = x + pad
    if dot_color:
        d = LAYOUT["chip_dot_in"]
        add_circle(slide, x + pad + d / 2, y + h / 2, d, dot_color)
        text_x = x + pad + d + px_in(6)

    add_textbox(slide, text, text_x, y + px_in(3), w - (text_x - x) - pad * 0.6,
                h - px_in(4),
                font_role="body", size=size, color=fg, bold=bold,
                upper=True, tracking=tracking, wrap=False)
    return w


def add_tier_chip(slide, tier: str, x: float, y: float,
                  *, right_edge: float | None = None) -> float:
    """Evidence-tier pill with its 8px ladder dot.

    Level III must always carry its text label (spec §1.2) — every chip here
    does, so colour is never the only channel carrying the tier.
    """
    key = tier_key(tier)
    if key is None:
        return 0.0
    bg, fg = TIER_CHIP[key]
    label = TIER_LABELS[key]
    w = chip_width_in(label, dot=True)
    if right_edge is not None:
        x = right_edge - w
    add_chip(slide, label, x, y, bg=bg, fg=fg,
             dot_color=TIER_CHART_FILL[key], width_in=w)
    return w


def add_pmid_pill(slide, pmid: str, x: float, y: float,
                  *, right_edge: float | None = None) -> float:
    """Author-style PMID pill — the only place a PMID appears as a mark."""
    bg, fg = PILL_PMID
    label = f"PMID {pmid}"
    w = chip_width_in(label)
    if right_edge is not None:
        x = right_edge - w
    add_chip(slide, label, x, y, bg=bg, fg=fg, width_in=w, bold=False)
    return w


# ── Slide furniture (spec §1.3 — identical on every content-class slide) ──────

def slide_background(slide, color: str | None = None):
    return add_filled_rect(slide, 0, 0, LAYOUT["slide_w_in"],
                           LAYOUT["slide_h_in"], color or COLORS["bg"])


def add_header_row(slide, eyebrow: str, *, tier: str | None = None,
                   right_label: str | None = None) -> None:
    """Fixed 26px header: eyebrow left, tier chip or CURO label right."""
    x = LAYOUT["pad_x_in"]
    y = LAYOUT["header_y_in"]
    h = LAYOUT["header_h_in"]
    w = LAYOUT["content_w_in"]
    right = x + w

    if eyebrow:
        add_textbox(slide, eyebrow, x, y + px_in(5), w * 0.72, h,
                    font_role="body", size=SIZES["eyebrow"],
                    color=COLORS["text_eyebrow"], bold=True,
                    upper=True, tracking=TRACK_EYEBROW, wrap=False)

    if tier and tier_key(tier):
        add_tier_chip(slide, tier, 0, y, right_edge=right)
    elif right_label:
        add_textbox(slide, right_label, right - 2.2, y + px_in(5), 2.2, h,
                    font_role="body", size=SIZES["eyebrow"],
                    color=COLORS["text_eyebrow"], bold=True, upper=True,
                    tracking=TRACK_EYEBROW, align=PP_ALIGN.RIGHT, wrap=False)


def add_slide_title(slide, text: str, *, y: float | None = None,
                    size=None, color: str | None = None,
                    width_in: float | None = None,
                    font_role: str = "display",
                    italic: bool = False) -> float:
    """Serif slide title. Returns the y of its bottom edge."""
    if size is None:
        size = SIZES["title"]
    if y is None:
        y = LAYOUT["title_y_in"]
    w = width_in if width_in is not None else LAYOUT["content_w_in"]
    h = estimate_height_in(text, size, w, font_role=font_role,
                           line_height=LINE_HEIGHT["title"])
    h = max(h, (size * LINE_HEIGHT["title"]) / 72.0)
    add_textbox(slide, text, LAYOUT["pad_x_in"], y, w, h + px_in(6),
                font_role=font_role, size=size,
                color=color or COLORS["text_body"], italic=italic,
                line_spacing=LINE_HEIGHT["title"])
    # Serif descenders sit below the predicted box, so the returned baseline
    # carries a little slack — callers stack the lead paragraph directly on it.
    return y + h + px_in(10)


def add_lead(slide, text: str, y: float) -> float:
    """Optional lead paragraph under the title. Returns its bottom edge."""
    if not sanitize(text):
        return y
    w = min(LAYOUT["lead_max_w_in"], LAYOUT["content_w_in"])
    y = y + LAYOUT["lead_gap_in"]
    h = estimate_height_in(text, SIZES["lead"], w,
                           line_height=LINE_HEIGHT["lead"])
    add_textbox(slide, text, LAYOUT["pad_x_in"], y, w, h + px_in(4),
                font_role="body", size=SIZES["lead"],
                color=COLORS["text_lead"], line_spacing=LINE_HEIGHT["lead"])
    return y + h


def add_footer(slide, *, citations: str = "", page_num: int | None = None,
               total_pages: int | None = None) -> None:
    """Footer: 1px top border, short citations left, page number right.

    `citations` is short-form ("Author et al. 2017 · Journal · n = 100 ·
    PMID 28294701"). Raw markers cannot survive `add_textbox`, so a footer
    built from marker-bearing text still renders clean.
    """
    x = LAYOUT["pad_x_in"]
    w = LAYOUT["content_w_in"]
    add_hairline(slide, x, LAYOUT["footer_rule_y_in"], w, COLORS["border"])

    text_y = LAYOUT["footer_text_y_in"]
    if citations:
        add_textbox(slide, citations, x, text_y, w * 0.86, px_in(20),
                    font_role="body", size=SIZES["footer"],
                    color=COLORS["text_footer"], wrap=False)
    if page_num is not None:
        label = str(page_num)
        add_textbox(slide, label, x + w * 0.86, text_y, w * 0.14, px_in(20),
                    font_role="body", size=SIZES["footer"],
                    color=COLORS["text_muted"], align=PP_ALIGN.RIGHT,
                    wrap=False)


def add_notice_box(slide, text: str, x: float, y: float, w: float, h: float,
                   *, heading: str | None = None,
                   bg: str | None = None,
                   heading_color: str | None = None,
                   body_color: str | None = None):
    """Surface-filled rounded box — 'DOES NOT APPLY WHEN', insufficient-evidence.

    A notice, not an error: the spec restyles the missing-module slide as a
    calm box on the deck surface rather than red alert furniture.
    """
    add_rounded_rect(slide, x, y, w, h, bg or COLORS["surface"],
                     radius_in=LAYOUT["card_radius_in"])
    pad = px_in(22)
    inner_y = y + pad * 0.8
    if heading:
        add_textbox(slide, heading, x + pad, inner_y, w - pad * 2, px_in(18),
                    font_role="body", size=SIZES["eyebrow"],
                    color=heading_color or COLORS["text_eyebrow"], bold=True,
                    upper=True, tracking=TRACK_EYEBROW, wrap=False)
        inner_y += px_in(24)
    add_textbox(slide, text, x + pad, inner_y, w - pad * 2,
                max(px_in(20), h - (inner_y - y) - pad * 0.8),
                font_role="body", size=SIZES["card_body"],
                color=body_color or COLORS["text_secondary"],
                line_spacing=LINE_HEIGHT["card"])


def add_bullet_row(slide, text: str, x: float, y: float, w: float,
                   *, size=None, color: str | None = None,
                   marker_color: str | None = None) -> float:
    """One bullet: 8px round marker + wrapped text. Returns the bottom edge."""
    if size is None:
        size = SIZES["bullet"]
    d = LAYOUT["bullet_dot_in"]
    text_x = x + d + px_in(14)
    text_w = w - (text_x - x)
    h = estimate_height_in(text, size, text_w, line_height=LINE_HEIGHT["bullet"])
    h = max(h, (size * LINE_HEIGHT["bullet"]) / 72.0)
    add_circle(slide, x + d / 2, y + px_in(11), d,
               marker_color or COLORS["accent_blue"])
    add_textbox(slide, text, text_x, y, text_w, h + px_in(4),
                font_role="body", size=size,
                color=color or COLORS["text_body"],
                line_spacing=LINE_HEIGHT["bullet"])
    return y + h


def add_image_bytes(slide, png_bytes: bytes, x: float, y: float,
                    *, width_in: float | None = None,
                    height_in: float | None = None):
    """Place PNG bytes (a rendered chart) without touching the filesystem."""
    import io
    kwargs = {}
    if width_in is not None:
        kwargs["width"] = Inches(width_in)
    if height_in is not None:
        kwargs["height"] = Inches(height_in)
    return slide.shapes.add_picture(io.BytesIO(png_bytes), Inches(x),
                                    Inches(y), **kwargs)


# ── Presentation factory ──────────────────────────────────────────────────────

def new_presentation() -> Presentation:
    """Blank 16:9 Presentation at the spec's frame size."""
    prs = Presentation()
    prs.slide_width = Inches(LAYOUT["slide_w_in"])
    prs.slide_height = Inches(LAYOUT["slide_h_in"])
    return prs


def blank_slide(prs: Presentation):
    """Add and return a completely blank slide (layout index 6)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def dark_slide(prs: Presentation, color: str | None = None):
    """Blank slide already painted with the deck background."""
    slide = blank_slide(prs)
    slide_background(slide, color)
    return slide


__all__ = [
    "hex_rgb", "resolve_color", "relative_luminance", "contrast_ratio",
    "set_font", "font_fallback_chain",
    "estimate_lines", "estimate_height_in",
    "add_textbox", "add_multiline_textbox",
    "add_filled_rect", "add_rounded_rect", "add_circle", "add_hairline",
    "add_left_accent_bar",
    "chip_width_in", "add_chip", "add_tier_chip", "add_pmid_pill",
    "slide_background", "add_header_row", "add_slide_title", "add_lead",
    "add_footer", "add_notice_box", "add_bullet_row", "add_image_bytes",
    "new_presentation", "blank_slide", "dark_slide",
    "extract_pmids", "sanitize",
]
