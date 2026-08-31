"""
Curo — Deck Dispatcher
=======================
Maps the "pattern" key in a slide spec to a layout function and builds a
complete python-pptx Presentation. Also handles the legacy "type" key format
produced by `generate_slides_content()`.

Usage
-----
    from presentations.build_deck import build_deck_from_specs

    deck = generate_slides_specs(answer, question, length_minutes)
    prs, slides_queue = build_deck_from_specs(deck, source_text=answer)
    prs.save(output_path)
    # slides_queue: [(slide_obj, speaker_notes_str, slide_num_1based), ...]

Backward compatibility
----------------------
`app.py` calls `build_deck_from_specs(deck)` with one positional argument and
must keep working unmodified, so every parameter added here is keyword-only
with a safe default. Two of them matter:

* **`source_text`** — the canonical answer/curriculum text. Charts are gated on
  it: spec §1.5 requires every plotted value to appear verbatim in the cited
  source, so with no source there is nothing to verify against and no chart is
  drawn. It can also travel inside the deck dict as `deck["source_text"]`.

* **`tier_counts`** — per-tier paper counts for the title slide's
  evidence-shape bar (`deck["tier_counts"]` also works). Absent, the card is
  still drawn — it is the product signature and mandatory — but says the
  breakdown is unavailable instead of inventing a shape.

A layout may return a LIST of slides when its body overflows the spec's body
budget and auto-splits onto continuation slides. The dispatcher flattens the
result and keeps page numbers running across the split, so a deck can finish
with more slides than the generator asked for.
"""

from __future__ import annotations

import warnings

from pptx.presentation import Presentation

from presentations.slide_helpers import new_presentation, dark_slide, add_footer
from presentations.design_tokens import COLORS, LAYOUT, px_in
from presentations.text_budget import sanitize
from presentations.slide_patterns import (
    # the eight approved layouts + the notice slide
    title_slide, section_divider, content_slide, table_slide, decision_tree,
    chart_slide, takeaways_slide, references_slide, notice_slide,
    # adapters for the vocabulary generate_slides_specs() emits
    objectives_slide, two_column_compare, cascade_slide, decision_table,
    three_route_grid, stat_panel, evidence_summary,
)

# ── Pattern dispatch table ───────────────────────────────────────────────────
PATTERN_DISPATCH: dict[str, callable] = {
    # spec §1.4
    "title_slide":       title_slide,
    "section_divider":   section_divider,
    "content_slide":     content_slide,
    "table_slide":       table_slide,
    "decision_tree":     decision_tree,
    "chart_slide":       chart_slide,
    "takeaways_slide":   takeaways_slide,
    "references_slide":  references_slide,
    "notice_slide":      notice_slide,
    # generator vocabulary — adapters onto the layouts above
    "objectives_slide":   objectives_slide,
    "two_column_compare": two_column_compare,
    "cascade_slide":      cascade_slide,
    "decision_table":     decision_table,
    "three_route_grid":   three_route_grid,
    "stat_panel":         stat_panel,
    "evidence_summary":   evidence_summary,
}

# ── Legacy "type" → pattern fallback mapping ─────────────────────────────────
_LEGACY_TYPE_MAP: dict[str, str] = {
    "title":            "title_slide",
    "summary":          "takeaways_slide",
    "stat_cards":       "stat_panel",
    "type_cards":       "content_slide",
    "numbered_grid":    "content_slide",
    "bullets":          "content_slide",
    "comparison_table": "decision_table",
    "references":       "references_slide",
    "chart_bar":        "stat_panel",
}

_DEFAULT_PATTERN = "content_slide"

# A module the evidence gate refused to write reaches the deck as prose. It
# gets the restyled notice, never a bullet slide implying content exists.
_NOTICE_MARKERS = ("module not generated", "insufficient evidence")


def _coerce_legacy_slide(spec: dict, deck: dict) -> dict:
    """Convert an old-format ("type"-keyed) slide dict to a pattern spec."""
    stype = str(spec.get("type", "")).lower()
    pattern = _LEGACY_TYPE_MAP.get(stype, _DEFAULT_PATTERN)

    if pattern == "title_slide":
        return {
            "pattern": "title_slide",
            "eyebrow": spec.get("eyebrow", deck.get("footer", "CURO")),
            "title": spec.get("title", deck.get("title", "")),
            "subtitle": spec.get("subtitle", deck.get("subtitle", "")),
            "footer_metadata": spec.get("footer_metadata", ""),
            "speaker_notes": spec.get("speaker_notes", ""),
        }

    if pattern == "stat_panel":
        cards = spec.get("cards") or spec.get("categories") or []
        values = spec.get("values") or []
        def _card(i):
            if i < len(cards) and isinstance(cards[i], dict):
                return str(cards[i].get("value", "")), str(cards[i].get("label", ""))
            if i < len(values):
                label = ""
                cats = spec.get("categories") or []
                if i < len(cats) and isinstance(cats[i], str):
                    label = cats[i]
                return str(values[i]), label
            return None, None
        p_stat, p_label = _card(0)
        s_stat, s_label = _card(1)
        return {
            "pattern": "stat_panel",
            "eyebrow": spec.get("eyebrow", ""),
            "title": spec.get("title", ""),
            "primary_stat": p_stat or "",
            "primary_label": p_label or "",
            "secondary_stat": s_stat,
            "secondary_label": s_label,
            "unit": spec.get("unit", ""),
            "categories": spec.get("categories"),
            "values": spec.get("values"),
            "speaker_notes": spec.get("speaker_notes", ""),
        }

    if pattern == "decision_table":
        rows = []
        for row in spec.get("rows", []) or []:
            if isinstance(row, (list, tuple)):
                rows.append({
                    "finding": row[0] if len(row) > 0 else "",
                    "implication": row[1] if len(row) > 1 else "",
                    "path": row[2] if len(row) > 2 else "",
                })
            elif isinstance(row, dict):
                rows.append(row)
        return {
            "pattern": "decision_table",
            "eyebrow": spec.get("eyebrow", ""),
            "title": spec.get("title", ""),
            "rows": rows,
            "speaker_notes": spec.get("speaker_notes", ""),
        }

    if pattern == "takeaways_slide":
        items = []
        for j, b in enumerate(spec.get("bullets") or spec.get("items") or []):
            if isinstance(b, dict):
                items.append(b)
            else:
                items.append({"number": f"{j + 1:02d}", "header": str(b),
                              "body": ""})
        return {
            "pattern": "takeaways_slide",
            "eyebrow": spec.get("eyebrow", "KEY TAKEAWAYS"),
            "title": spec.get("title", ""),
            "items": items,
            "speaker_notes": spec.get("speaker_notes", ""),
        }

    if pattern == "references_slide":
        refs = []
        for item in spec.get("items") or spec.get("references") or []:
            refs.append(item if isinstance(item, dict)
                        else {"citation": str(item)})
        return {
            "pattern": "references_slide",
            "eyebrow": spec.get("eyebrow", ""),
            "title": spec.get("title", "References"),
            "references": refs,
            "speaker_notes": spec.get("speaker_notes", ""),
        }

    bullets = (spec.get("bullets") or spec.get("items") or spec.get("cards")
               or [])
    return {
        "pattern": "content_slide",
        "eyebrow": spec.get("eyebrow", ""),
        "title": spec.get("title", ""),
        "bullets": list(bullets),
        "speaker_notes": spec.get("speaker_notes", ""),
    }


def _looks_like_notice(spec: dict) -> bool:
    blob = " ".join(
        str(v) for k, v in spec.items()
        if isinstance(v, str) and k not in ("speaker_notes",)
    ).lower()
    return any(marker in blob for marker in _NOTICE_MARKERS)


def _module_positions(slides_list: list) -> dict[int, tuple[int, int]]:
    """Index every section divider so it can say 'MODULE N OF M'."""
    idx = [i for i, s in enumerate(slides_list)
           if str((s or {}).get("pattern", "")).lower() == "section_divider"]
    total = len(idx)
    return {i: (n + 1, total) for n, i in enumerate(idx)}


def build_deck_from_specs(
    deck: dict,
    *,
    section_label: str | None = None,
    source_text: str | None = None,
    tier_counts: dict | None = None,
) -> tuple[Presentation, list[tuple]]:
    """Build a Presentation from a pattern-based (or legacy) deck spec.

    Returns (prs, slides_queue) where slides_queue is
    [(slide, speaker_notes, page_num_1based), ...] for the TTS pipeline.

    Speaker notes are attached to the FIRST slide of a split group only, so a
    body that auto-splits is narrated once rather than repeated per
    continuation slide. The notes text itself is passed through verbatim: it is
    the canonical narration and this phase does not author or edit it.
    """
    slides_list = list(deck.get("slides", []) or [])
    footer_base = section_label or deck.get("footer", "CURO")
    source_text = source_text or deck.get("source_text") or ""
    tier_counts = tier_counts or deck.get("tier_counts") or None
    modules = _module_positions(slides_list)

    prs = new_presentation()
    slides_queue: list[tuple] = []
    page = 1

    for i, raw_spec in enumerate(slides_list):
        spec = dict(raw_spec or {})

        pattern_name = spec.pop("pattern", None)
        if pattern_name is None:
            spec = _coerce_legacy_slide(spec, deck)
            pattern_name = spec.pop("pattern", _DEFAULT_PATTERN)

        if _looks_like_notice(spec):
            pattern_name = "notice_slide"
            spec = {
                "title": sanitize(spec.get("title") or "Module not generated"),
                "eyebrow": spec.get("eyebrow", footer_base),
                "message": sanitize(spec.get("body") or spec.get("subtitle")
                                    or "Insufficient evidence retrieved for "
                                       "this module."),
                "speaker_notes": spec.get("speaker_notes", ""),
            }

        pattern_fn = PATTERN_DISPATCH.get(pattern_name)
        if pattern_fn is None:
            warnings.warn(
                f"[build_deck] Unknown pattern '{pattern_name}' on slide "
                f"{i + 1} — falling back to {_DEFAULT_PATTERN}",
                stacklevel=2,
            )
            pattern_fn = PATTERN_DISPATCH[_DEFAULT_PATTERN]

        notes = str(spec.pop("speaker_notes", "") or "")

        spec.setdefault("eyebrow", footer_base)
        spec["_page_num"] = page
        spec["_total_pages"] = len(slides_list)
        spec["_source_text"] = source_text
        if pattern_name == "title_slide" and tier_counts:
            spec.setdefault("tier_counts", tier_counts)
        if pattern_name == "section_divider" and i in modules:
            spec["_module_index"], spec["_module_total"] = modules[i]

        try:
            result = pattern_fn(prs, **spec)
        except Exception as exc:            # pragma: no cover - defensive
            warnings.warn(
                f"[build_deck] Pattern '{pattern_name}' failed on slide "
                f"{i + 1}: {exc}", stacklevel=2,
            )
            result = _error_slide(prs, spec, exc, page)

        produced = result if isinstance(result, list) else [result]
        for j, slide in enumerate(produced):
            slide_notes = notes if j == 0 else ""
            if slide_notes:
                try:
                    slide.notes_slide.notes_text_frame.text = slide_notes
                except Exception:
                    pass
            slides_queue.append((slide, slide_notes, page))
            page += 1

    return prs, slides_queue


def _error_slide(prs, spec: dict, exc: Exception, page: int):
    """Last-resort slide. Dark theme, notice framing — never a red error page."""
    from presentations.slide_helpers import add_slide_title, add_notice_box
    slide = dark_slide(prs)
    add_slide_title(slide, sanitize(spec.get("title") or f"Slide {page}"))
    add_notice_box(
        slide,
        "This slide could not be laid out and has been left blank rather than "
        "rendered incorrectly.",
        LAYOUT["pad_x_in"], px_in(300), LAYOUT["content_w_in"], px_in(120),
        heading="Layout unavailable",
    )
    add_footer(slide, citations="", page_num=page)
    return slide


__all__ = ["build_deck_from_specs", "PATTERN_DISPATCH"]
