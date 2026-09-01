"""
Curo — deck invariants (PRESENTATION_WORKLIST §2.5)

Five rules, each of which has a way of failing silently, which is why each is
pinned here rather than left to review:

  1. the body-budget split                (a slide that overflows must split)
  2. no raw citation marker on a slide    (`[[PMID:N]]` must never render)
  3. the footer is on every content slide (rule + page number)
  4. a tier chip's colours ARE the tokens (not a lookalike hex)
  5. the chart hard rules of spec §1.5    (uncited -> no chart; cited -> chart)

Every test here was mutation-checked: the bug it guards was reintroduced, the
test was confirmed to fail, and the code was restored. A test that could not
be made to fail was deleted rather than kept.

NOTE ON DISCOVERY: pytest.ini sets `testpaths = tests presentations`, so a
bare `python -m pytest` collects this file. It did not, once — a guard that
silently never runs is the fail-open bug class itself.
"""

from __future__ import annotations

import json
import pathlib

import pytest
from pptx.util import Emu

from presentations.build_deck import build_deck_from_specs
from presentations.chart_data import detect_chartable, evidence_shape, verbatim_in
from presentations.design_tokens import (
    BODY_BUDGET, COLORS, LAYOUT, TIER_CHART_FILL, TIER_CHIP, TIER_LABELS,
    TIER_ORDER,
)
from presentations.slide_helpers import (
    add_textbox, add_tier_chip, contrast_ratio, dark_slide, hex_rgb,
    new_presentation,
)
from presentations.slide_patterns import content_slide, table_slide
from presentations.text_budget import (
    bullet_cost, has_raw_marker, sanitize, split_bullets, split_rows,
)

EMU_PER_IN = 914400


def _emu(inches: float) -> int:
    return int(round(inches * EMU_PER_IN))


def _texts(slide) -> list[str]:
    return [sh.text_frame.text for sh in slide.shapes
            if sh.has_text_frame and sh.text_frame.text.strip()]


def _fills(slide) -> list[str]:
    out = []
    for sh in slide.shapes:
        try:
            if sh.fill.type is not None and sh.fill.fore_color.rgb is not None:
                out.append(str(sh.fill.fore_color.rgb).upper())
        except Exception:
            continue
    return out


def _slides(result) -> list:
    return result if isinstance(result, list) else [result]


# ─────────────────────────────────────────────────────────────────────────────
# 1. The body-budget split
# ─────────────────────────────────────────────────────────────────────────────

class TestBodyBudgetSplit:
    def test_more_than_five_bullets_splits(self):
        bullets = [f"Short bullet number {i}." for i in range(9)]
        pages = split_bullets(bullets)
        assert len(pages) > 1
        assert all(len(p) <= BODY_BUDGET["max_bullets"] for p in pages)
        assert sum(len(p) for p in pages) == len(bullets)

    def test_five_short_bullets_stay_on_one_slide(self):
        bullets = [f"Bullet {i}." for i in range(BODY_BUDGET["max_bullets"])]
        assert len(split_bullets(bullets)) == 1

    def test_an_over_long_bullet_costs_more_than_one_slot(self):
        long_bullet = " ".join(["word"] * 60)
        assert bullet_cost(long_bullet) == 3          # ceil(60 / 25)
        assert bullet_cost("five words in this one") == 1

    def test_an_over_long_bullet_is_never_truncated(self):
        """The prime rule: rendering only. Splitting moves text between
        slides; it must never shorten a clinical sentence."""
        long_bullet = " ".join(f"w{i}" for i in range(80))
        pages = split_bullets([long_bullet, "Second bullet."])
        rendered = [b for page in pages for b in page]
        assert long_bullet in rendered

    def test_split_never_drops_a_bullet(self):
        bullets = [" ".join(["word"] * n) for n in (5, 40, 5, 70, 5, 5, 5)]
        pages = split_bullets(bullets)
        flat = [b for page in pages for b in page]
        assert flat == bullets

    def test_content_slide_emits_continuation_slides(self):
        prs = new_presentation()
        result = content_slide(prs, title="Overflowing slide", eyebrow="MODULE 1",
                               bullets=[f"Bullet {i} text here." for i in range(11)],
                               _page_num=1)
        slides = _slides(result)
        assert len(slides) > 1
        # The continuation is marked in the eyebrow's slide-role slot, and the
        # clinical title is left exactly as authored.
        assert any("CONTINUED" in t.upper() for t in _texts(slides[1]))
        assert any("Overflowing slide" in t for t in _texts(slides[1]))

    def test_long_table_splits_with_the_header_repeated(self):
        headers = ["Parameter", "Finding", "Action"]
        rows = [[f"row {i}", "finding", "action"] for i in range(16)]
        assert len(split_rows(rows)) > 1

        prs = new_presentation()
        slides = _slides(table_slide(prs, title="Long table", headers=headers,
                                     rows=rows, _page_num=1))
        assert len(slides) > 1
        for slide in slides:
            texts = [t.upper() for t in _texts(slide)]
            for header in headers:
                assert header.upper() in texts, (
                    f"header {header!r} missing from a split table page")


# ─────────────────────────────────────────────────────────────────────────────
# 2. No raw citation marker reaches a slide
# ─────────────────────────────────────────────────────────────────────────────

MARKER_SPEC = {
    "slides": [
        {
            "pattern": "objectives_slide",
            "eyebrow": "MODULE 1 [[PMID:11111]]",
            "title": "Title with [[PMID:22222]] marker",
            "items": [
                {"header": "Header [[PMID:33333]]",
                 "body": "Body text [PMID:44444] single bracket too."},
            ],
            "closing_callout": "Lead line [[PMID: 55555]] with a space.",
            "speaker_notes": "Narration keeps its [[PMID:66666]] markers.",
        },
        {
            "pattern": "decision_table",
            "title": "Table [[PMID:77777]]",
            "rows": [{"finding": "F [[PMID:88888]]", "implication": "I",
                      "path": "P [PMID:99999]"}],
            "footer_caption": "Derived from ESE SR 2023 [[PMID:36156804]].",
        },
    ]
}


class TestNoRawMarkers:
    def test_sanitize_strips_every_bracketed_shape(self):
        for raw in ("[[PMID:12345]]", "[PMID:12345]", "[[PMID: 12345]]",
                    "[[pmid:12345]]", "[[PMID:123, 456]]"):
            assert not has_raw_marker(sanitize(f"text {raw} more"))

    def test_a_bare_pmid_is_the_footer_format_and_survives(self):
        """Spec §1.3's own footer example ends '... n = 100 · PMID 28294701'.
        Stripping that would delete the citation the footer exists to show."""
        out = sanitize("Schulte-Lünzum et al. 2017 · n = 100 · PMID 28294701")
        assert "PMID 28294701" in out
        assert not has_raw_marker(out)

    def test_add_textbox_is_the_chokepoint(self):
        prs = new_presentation()
        slide = dark_slide(prs)
        add_textbox(slide, "A claim [[PMID:12345]] with a marker.",
                    1.0, 1.0, 5.0, 0.5)
        assert not any(has_raw_marker(t) for t in _texts(slide))

    def test_no_marker_survives_a_full_build(self):
        prs, queue = build_deck_from_specs(MARKER_SPEC)
        offenders = [
            (n, t) for slide, _, n in queue for t in _texts(slide)
            if has_raw_marker(t)
        ]
        assert offenders == [], f"raw markers rendered: {offenders}"

    def test_markers_become_footer_citations_rather_than_vanishing(self):
        prs, queue = build_deck_from_specs(MARKER_SPEC)
        blob = " ".join(t for slide, _, _ in queue for t in _texts(slide))
        assert "PMID 36156804" in blob, (
            "a cited PMID was stripped instead of being moved to the footer")


# ─────────────────────────────────────────────────────────────────────────────
# 3. The footer is on every content-class slide
# ─────────────────────────────────────────────────────────────────────────────

FOOTER_SPEC = {
    "slides": [
        {"pattern": "objectives_slide", "title": "Objectives", "eyebrow": "M1",
         "items": [{"header": "One", "body": "Body."}]},
        {"pattern": "decision_table", "title": "Table", "eyebrow": "M1",
         "rows": [{"finding": "F", "implication": "I", "path": "P"}]},
        {"pattern": "three_route_grid", "title": "Routes", "eyebrow": "M1",
         "routes": [{"name": "A", "when": "w", "how": "h"},
                    {"name": "B", "when": "w", "how": "h"}]},
        {"pattern": "takeaways_slide", "title": "Takeaways", "eyebrow": "M1",
         "items": [{"number": "01", "header": "H", "body": "B"}]},
        {"pattern": "references_slide", "title": "References", "eyebrow": "M1",
         "references": [{"citation": "Author et al.", "journal": "J Endod",
                         "year": 2023, "pmid": "36156804", "tier": "level1",
                         "score": 70}]},
    ]
}


def _has_footer_rule(slide) -> bool:
    want_y = _emu(LAYOUT["footer_rule_y_in"])
    want_h = _emu(LAYOUT["hairline_in"])
    for sh in slide.shapes:
        if sh.top is None or sh.height is None:
            continue
        if abs(sh.top - want_y) < 8000 and abs(sh.height - want_h) < 8000:
            return True
    return False


def _has_page_number(slide, page_num: int) -> bool:
    want_y = _emu(LAYOUT["footer_text_y_in"])
    for sh in slide.shapes:
        if not sh.has_text_frame or sh.top is None:
            continue
        if abs(sh.top - want_y) < 20000 and sh.text_frame.text.strip() == str(page_num):
            return True
    return False


class TestFooterPresence:
    def test_every_content_slide_has_the_footer_rule_and_page_number(self):
        prs, queue = build_deck_from_specs(FOOTER_SPEC)
        assert len(queue) == len(FOOTER_SPEC["slides"])
        for slide, _, page_num in queue:
            assert _has_footer_rule(slide), f"slide {page_num} has no footer rule"
            assert _has_page_number(slide, page_num), (
                f"slide {page_num} has no page number in the footer")

    def test_the_footer_rule_uses_the_border_token(self):
        prs, queue = build_deck_from_specs(FOOTER_SPEC)
        slide = queue[0][0]
        want_y = _emu(LAYOUT["footer_rule_y_in"])
        rules = [sh for sh in slide.shapes
                 if sh.top is not None and abs(sh.top - want_y) < 8000]
        assert rules
        assert str(rules[0].fill.fore_color.rgb).upper() == \
            COLORS["border"].lstrip("#").upper()

    def test_page_numbers_run_across_a_split(self):
        spec = {"slides": [
            {"pattern": "objectives_slide", "title": "Long", "eyebrow": "M1",
             "items": [{"header": f"H{i}", "body": "Body text."}
                       for i in range(12)]},
            {"pattern": "takeaways_slide", "title": "End", "eyebrow": "M1",
             "items": [{"number": "01", "header": "H", "body": "B"}]},
        ]}
        prs, queue = build_deck_from_specs(spec)
        assert len(queue) > 2, "the overflowing slide did not split"
        assert [n for _, _, n in queue] == list(range(1, len(queue) + 1))
        for slide, _, page_num in queue:
            assert _has_page_number(slide, page_num)


# ─────────────────────────────────────────────────────────────────────────────
# 4. A tier chip's colours ARE the token values
# ─────────────────────────────────────────────────────────────────────────────

class TestTierChipColors:
    @pytest.mark.parametrize("tier", TIER_ORDER)
    def test_chip_background_and_dot_match_the_ladder(self, tier):
        prs = new_presentation()
        slide = dark_slide(prs)
        add_tier_chip(slide, tier, 1.0, 1.0)

        want_bg = TIER_CHIP[tier][0].lstrip("#").upper()
        want_dot = TIER_CHART_FILL[tier].lstrip("#").upper()
        fills = _fills(slide)
        assert want_bg in fills, f"{tier} chip background is not {want_bg}"
        assert want_dot in fills, f"{tier} chip dot is not the ladder fill"

    @pytest.mark.parametrize("tier", TIER_ORDER)
    def test_every_chip_carries_its_text_label(self, tier):
        """Spec §1.2 requires Level III to always carry its label; the whole
        ladder does, so colour is never the only channel carrying the tier."""
        prs = new_presentation()
        slide = dark_slide(prs)
        add_tier_chip(slide, tier, 1.0, 1.0)
        texts = [t.upper() for t in _texts(slide)]
        assert TIER_LABELS[tier].upper() in texts

    def test_chip_text_colour_matches_the_token(self):
        prs = new_presentation()
        slide = dark_slide(prs)
        add_tier_chip(slide, "level3", 1.0, 1.0)
        want = hex_rgb(TIER_CHIP["level3"][1])
        colours = [
            run.font.color.rgb
            for sh in slide.shapes if sh.has_text_frame
            for para in sh.text_frame.paragraphs for run in para.runs
        ]
        assert want in colours

    def test_alert_red_is_never_a_tier_colour(self):
        """Spec §1.2 reserves the COI/alert red family. If it ever appears in
        the ladder, a contraindication and an evidence tier become the same
        signal."""
        assert COLORS["alert_red"] not in TIER_CHART_FILL.values()
        assert COLORS["alert_red"] not in [bg for bg, _ in TIER_CHIP.values()]
        assert COLORS["alert_red"] not in [fg for _, fg in TIER_CHIP.values()]

    def test_body_text_clears_4_5_to_1_on_the_deck_background(self):
        for key in ("text_title", "text_body", "text_secondary", "text_lead",
                    "text_eyebrow", "text_footer", "text_muted"):
            ratio = contrast_ratio(COLORS[key], COLORS["bg"])
            assert ratio >= 4.5, f"{key} is {ratio:.2f}:1 on {COLORS['bg']}"


# ─────────────────────────────────────────────────────────────────────────────
# 5. Chart hard rules (spec §1.5)
# ─────────────────────────────────────────────────────────────────────────────

SOURCE = (
    "Periapical healing reached 96.1% in the laser-assisted arm versus 74.5% "
    "with conventional irrigation at 24 months (n = 412). The pooled "
    "standardised mean difference for 24-hour pain was -1.57."
)


class TestChartHardRules:
    def test_a_cited_comparison_produces_a_chart(self):
        spec = {
            "title": "Healing at 24 months",
            "primary_stat": "96.1%", "primary_label": "Laser-assisted arm",
            "secondary_stat": "74.5%", "secondary_label": "Conventional irrigation",
        }
        chart = detect_chartable(spec, SOURCE)
        assert chart is not None
        assert chart.literals == ["96.1", "74.5"]
        assert chart.unit == "%"

    def test_uncited_numbers_produce_no_chart(self):
        """The numbers are well-formed and plausible; they are simply not in
        the source. That alone must be enough to drop the chart."""
        spec = {
            "title": "Healing at 24 months",
            "primary_stat": "88.4%", "primary_label": "Laser-assisted arm",
            "secondary_stat": "62.2%", "secondary_label": "Conventional irrigation",
        }
        assert detect_chartable(spec, SOURCE) is None

    def test_one_uncited_value_drops_the_whole_chart(self):
        spec = {
            "title": "Healing", "primary_stat": "96.1%", "primary_label": "A",
            "secondary_stat": "62.2%", "secondary_label": "B",
        }
        assert detect_chartable(spec, SOURCE) is None

    def test_a_rounded_value_is_not_verbatim(self):
        assert verbatim_in("96.1", SOURCE)
        assert not verbatim_in("96", SOURCE)     # rounded, not what was written
        assert not verbatim_in("6.1", SOURCE)    # substring of 96.1

    def test_no_source_text_means_no_chart(self):
        spec = {"primary_stat": "96.1%", "primary_label": "A",
                "secondary_stat": "74.5%", "secondary_label": "B"}
        assert detect_chartable(spec, None) is None
        assert detect_chartable(spec, "") is None

    def test_a_single_value_is_not_a_chart(self):
        spec = {"primary_stat": "96.1%", "primary_label": "A"}
        assert detect_chartable(spec, SOURCE) is None

    def test_a_truncated_axis_carries_an_explicit_note(self):
        source = "Success was 98.2% in one arm and 99.1% in the other."
        spec = {"primary_stat": "98.2%", "primary_label": "A",
                "secondary_stat": "99.1%", "secondary_label": "B"}
        chart = detect_chartable(spec, source)
        assert chart is not None
        assert chart.kind == "dot", "near-equal percentages should be a dot plot"
        assert chart.axis_note, "a truncated axis must state where it starts"

    def test_a_magnitude_comparison_is_a_bar_from_zero(self):
        chart = detect_chartable(
            {"primary_stat": "96.1%", "primary_label": "A",
             "secondary_stat": "74.5%", "secondary_label": "B"}, SOURCE)
        assert chart.kind == "bar"
        assert chart.axis_note is None

    def test_the_evidence_shape_needs_real_counts(self):
        assert evidence_shape({}) is None
        assert evidence_shape({"level1": 0, "level2": 0}) is None
        chart = evidence_shape({"cochrane": 3, "level1": 22})
        assert chart is not None
        assert chart.values == [3.0, 22.0]
        assert chart.tier_keys == ["cochrane", "level1"]

    def test_a_deck_built_without_source_text_renders_no_chart(self):
        spec = {"slides": [{
            "pattern": "stat_panel", "title": "Healing", "eyebrow": "M1",
            "primary_stat": "96.1%", "primary_label": "Laser arm",
            "secondary_stat": "74.5%", "secondary_label": "Control",
        }]}
        prs, queue = build_deck_from_specs(spec)
        assert not any(sh.shape_type == 13 for sh in queue[0][0].shapes), (
            "a chart was drawn with nothing to verify the values against")

    def test_a_deck_built_with_source_text_renders_the_chart(self):
        pytest.importorskip("matplotlib")
        spec = {"slides": [{
            "pattern": "stat_panel", "title": "Healing", "eyebrow": "M1",
            "primary_stat": "96.1%", "primary_label": "Laser arm",
            "secondary_stat": "74.5%", "secondary_label": "Control",
        }]}
        prs, queue = build_deck_from_specs(spec, source_text=SOURCE)
        pictures = [sh for sh in queue[0][0].shapes if sh.shape_type == 13]
        assert pictures, "a verified comparison did not produce a chart"


class TestChartUnitConsistency:
    """A shared axis asserts the values are comparable magnitudes. Two real,
    correctly-cited numbers in different units are not.

    Both cases below are from the Phase 4 laser deck, where they rendered as
    charts that passed every provenance check and still misled: an SMD of
    -0.551 beside an I-squared of 23.89% (the bar implied heterogeneity was
    ~43x the effect), and a 24-hour window beside 0 adverse events (one bar
    had no length). Found by looking at the PNGs, not by any assertion.
    """

    def test_matching_units_are_chartable(self):
        from presentations.chart_data import consistent_unit
        assert consistent_unit(["78%", "64%"]) == "%"
        assert consistent_unit(["12 months", "24 months"]) == "months"

    def test_bare_numbers_are_refused(self):
        """Changed deliberately. This used to assert that two unitless values
        share the empty unit and may be charted together. That is what let a
        network-meta P-score of 0.993 chart beside an SMD of -0.58 on a real
        deck: both map to "", and {""} has length one, so the gate saw
        agreement where there was only ignorance.

        A bare number is not evidence of a shared quantity. It costs the
        occasional legitimate count-vs-count chart, which is the cheaper error:
        a suppressed chart omits information, a false one asserts something
        untrue."""
        from presentations.chart_data import consistent_unit
        assert consistent_unit(["12", "34"]) is None
        assert consistent_unit(["0.993", "-0.58"]) is None

    def test_a_single_unitless_value_still_resolves(self):
        """The refusal is about PAIRING unitless values, not about the unit
        lookup itself — a lone literal has nothing to disagree with."""
        from presentations.chart_data import consistent_unit
        assert consistent_unit(["12"]) == ""

    def test_effect_size_beside_a_percentage_is_refused(self):
        from presentations.chart_data import consistent_unit
        assert consistent_unit(["-0.551", "23.89%"]) is None

    def test_a_duration_beside_a_bare_count_is_refused(self):
        from presentations.chart_data import consistent_unit
        assert consistent_unit(["24 hours", "0"]) is None

    def test_stat_panel_with_mixed_units_produces_no_chart(self):
        """End of the path: the builder must return None, not a chart."""
        from presentations.chart_data import detect_chartable
        spec = {"pattern": "stat_panel", "title": "Mixed",
                "primary_stat": "-0.551", "primary_label": "SMD",
                "secondary_stat": "23.89%", "secondary_label": "I2"}
        src = "The SMD was -0.551 and heterogeneity was 23.89%."
        assert detect_chartable(spec, src) is None

    def test_stat_panel_with_one_unit_still_charts(self):
        """The gate must not have suppressed charting altogether."""
        from presentations.chart_data import detect_chartable
        spec = {"pattern": "stat_panel", "title": "Success",
                "primary_stat": "78%", "primary_label": "Laser",
                "secondary_stat": "64%", "secondary_label": "Control"}
        src = "Success was 78% with the laser and 64% in controls."
        assert detect_chartable(spec, src) is not None


class TestChartRangeValues:
    """A bar draws one number. Charting a range plots its lower bound and
    silently discards the rest, so the value on the slide stops being the value
    in the source — the one thing 1.5 does not allow.

    From the Phase 4 laser deck: primary_stat was "24-48 h" and rendered as a
    bar of 24, beside a bar of 0 severe adverse events. Every number was real
    and cited; the chart still misrepresented both of them.
    """

    def test_a_span_is_detected(self):
        from presentations.chart_data import is_range
        assert is_range("24-48 h")
        assert is_range("24–48 h")      # en dash, as the generator emits
        assert is_range("3 to 5 mm")

    def test_a_scalar_is_not_a_span(self):
        from presentations.chart_data import is_range
        assert not is_range("78%")
        assert not is_range("-0.551")
        assert not is_range("0")

    def test_abbreviated_units_are_recognised(self):
        from presentations.chart_data import consistent_unit
        assert consistent_unit(["15 min", "30 min"]) == "min"
        assert consistent_unit(["24 h", "12"]) is None

    def test_a_unit_letter_inside_a_word_is_not_a_unit(self):
        """The word boundaries are load-bearing. Without them the "h" in
        "Charter" reads as hours, the two literals disagree, and a perfectly
        chartable comparison is silently dropped."""
        from presentations.chart_data import _unit_of
        assert _unit_of("Charter") == ""
        assert _unit_of("18 charts") == ""
        # Asserted at the _unit_of level: consistent_unit now refuses unitless
        # PAIRS outright, so it can no longer distinguish "the h in charts was
        # correctly ignored" from "both values are unitless".

    def test_a_range_produces_no_chart(self):
        """Both literals carry the SAME unit here, deliberately: with mismatched
        units the unit gate fires first and this test would pass even with the
        range gate deleted. Mutation-checking caught exactly that."""
        from presentations.chart_data import detect_chartable
        spec = {"pattern": "stat_panel", "title": "T",
                "primary_stat": "24–48 h", "primary_label": "Window",
                "secondary_stat": "12 h", "secondary_label": "Comparator"}
        src = "superior for 24–48 h against 12 h for the comparator"
        assert detect_chartable(spec, src) is None

    def test_two_scalars_still_chart(self):
        """The gate must not have suppressed charting altogether."""
        from presentations.chart_data import detect_chartable
        spec = {"pattern": "stat_panel", "title": "T",
                "primary_stat": "78%", "primary_label": "Laser",
                "secondary_stat": "64%", "secondary_label": "Control"}
        assert detect_chartable(spec, "78% with laser versus 64% control") is not None


# ─────────────────────────────────────────────────────────────────────────────
# 6. The body budget is a HEIGHT budget too  (CURO_HANDOVER §5[D])
# ─────────────────────────────────────────────────────────────────────────────

_SPECS_DIR = pathlib.Path(__file__).resolve().parent.parent / "slide_specs"


def _cached_specs(pattern: str):
    """Every cached spec of one pattern — real generator output, not fixtures.

    The overflow this guards was invisible to hand-written fixtures: it needs a
    real cascade whose five steps each fit the WORD budget while their rendered
    header+body pairs do not fit the frame.
    """
    out = []
    for path in sorted(_SPECS_DIR.glob("*.json")):
        doc = json.loads(path.read_text(encoding="utf-8"))
        slides = doc.get("slides") if isinstance(doc, dict) else doc
        for i, spec in enumerate(slides or []):
            if isinstance(spec, dict) and spec.get("pattern") == pattern:
                out.append((f"{path.stem[:8]}#{i}", spec))
    return out


def _body_bottom(slide) -> float:
    """Lowest bottom edge of any BODY shape, in inches.

    The full-bleed background is excluded by height, and footer furniture by
    position — both live at or below the rule by design.
    """
    worst = 0.0
    for shp in slide.shapes:
        top = Emu(shp.top).inches
        h = Emu(shp.height).inches
        if h >= LAYOUT["slide_h_in"] - 0.01:
            continue
        if top >= LAYOUT["footer_rule_y_in"] - 0.01:
            continue
        worst = max(worst, top + h)
    return worst


class TestBodyFitsAboveTheFooterRule:
    """`content_slide` was handed `avail` by `_content_frame` and ignored it.

    `split_bullets` packs by words, and five cascade steps of under 25 words
    each cost five slots and "fit" — while each one renders as a bold header
    line PLUS a body paragraph. Spec 01e071f7 slide 9 drew its body down to
    7.385in against a 7.000in footer rule.
    """

    @pytest.mark.parametrize("name,spec", _cached_specs("cascade_slide"),
                             ids=lambda v: v if isinstance(v, str) else "")
    def test_no_cascade_page_crosses_the_rule(self, name, spec):
        from presentations.slide_patterns import cascade_slide
        prs = new_presentation()
        out = cascade_slide(prs, **{k: v for k, v in spec.items()
                                    if k != "pattern"})
        pages = out if isinstance(out, list) else [out]
        for n, slide in enumerate(pages, 1):
            bottom = _body_bottom(slide)
            assert bottom <= LAYOUT["footer_rule_y_in"], (
                f"{name} page {n}/{len(pages)} body reaches {bottom:.3f}in, "
                f"footer rule is at {LAYOUT['footer_rule_y_in']:.3f}in")

    @pytest.mark.parametrize("name,spec", _cached_specs("cascade_slide"),
                             ids=lambda v: v if isinstance(v, str) else "")
    def test_splitting_never_drops_a_step(self, name, spec):
        """Fitting the frame must never be achieved by dropping content."""
        from presentations.slide_patterns import cascade_slide
        prs = new_presentation()
        out = cascade_slide(prs, **{k: v for k, v in spec.items()
                                    if k != "pattern"})
        pages = out if isinstance(out, list) else [out]
        rendered = " ".join(t for slide in pages for t in _texts(slide))
        for step in spec.get("steps") or []:
            header = sanitize(str(step.get("header") or ""))
            if header:
                assert header in rendered, f"{name}: a step vanished: {header}"

    def test_a_page_still_holds_more_than_one_bullet(self):
        """Guards the opposite failure: a height check that splits everything
        onto its own slide also never overflows."""
        prs = new_presentation()
        steps = [{"number": f"0{i}", "header": f"Step {i}",
                  "body": "One short sentence."} for i in range(1, 6)]
        out = content_slide(prs, title="T", bullets=steps)
        pages = out if isinstance(out, list) else [out]
        assert len(pages) == 1, "five one-line bullets must share one slide"


class TestEvidenceSummaryStatColumn:
    """A "Reported" header over an entirely blank column reads as "these
    studies reported nothing". Every evidence_summary in the cached specs is
    stat-free, so this was the normal rendering, not an edge case."""

    @pytest.mark.parametrize("name,spec", _cached_specs("evidence_summary"),
                             ids=lambda v: v if isinstance(v, str) else "")
    def test_no_reported_header_when_no_row_carries_a_stat(self, name, spec):
        from presentations.slide_patterns import evidence_summary
        rows = [r for r in (spec.get("hierarchy_rows") or [])
                if isinstance(r, dict)]
        if any(str(r.get("stat") or "").strip() for r in rows):
            pytest.skip("this cached spec does carry a stat")
        prs = new_presentation()
        out = evidence_summary(prs, **{k: v for k, v in spec.items()
                                       if k != "pattern"})
        slide = out[0] if isinstance(out, list) else out
        headers = [t.strip().upper() for t in _texts(slide)]
        assert "REPORTED" not in headers, f"{name}: blank REPORTED column drawn"
        assert "TIER" in headers, f"{name}: the table itself disappeared"

    def test_the_column_returns_when_a_row_has_a_stat(self):
        """The column must be dropped for absence of data, not deleted."""
        from presentations.slide_patterns import evidence_summary
        prs = new_presentation()
        out = evidence_summary(prs, title="T", hierarchy_rows=[
            {"tier_label": "LEVEL I", "description": "RCTs", "stat": "86.9%"},
            {"tier_label": "LEVEL II", "description": "Cohorts", "stat": "74.5%"},
        ])
        slide = out[0] if isinstance(out, list) else out
        headers = [t.strip().upper() for t in _texts(slide)]
        assert "REPORTED" in headers


class TestMultiArmComparison:
    """`primary_stat`/`secondary_stat` hold two arms. A three-concentration or
    four-wavelength comparison had nowhere to go: `stat_panel` assembled a spec
    dict of only the two-arm keys, so `categories`/`values` never reached
    `detect_chartable` and `_from_explicit_chart` never fired."""

    SOURCE = ("1% NaOCl achieved 78.4% reduction, 2.5% NaOCl achieved 88.1% "
              "reduction, and 5.25% NaOCl achieved 96.2% reduction.")
    ARMS = [{"label": "1% NaOCl", "stat": "78.4%"},
            {"label": "2.5% NaOCl", "stat": "88.1%"},
            {"label": "5.25% NaOCl", "stat": "96.2%"}]

    def _render(self, arms, source=None):
        from presentations.slide_patterns import stat_panel
        prs = new_presentation()
        out = stat_panel(prs, title="Bacterial reduction", arms=arms,
                         citation="Vertucci et al. 2024 [[PMID:12345678]]",
                         _source_text=source or self.SOURCE)
        slide = out[0] if isinstance(out, list) else out
        picture = [sh for sh in slide.shapes if sh.shape_type == 13]
        return slide, picture

    def test_three_arms_reach_a_chart(self):
        slide, picture = self._render(self.ARMS)
        assert picture, "a verified three-arm comparison must plot"

    def test_every_arm_value_must_be_verbatim_in_the_source(self):
        arms = self.ARMS + [{"label": "6% NaOCl", "stat": "99.9%"}]
        slide, picture = self._render(arms)
        assert not picture, "an arm absent from the source must kill the chart"
        rendered = " ".join(_texts(slide))
        assert "99.9%" in rendered, "the refused value is still slide content"

    def test_mixed_units_across_arms_produce_no_chart(self):
        arms = self.ARMS[:2] + [{"label": "contact time", "stat": "30 min"}]
        src = self.SOURCE + " Contact time was 30 min."
        slide, picture = self._render(arms, src)
        assert not picture

    def test_a_range_arm_produces_no_chart(self):
        """Every arm carries the SAME unit here, deliberately. With a mixed
        unit the unit gate fires first and this test passes with the range
        gate deleted — mutation-checking caught exactly that, as it did once
        before for the two-arm version above."""
        arms = self.ARMS[:2] + [{"label": "pooled", "stat": "80-90%"}]
        src = self.SOURCE + " Pooled estimates spanned 80-90%."
        slide, picture = self._render(arms, src)
        assert not picture

    def test_unitless_arms_produce_no_chart(self):
        arms = [{"label": "P-score", "stat": "0.993"},
                {"label": "SMD", "stat": "0.58"},
                {"label": "OR", "stat": "1.42"}]
        src = "P-score 0.993, SMD 0.58 and OR 1.42 were reported."
        slide, picture = self._render(arms, src)
        assert not picture

    def test_the_pmid_reaches_the_footer_not_the_body(self):
        """Two independent routes put the PMID in the footer — the citation
        harvest in `_citation_fields` and `chart.pmids` via `chart_slide` — so
        killing this test takes removing BOTH (verified: it fails when
        `_spec_pmids` returns [] and the harvest loop is emptied). One route
        alone going missing is invisible here, and deliberately so: the
        invariant is that the marker arrives, not which path carried it."""
        slide, _ = self._render(self.ARMS)
        rendered = " ".join(_texts(slide))
        assert "12345678" in rendered
        assert not has_raw_marker(rendered)
