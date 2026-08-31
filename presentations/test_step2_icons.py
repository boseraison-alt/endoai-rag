"""
Step 2 checkpoint — render a PPTX proof slide for every icon in ICON_MAP.

Run:
    python presentations/test_step2_icons.py

Output:
    presentations/test_step2_icons.pptx   (open in PowerPoint to verify)
    presentations/test_step2_icons.png    (thumbnail if python-pptx + Pillow available)
"""

import sys, os
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from presentations.design_tokens import (
    COLORS, FONTS, SIZES, LAYOUT, CHART_SERIES_SINGLE,
)
from presentations.icons import ICON_MAP, get_icon_char

OUT_PPTX = os.path.join(os.path.dirname(__file__), "test_step2_icons.pptx")

# ── helpers ──────────────────────────────────────────────────────────────────

def rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_text_box(slide, text, x, y, w, h, *, font_name, size, color, bold=False,
                 align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return txb


def add_rect(slide, x, y, w, h, fill_hex):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    shape.line.fill.background()
    return shape


# ── build slide ───────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = Inches(LAYOUT["slide_w_in"])
    prs.slide_height = Inches(LAYOUT["slide_h_in"])

    blank_layout = prs.slide_layouts[6]  # completely blank
    slide = prs.slides.add_slide(blank_layout)

    # Dark background
    add_rect(slide, 0, 0, LAYOUT["slide_w_in"], LAYOUT["slide_h_in"], COLORS["bg"])

    # Title
    add_text_box(slide, "Step 2 — icon glyph proof sheet",
                 LAYOUT["pad_x_in"], 0.20,
                 12.0, 0.60,
                 font_name=FONTS["display"],
                 size=SIZES["subtitle"],
                 color=COLORS["text_title"],
                 bold=True)

    add_text_box(slide, "Every icon rendered in Segoe UI Symbol, on the dark deck background",
                 LAYOUT["pad_x_in"], 0.72,
                 12.0, 0.35,
                 font_name=FONTS["body"],
                 size=SIZES["footer"],
                 color=COLORS["text_lead"])

    # Subtle rule
    rule = slide.shapes.add_shape(1,
        Inches(LAYOUT["pad_x_in"]), Inches(1.12),
        Inches(LAYOUT["slide_w_in"] - 2 * LAYOUT["pad_x_in"]), Inches(0.012))
    rule.fill.solid()
    rule.fill.fore_color.rgb = rgb(COLORS["border"])
    rule.line.fill.background()

    # Layout icons in a grid: 4 columns
    names = sorted(ICON_MAP.keys())
    col_count  = 4
    col_w      = (LAYOUT["slide_w_in"] - 2 * LAYOUT["pad_x_in"]) / col_count
    start_y    = 1.25
    row_h      = 0.40
    max_rows   = 13   # fits within 7.5-inch slide height

    for idx, name in enumerate(names[:col_count * max_rows]):
        col = idx % col_count
        row = idx // col_count
        x = LAYOUT["pad_x_in"] + col * col_w
        y = start_y + row * row_h

        # Glyph — Segoe UI Symbol for best Windows coverage
        glyph_box = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(0.45), Inches(row_h - 0.04))
        glyph_box.text_frame.word_wrap = False
        gp = glyph_box.text_frame.paragraphs[0]
        gr = gp.add_run()
        gr.text = get_icon_char(name)
        gr.font.name = "Segoe UI Symbol"
        gr.font.size = Pt(16)
        gr.font.color.rgb = rgb(CHART_SERIES_SINGLE)

        # Name label
        add_text_box(slide, name,
                     x + 0.48, y + 0.04,
                     col_w - 0.52, row_h - 0.04,
                     font_name=FONTS["body"],
                     size=10,
                     color=COLORS["text_lead"])

    # Footer
    add_text_box(slide, f"Total icons: {len(ICON_MAP)}  •  presentations/icons.py",
                 LAYOUT["pad_x_in"], LAYOUT["footer_text_y_in"],
                 12.0, 0.35,
                 font_name=FONTS["body"],
                 size=SIZES["footer"],
                 color=COLORS["text_lead"])

    prs.save(OUT_PPTX)
    print(f"Saved: {OUT_PPTX}")
    return OUT_PPTX


if __name__ == "__main__":
    build()
    print("Open the PPTX in PowerPoint to inspect glyph rendering.")
