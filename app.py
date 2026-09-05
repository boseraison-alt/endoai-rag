"""
Endo AI — Flask Web Server
Wraps the original endo_ai.py engine with a browser UI.
Background threading so long PubMed fetches don't block the page.
"""

import os
import sys
import uuid
import threading
from concurrent.futures import ThreadPoolExecutor, as_completed

import narration
import tempfile
from datetime import datetime
from flask import Flask, request, jsonify, render_template, send_file, session
from dotenv import load_dotenv

# Load .env FIRST so all os.getenv() calls below (and in imported modules) see the keys
# override=True needed because Claude Code pre-sets ANTHROPIC_API_KEY='' in environment
load_dotenv(override=True)

# Force UTF-8 stdout/stderr on Windows so Unicode in print() never raises UnicodeEncodeError.
# Reconfigure first; if that fails or stream is non-standard, re-wrap the underlying buffer.
import io as _io_init
for _name in ('stdout', 'stderr'):
    _stream = getattr(sys, _name, None)
    if _stream is None:
        continue
    try:
        _stream.reconfigure(encoding='utf-8', errors='replace')
    except Exception:
        try:
            buf = getattr(_stream, 'buffer', None)
            if buf is not None:
                setattr(sys, _name,
                        _io_init.TextIOWrapper(buf, encoding='utf-8',
                                                errors='replace', line_buffering=True))
        except Exception:
            pass

# ── TTS backends (prefer OpenAI, fall back to gTTS) ──────
try:
    from openai import OpenAI as _OpenAI
    _oai_tts = _OpenAI(api_key=os.getenv("OPENAI_API_KEY", ""))
    OPENAI_TTS_AVAILABLE = bool(os.getenv("OPENAI_API_KEY"))
except Exception:
    OPENAI_TTS_AVAILABLE = False

try:
    from gtts import gTTS
    GTTS_AVAILABLE = True
except ImportError:
    GTTS_AVAILABLE = False

TTS_AVAILABLE = OPENAI_TTS_AVAILABLE or GTTS_AVAILABLE
if not TTS_AVAILABLE:
    print("Warning: No TTS backend -- set OPENAI_API_KEY or run: pip install gTTS")

# ── moviepy (narrated video) ─────────────────────────────
# Point imageio-ffmpeg (used by moviepy) at the native ffmpeg binary if present.
# A native install (e.g. via `winget install Gyan.FFmpeg`) is several times faster
# than imageio's bundled fallback.
def _find_native_ffmpeg():
    import shutil
    p = shutil.which("ffmpeg")
    if p: return p
    candidates = [
        os.path.expandvars(r"%LOCALAPPDATA%\Microsoft\WinGet\Packages\Gyan.FFmpeg_Microsoft.Winget.Source_8wekyb3d8bbwe\ffmpeg-8.1-full_build\bin\ffmpeg.exe"),
        r"C:\ffmpeg\bin\ffmpeg.exe",
        r"C:\Program Files\ffmpeg\bin\ffmpeg.exe",
    ]
    for c in candidates:
        if os.path.exists(c): return c
    return None

_native_ffmpeg = _find_native_ffmpeg()
if _native_ffmpeg:
    os.environ["IMAGEIO_FFMPEG_EXE"] = _native_ffmpeg
    print(f"  Using native ffmpeg: {_native_ffmpeg}")

try:
    from moviepy import ImageClip, AudioFileClip, concatenate_videoclips
    MOVIEPY_AVAILABLE = True
except ImportError:
    MOVIEPY_AVAILABLE = False
    print("Warning: moviepy not installed -- run: pip install moviepy")

# ── PPTX ─────────────────────────────────────────────────
try:
    from pptx import Presentation as _Prs
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not installed -- run: pip install python-pptx")

from endo_ai import (coverage_groups as endo_ai_coverage_groups,
                     question_coverage as endo_ai_question_coverage,
                     build_evidence_base, ask_clinical_question, ask_learn_question,
                     build_deep_learning_module, StreamAborted,
                     finalise_answer_text, assemble_bibliography,
                     display_title,
                     save_answer, generate_clarifying_questions,
                     classify_question_intent,
                     analyze_radiograph, _analysis_to_prefill,
                     calculate_case_difficulty, match_case_to_profile,
                     generate_referral_letter, generate_podcast_script,
                     generate_audio_script,
                     build_context_block, context_prior_pmids,
                     extract_clinical_recommendation, MAX_CONTEXT_EXCHANGES)
from rag import (setup_query_cache, get_cached_answer, save_query_cache,
                 setup_abstract_cache, get_cached_abstract, cache_abstract,
                 context_fingerprint)

setup_query_cache()
setup_abstract_cache()

# ── Which commit is this process actually running? ───────
# Resolved ONCE, at import time, and never again. That is the whole point:
# `endo-ai-noreload` does not pick up code changes by design, so a process
# can serve a commit that is weeks behind the working tree and nothing about
# it says so. PID 35820 served `grounding-v1` code throughout the whole of
# `grounding-v2` and two batches paid for the confusion.
#
# A REQUEST-time shell-out would report the working tree — i.e. the code on
# disk right now, which is exactly the thing you already know and exactly not
# the thing you are asking about. Import time answers "what did this process
# load?", which is the question.
#
# Failure is not fatal: a deployment from a tarball has no .git, and the app
# must still serve. The field then reads "unknown" rather than being absent,
# so a caller can tell "no git here" from "field not implemented".
def _resolve_git_revision():
    import subprocess
    here = os.path.dirname(os.path.abspath(__file__))
    def _git(*args):
        return subprocess.run(("git", "-C", here) + args,
                              capture_output=True, text=True, timeout=10)
    try:
        r = _git("rev-parse", "--short", "HEAD")
        if r.returncode != 0:
            return "unknown", False
        rev = (r.stdout or "").strip() or "unknown"
        d = _git("status", "--porcelain", "--untracked-files=no")
        dirty = bool((d.stdout or "").strip()) if d.returncode == 0 else False
        return rev, dirty
    except Exception:
        return "unknown", False


GIT_REVISION, GIT_DIRTY = _resolve_git_revision()
IMPORT_TIME = datetime.now().isoformat(timespec="seconds")
print(f"  Serving git {GIT_REVISION}{'+dirty' if GIT_DIRTY else ''} "
      f"(imported {IMPORT_TIME})")

app = Flask(__name__)

# ── Admin authentication ─────────────────────────────────
# Shared-secret gate for operator-only / destructive routes. The token is
# checked at REQUEST time (not import time) so tests and deployments can set
# or rotate ADMIN_TOKEN without restarting differently-configured workers.
#
# Deny by default: if ADMIN_TOKEN is unset, the gated routes return 403 —
# they never fail open. This is bug class (d) in HANDOVER.md (a check that
# fails open) applied to auth.
#
# TWO ways to pass the gate (WORKLIST C4):
#   1. The X-Admin-Token header, matching env ADMIN_TOKEN — unchanged, for
#      curl and the tests.
#   2. A signed admin SESSION, established once via POST /admin/login with
#      that same header. The browser then holds an HttpOnly session cookie
#      and the token itself never appears in page source. (It used to be
#      rendered into a <meta> tag on `/`, where anyone who could load the
#      page could read it.)
#
# The session stores an HMAC fingerprint of the token, keyed on the server's
# secret key — NOT the token or any bare hash of it, because Flask session
# cookies are signed but READABLE client-side, and a readable sha256(token)
# would hand out an offline brute-force target. Recomputing the fingerprint
# per request means rotating ADMIN_TOKEN (or the secret key) invalidates
# every session already issued.
#
# app.secret_key comes from env FLASK_SECRET_KEY and FAILS CLOSED: when it is
# unset there is no signing key, no session can be issued or read, and the
# session path simply never authenticates — the header path still works. A
# hard-coded fallback key would let anyone who has read the source forge the
# cookie.
import hmac as _admin_hmac
import hashlib as _admin_hashlib
from functools import wraps as _admin_wraps

_secret = (os.getenv("FLASK_SECRET_KEY") or "").strip()
if _secret:
    app.secret_key = _secret
else:
    print("Warning: FLASK_SECRET_KEY is not set -- admin sessions disabled "
          "(the X-Admin-Token header still works).")
app.config.setdefault("SESSION_COOKIE_HTTPONLY", True)
app.config.setdefault("SESSION_COOKIE_SAMESITE", "Lax")
del _secret


def _admin_session_fingerprint(token: str) -> str:
    """HMAC(secret_key, token) — what an authenticated session stores."""
    return _admin_hmac.new(app.secret_key.encode("utf-8"),
                           token.encode("utf-8"),
                           _admin_hashlib.sha256).hexdigest()


def _admin_session_valid(expected: str) -> bool:
    """True when this request carries a session whose fingerprint matches the
    CURRENT token. Fails closed: no secret key, no session support, no
    unexpected exception ever authenticates."""
    try:
        if not app.secret_key:
            return False
        stored = session.get("admin_fp") or ""
        if not stored:
            return False
        return _admin_hmac.compare_digest(stored,
                                          _admin_session_fingerprint(expected))
    except Exception:
        return False


def require_admin_token(fn):
    """403 unless the request carries X-Admin-Token matching env ADMIN_TOKEN,
    or an admin session established through POST /admin/login.

    Comparisons are constant-time (hmac.compare_digest) so the token can't be
    recovered byte-by-byte from response timing.
    """
    @_admin_wraps(fn)
    def _admin_guard(*args, **kwargs):
        expected = (os.getenv("ADMIN_TOKEN") or "").strip()
        if not expected:
            return jsonify({
                "error": "Admin routes are disabled: ADMIN_TOKEN is not set "
                         "on the server. Set ADMIN_TOKEN in .env to enable "
                         "them (see README)."
            }), 403
        provided = request.headers.get("X-Admin-Token", "")
        if _admin_hmac.compare_digest(provided.encode("utf-8"),
                                      expected.encode("utf-8")):
            return fn(*args, **kwargs)
        if _admin_session_valid(expected):
            return fn(*args, **kwargs)
        return jsonify({"error": "Invalid or missing X-Admin-Token header, "
                                 "and no admin session."}), 403
    return _admin_guard


@app.route("/admin/login", methods=["POST"])
def admin_login():
    """Trade the admin token (sent ONCE, as a header) for a signed session.

    The browser UI calls this so the token never has to live in page source;
    afterwards the HttpOnly cookie authenticates the gated routes. Deny by
    default on every axis: no ADMIN_TOKEN => 403; no FLASK_SECRET_KEY => 403
    (nothing to sign the cookie with — never issue an unsigned/forgeable
    one); wrong token => 403.
    """
    expected = (os.getenv("ADMIN_TOKEN") or "").strip()
    if not expected:
        return jsonify({
            "error": "Admin routes are disabled: ADMIN_TOKEN is not set "
                     "on the server. Set ADMIN_TOKEN in .env to enable "
                     "them (see README)."
        }), 403
    if not app.secret_key:
        return jsonify({
            "error": "Admin sessions are disabled: FLASK_SECRET_KEY is not "
                     "set on the server, so there is nothing to sign the "
                     "session cookie with. Set it in .env (see README) or "
                     "send X-Admin-Token per request."
        }), 403
    provided = request.headers.get("X-Admin-Token", "")
    if not _admin_hmac.compare_digest(provided.encode("utf-8"),
                                      expected.encode("utf-8")):
        return jsonify({"error": "Invalid or missing X-Admin-Token header."}), 403
    session["admin_fp"] = _admin_session_fingerprint(expected)
    return jsonify({"ok": True})


# ── In-memory job store ──────────────────────────────────
jobs      = {}
jobs_lock = threading.Lock()

# ── Audio job store ───────────────────────────────────────
audio_jobs      = {}
audio_jobs_lock = threading.Lock()

# ── Review-mode conversation memory ──────────────────────
# thread_id → [{question, recommendation, pmids}], oldest first.
#
# Server-side rather than client-side on purpose. The recommendation has to be
# extracted from the answer MARKDOWN, and the "continues from" line the UI
# shows must be a report of what the server actually used — a client that
# believed it was in a thread while the server answered cold would put a
# continuity claim on an answer that has none.
review_threads      = {}
review_threads_lock = threading.Lock()

# A long-lived server accumulates threads from every browser tab that ever
# asked a question; only the last few exchanges of each are ever read.
REVIEW_THREADS_MAX = 500


def _thread_exchanges(thread_id: str) -> list:
    """The carried exchanges for a thread, oldest first. Never more than
    MAX_CONTEXT_EXCHANGES — the cap is applied on write as well as on read, so
    the store cannot grow a long tail nothing will ever look at."""
    if not thread_id:
        return []
    with review_threads_lock:
        return list(review_threads.get(thread_id) or [])


# A21b — a curriculum carries far more evidence than a review answer, and a
# follow-up to one is supposed to be answered over THAT evidence. 12 PMIDs is
# the right carry for a review exchange and far too thin for a curriculum, so
# the cap is per-kind rather than one number pretending to suit both.
THREAD_PMIDS_REVIEW = 12
THREAD_PMIDS_LEARN  = 60


def _thread_record(thread_id: str, question: str, answer: str, papers: list,
                   mode: str = "review") -> None:
    """Append one completed exchange. Recommendation only — see
    endo_ai.extract_clinical_recommendation for why the full answer stays out."""
    if not thread_id:
        return
    from endo_ai import _extract_cited_pmids
    cited = []
    for p in _extract_cited_pmids(answer or ""):
        if p not in cited:
            cited.append(p)
    cap = THREAD_PMIDS_LEARN if mode == "learn" else THREAD_PMIDS_REVIEW
    if len(cited) > cap:
        # Standing rule 5: anything that caps says what it dropped.
        print(f"  [thread] carrying {cap} of {len(cited)} cited paper(s) "
              f"from this {mode} answer into the thread")
    entry = {
        "question":       (question or "").strip(),
        "recommendation": extract_clinical_recommendation(answer or ""),
        "pmids":          cited[:cap],
    }
    if not entry["question"]:
        return
    with review_threads_lock:
        thread = review_threads.setdefault(thread_id, [])
        thread.append(entry)
        del thread[:-MAX_CONTEXT_EXCHANGES]
        if len(review_threads) > REVIEW_THREADS_MAX:
            for stale in list(review_threads)[:len(review_threads) - REVIEW_THREADS_MAX]:
                if stale != thread_id:
                    review_threads.pop(stale, None)


def _thread_clear(thread_id: str) -> None:
    """"New topic" — the thread is gone, and the next question is answered as
    cold as the first one was."""
    with review_threads_lock:
        review_threads.pop(thread_id, None)


def create_job(question: str, mode: str = "review") -> str:
    job_id = str(uuid.uuid4())
    with jobs_lock:
        jobs[job_id] = {
            "id":         job_id,
            "question":   question,
            "mode":       mode,
            "status":     "running",
            "progress":   0,
            "message":    "Starting...",
            # `answer` is the GUARDED text: it is only ever set once
            # validate_evidence_mapping and verify_citation_support have run on
            # the complete synthesis. `partial_answer` is the raw, unchecked
            # stream. The UI derives its trust chips from `answer` alone, so a
            # pass-state chip is structurally unreachable while streaming —
            # this split is the fix for bug class (d) in this path.
            "answer":     None,
            "partial_answer": "",
            "streaming":  False,
            # "pending" until the guardrails have actually run on the finished
            # text; the header chips read "checking…" for as long as it says so.
            "checks_status": "pending",
            "papers":     [],
            "images":     [],
            "cost_usd":   None,
            "error":      None,
            "abort":      False,
            "created_at": datetime.now().isoformat(),
        }
    return job_id


def is_aborted(job_id: str) -> bool:
    with jobs_lock:
        return jobs.get(job_id, {}).get("abort", False)


def update_job(job_id: str, **kwargs):
    with jobs_lock:
        if job_id in jobs:
            jobs[job_id].update(kwargs)


# ── Routes ───────────────────────────────────────────────

@app.route("/")
def index():
    # The admin token is deliberately NOT rendered into the page (it used to
    # be, in a <meta name="admin-token"> tag, where anyone who could load /
    # could read it — WORKLIST C4). The sidebar's admin actions authenticate
    # through POST /admin/login instead: the operator pastes the token once,
    # the server sets a signed HttpOnly session cookie, and page source stays
    # secret-free.
    return render_template("index.html")


@app.route("/tos")
def tos():
    return render_template("tos.html")


@app.route("/health")
def health():
    """Liveness plus the commit this PROCESS loaded.

    `git_revision` is resolved at import and frozen there, so a stale
    no-reload server reports the commit it is actually running rather than
    whatever is checked out now. Compare it against `git rev-parse --short
    HEAD` before trusting anything this server serves or writes.
    """
    return jsonify({
        "status":       "ok",
        "git_revision": GIT_REVISION,
        "git_dirty":    GIT_DIRTY,
        "imported_at":  IMPORT_TIME,
        "pid":          os.getpid(),
    })


@app.route("/clarify", methods=["POST"])
def clarify():
    """Return 2-3 clarifying questions for a clinical query, or [] if none needed."""
    data     = request.json or {}
    question = data.get("question", "").strip()
    if not question:
        return jsonify({"questions": []}), 400
    try:
        questions = generate_clarifying_questions(question)
        return jsonify({"questions": questions})
    except Exception as e:
        return jsonify({"questions": [], "error": str(e)})


@app.route("/ask", methods=["POST"])
def ask():
    data     = request.json or {}
    question = data.get("question", "").strip()
    mode     = data.get("mode", "review")
    context  = data.get("context", "")   # clarification Q&A — set after user answers
    skip_clarify = data.get("skip_clarify", False)  # true when context already provided
    thread_id = (data.get("thread_id") or "").strip()[:64]
    if mode not in ("review", "learn"):
        mode = "review"
    if not question:
        return jsonify({"error": "Question required"}), 400

    # "New topic": drop the thread BEFORE reading it, so this question is
    # answered with no context and shows no continuity line.
    if data.get("new_topic"):
        _thread_clear(thread_id)

    # ── Conversation memory (Review only) ────────────────
    # Deep Learning builds a whole curriculum from one question and has no
    # follow-up turn; carrying a prior recommendation into it would be context
    # nothing asked for.
    exchanges     = _thread_exchanges(thread_id) if mode == "review" else []
    context_block = build_context_block(exchanges)
    prior_pmids   = context_prior_pmids(exchanges)
    continues_from = exchanges[-1]["question"] if (exchanges and context_block) else ""

    # ── Clarify gate ─────────────────────────────────────
    # A20. Literature answers; it does not interview. This gate was the ONLY
    # way the review route could put a question back to the clinician — the
    # synthesis prompt has forbidden questions in the answer body since
    # `trust-surface-v1`, and 0 of 10 stored review answers contain one.
    #
    # A20 (revision), RB 2026-09-03. Curriculum keeps the ability to ask, but
    # narrowed to its one legitimate use: a topic too broad to teach from.
    # `generate_clarifying_questions` asked 2-3 clinical questions on
    # principle, which on "apicoectomy of mandibular teeth" is the same
    # interrogation Literature just stopped doing.
    if mode == "learn" and not skip_clarify and not context:
        try:
            from endo_ai import generate_curriculum_narrowing
            questions = generate_curriculum_narrowing(question)
            if questions:
                return jsonify({"needs_clarification": True, "questions": questions})
        except Exception:
            pass   # Fail open, and open means build.

    # Build enriched question if user answered clarifying questions
    full_question = question
    if context:
        full_question = f"{question}\n\nAdditional clinical context provided by the clinician:\n{context}"

    job_id = create_job(question, mode)
    # The UI's "continues from" line is published with the job, so it states
    # what the server actually did rather than what the browser assumed.
    update_job(job_id, continues_from=continues_from)
    thread = threading.Thread(
        target=run_question,
        args=(job_id, full_question, mode),
        kwargs={"context_block": context_block, "prior_pmids": prior_pmids,
                "thread_id": thread_id},
        daemon=True
    )
    thread.start()
    return jsonify({"job_id": job_id, "continues_from": continues_from})


@app.route("/thread/clear", methods=["POST"])
def thread_clear():
    """"New topic" — forget this thread. Idempotent, and safe to call for a
    thread_id the server has never seen."""
    data = request.json or {}
    _thread_clear((data.get("thread_id") or "").strip()[:64])
    return jsonify({"ok": True})


@app.route("/thread/seed", methods=["POST"])
def thread_seed():
    """Opening a stored answer starts a thread FROM it.

    A21b/A21c. An answer opened out of history is the same thing as one just
    written, as far as a follow-up is concerned — but nothing was recorded for
    it, so without this the follow-up is answered cold AND inherits whatever
    thread the page happened to be on. Both halves are wrong in opposite
    directions, which is why this seeds and clears in one step.

    The client sends a REFERENCE, never content. The server reads the stored
    answer itself, so a page cannot dictate what its own context claims the
    earlier exchange said.
    """
    data      = request.json or {}
    thread_id = (data.get("thread_id") or "").strip()[:64]
    if not thread_id:
        return jsonify({"error": "thread_id required"}), 400
    # Seeding REPLACES. A thread seeded onto an existing one would carry two
    # unrelated topics and answer the follow-up out of both.
    _thread_clear(thread_id)

    learn_file = (data.get("learn_file") or "").strip()
    cache_id   = data.get("cache_id")

    if learn_file:
        if ("/" in learn_file or "\\" in learn_file or ".." in learn_file
                or not learn_file.endswith(".json")):
            return jsonify({"error": "invalid filename"}), 400
        path = os.path.join(_LEARN_HISTORY_DIR, learn_file)
        if not os.path.isfile(path):
            return jsonify({"error": "not found"}), 404
        with open(path, encoding="utf-8") as fh:
            rec = _audit_json.load(fh)
        _thread_record(thread_id, display_title(rec.get("question", "")),
                       rec.get("answer") or "", rec.get("papers") or [],
                       mode="learn")
    elif cache_id is not None:
        from rag import get_conn
        conn = get_conn()
        cur  = conn.cursor()
        try:
            cur.execute("SELECT question_text, answer, papers FROM query_cache "
                        "WHERE id = %s", (int(cache_id),))
            row = cur.fetchone()
        finally:
            cur.close()
            conn.close()
        if not row:
            return jsonify({"error": "not found"}), 404
        qt, answer, papers = row
        row_mode = "review"
        for tag in ("learn", "review", "case"):
            if (qt or "").startswith(f"[{tag}] "):
                row_mode = tag
                qt = qt[len(tag) + 3:]
                break
        _thread_record(thread_id, display_title(qt), answer or "",
                       papers or [], mode=row_mode)
    else:
        return jsonify({"error": "learn_file or cache_id required"}), 400

    carried = _thread_exchanges(thread_id)
    return jsonify({"ok": True,
                    "carried_papers": len(carried[-1]["pmids"]) if carried else 0})


@app.route("/abort/<job_id>", methods=["POST"])
def abort(job_id: str):
    with jobs_lock:
        if job_id in jobs:
            jobs[job_id]["abort"] = True
    return jsonify({"ok": True})


@app.route("/cache/clear", methods=["POST"])
@require_admin_token
def cache_clear():
    """Delete a specific cached answer so it gets regenerated fresh."""
    from rag import get_conn
    data     = request.json or {}
    question = data.get("question", "").strip()
    mode     = data.get("mode", "review")
    if not question:
        return jsonify({"error": "Question required"}), 400

    cache_key = f"[{mode}] {question}"
    from rag import embed
    try:
        q_vec = embed(cache_key)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    conn = get_conn()
    cur  = conn.cursor()
    try:
        cur.execute("""
            DELETE FROM query_cache
            WHERE 1 - (question_embedding <=> %s::vector) >= 0.99;
        """, (q_vec,))
        deleted = cur.rowcount
        conn.commit()
        return jsonify({"deleted": deleted})
    except Exception as e:
        conn.rollback()
        return jsonify({"error": str(e)}), 500
    finally:
        cur.close()
        conn.close()


def _safe_papers(papers: list) -> list:
    """Strip abstract/raw text from papers before sending to browser.
    Users see only metadata (PMID, title, authors, year, journal, score, level).
    Abstract text stays server-side as Claude context only — never exposed to clients.
    """
    # `impact_factor` is deliberately NOT here (invariant 11, Q3, addendum A6).
    # Q3 removed it from every rendered surface, but the field kept leaving the
    # server, so the browser was one line of code away from showing it again.
    # Verified on a live payload: the AAE position statements were still being
    # shipped `impact_factor: 8.0` — a number that cannot exist, since a
    # position statement is not a journal and has no impact factor.
    ALLOWED = {"pmid", "title", "authors", "year", "journal",
               "journal_abbrev", "volume", "issue", "pages",
               "sample_size", "followup_months",
               "citations", "level_key", "score",
               # Provenance badges — metadata only, no abstract text
               "has_coi", "coi_funder", "coi_status", "is_registered", "registry",
               "has_erratum", "has_retraction", "medline_indexed",
               # Superseded rows are filtered out of search, so this is empty in
               # practice today. It carries the PMID of the current version, so
               # the UI can point a clinician at it rather than just hiding the
               # stale one — and the whitelist has silently dropped provenance
               # fields twice before.
               "superseded_by", "is_reference_text"}
    return [{k: v for k, v in p.items() if k in ALLOWED} for p in (papers or [])]


@app.route("/status/<job_id>")
def status(job_id: str):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        return jsonify({"error": "Job not found"}), 404
    # Return a copy with abstracts stripped — never expose raw copyrighted text to client
    safe = dict(job)
    safe["papers"] = _safe_papers(job.get("papers", []))
    # `trust-surface-v1` Q5. The browser's bibliography is the CITATION set,
    # not the retrieval pool, and the split is computed here — from the answer
    # and the papers of this same job — so every consumer of /status gets one
    # answer to "what did this cite?" rather than each deriving its own.
    safe["cited_pmids"] = assemble_bibliography(
        job.get("answer") or "", job.get("papers") or [])["cited_pmids"]
    return jsonify(safe)


# ── Background worker ────────────────────────────────────

def run_question(job_id: str, question: str, mode: str = "review",
                 context_block: str = "", prior_pmids: list = None,
                 thread_id: str = ""):
    try:
        # Cache key includes mode so review/learn answers are stored separately
        cache_key = f"[{mode}] {question}"
        # ...and the conversation context, as a hard partition of the table.
        # A follow-up's TEXT alone is often within 0.92 of the same words asked
        # cold — "what about in immature teeth?" is the same string either way —
        # so without this the first follow-up in a thread would be served the
        # context-free answer. See rag.context_fingerprint.
        ctx_hash = context_fingerprint(context_block)

        # ── Cache check ───────────────────────────────────
        # Deep Learning: 7-day TTL (curricula go stale as literature evolves).
        # Other modes: no age limit on cache hits.
        update_job(job_id, message="Checking answer cache...", progress=3)
        # Review answers expire too. They previously had NO ttl, so a cached
        # clinical recommendation could be served at $0 indefinitely, never
        # re-validated against newer literature — the opposite of the intended
        # ordering, since a point-of-care recommendation is exactly the thing
        # that should go stale.
        ttl = LEARN_HISTORY_TTL_DAYS if mode == "learn" else REVIEW_CACHE_TTL_DAYS
        cached = get_cached_answer(cache_key, max_age_days=ttl,
                                   context_hash=ctx_hash)
        if cached:
            # `trust-surface-v1`: a cached answer is a rendered surface too,
            # and every answer in the cache was written before Q2 and Q3. Same
            # normalisation as a fresh one — one function, both paths.
            cached["answer"], _cq = finalise_answer_text(cached["answer"])
            age_days = cached.get("age_days")
            if mode == "learn" and age_days is not None:
                age_label = "today" if age_days == 0 else f"{age_days}d ago"
                msg = f"Done (from history — saved {age_label})"
            else:
                msg = "Done (from cache)"
            update_job(
                job_id,
                status   = "complete",
                progress = 100,
                message  = msg,
                answer   = cached["answer"],
                papers   = cached["papers"],
                images   = [],
                cost_usd = 0.0,
                # A cached answer was guarded before it was stored, and the
                # chips are read back off the markers inside it — including the
                # explicit "not available" case.
                streaming     = False,
                checks_status = "complete",
            )
            # A cache hit is still an exchange: the clinician asked, and an
            # answer was given. If it did not join the thread, the NEXT
            # follow-up would reach back past it to a stale one.
            if mode in ("review", "learn"):
                _thread_record(thread_id, question, cached["answer"],
                               cached.get("papers") or [], mode=mode)
            return

        # ── Intent routing (Haiku triage) ─────────────────
        # Runs ahead of retrieval so we can pick a cheaper pipeline
        # for trivial questions and tell the clinician what we're doing.
        update_job(job_id, message="Routing question...", progress=2)
        intent = classify_question_intent(question, context_block=context_block)
        intent_cost = float(intent.get("cost") or 0.0)
        print(f"[intent] kind={intent['kind']} retrieval={intent['retrieval']} "
              f"clarify={intent['needs_clarify']} reason={intent['reason']}")

        # ── Full pipeline ─────────────────────────────────
        if mode == "learn":
            # Agentic curriculum builder: syllabus → per-module retrieval →
            # per-module writing → stitch. Density-capped at ~20 minutes.
            update_job(job_id, message="Planning teaching curriculum...", progress=5)

            def _learn_progress(pct: int, msg: str):
                if is_aborted(job_id):
                    raise RuntimeError("Cancelled by user")
                update_job(job_id, message=msg, progress=pct)

            try:
                answer, cost, evidence = build_deep_learning_module(
                    question, progress_cb=_learn_progress
                )
            except RuntimeError as ce:
                if "Cancelled" in str(ce):
                    update_job(job_id, status="aborted", progress=100, message="Cancelled")
                    return
                raise
        else:
            kind_label = {"simple": "definition", "standard": "clinical question",
                          "complex": "complex multi-system question"}.get(intent["kind"], "question")
            update_job(job_id, message=f"Routing as {kind_label} — generating search terms...",
                       progress=5)
            evidence = build_evidence_base_with_progress(
                job_id, question, mode=mode,
                context_block=context_block, prior_pmids=prior_pmids)

            if is_aborted(job_id):
                update_job(job_id, status="aborted", progress=100, message="Cancelled")
                return

            # Publish the retrieved papers BEFORE synthesis so the inline
            # [[PMID:N]] pills rendered mid-stream can resolve to real author
            # names instead of bare numbers. `_safe_papers` still strips the
            # abstracts on the way out of /status.
            update_job(job_id,
                       message   = "Asking Claude to synthesize the evidence...",
                       progress  = 80,
                       papers    = evidence.get("_summary", {}).get("all_scored", []),
                       streaming = True,
                       partial_answer = "",
                       checks_status  = "pending")

            def _on_partial(text: str):
                # Raw, UNCHECKED model text. It goes to `partial_answer` — never
                # to `answer` — so nothing downstream can mistake it for a
                # validated result.
                update_job(job_id,
                           partial_answer = text,
                           streaming      = True,
                           checks_status  = "pending",
                           message        = "Writing the answer…",
                           progress       = min(95, 80 + len(text) // 400))

            def _on_phase(label: str):
                # The model has stopped writing but the guardrails have NOT
                # finished. checks_status stays "pending" — the chips keep
                # saying "checking…" for exactly this window.
                update_job(job_id, streaming=False, checks_status="pending",
                           progress=97,
                           message="Checking citations against the abstracts…")

            try:
                answer, cost = ask_clinical_question(
                    question, evidence,
                    stream_cb = _on_partial,
                    abort_cb  = lambda: is_aborted(job_id),
                    phase_cb  = _on_phase,
                    context_block = context_block,
                )
            except StreamAborted:
                update_job(job_id, status="aborted", progress=100,
                           message="Cancelled", streaming=False, partial_answer="")
                return

        cost = float(cost or 0.0) + intent_cost
        images = []

        if is_aborted(job_id):
            update_job(job_id, status="aborted", progress=100, message="Cancelled")
            return

        # Pull top papers for display
        summary = evidence.get("_summary", {})
        papers  = summary.get("all_scored", [])

        save_answer(question, answer, evidence)
        save_query_cache(cache_key, answer, papers, context_hash=ctx_hash)
        write_citation_audit(question, answer, mode)
        # A21b — a curriculum joins the thread too, so a follow-up to it is
        # answered over the evidence it was built from rather than rebuilding
        # it. The recorder stores the extracted recommendation and the cited
        # PMIDs, never the 12,000 words.
        if mode in ("review", "learn"):
            _thread_record(thread_id, question, answer, papers, mode=mode)

        # Deep Learning curricula get an additional persistent file archive
        # under learn_history/. The 7-day re-use window is enforced via the
        # query_cache age filter above; this folder is the durable record.
        if mode == "learn":
            # A15f.1 — the record's `question` is a TITLE, so it stores the
            # clinician's own words rather than the question plus the
            # clarification transcript appended to it.
            save_learn_output(display_title(question), answer, evidence, cost)

        update_job(
            job_id,
            status   = "complete",
            progress = 100,
            message  = "Done",
            answer   = answer,
            papers   = papers,
            images   = images,
            cost_usd = round(cost, 4),
            # Guardrails have now run on the complete text — and only now may
            # the header chips show a real pass/fail state.
            streaming      = False,
            partial_answer = "",
            checks_status  = "complete",
        )

    except Exception as e:
        update_job(job_id, status="error", progress=100, error=str(e), message=str(e),
                   streaming=False, checks_status="pending")


def multi_query_search(question: str, generated_terms: list, limit: int = 100) -> list:
    """KNN the library once per query string and union the results, keeping the
    best similarity seen for each PMID.

    Retrieval used to depend on a single generated boolean, and that coupled two
    things that pull in opposite directions. A well-formed PubMed query is
    mostly operators, quotes and truncation asterisks, so it embeds FURTHER from
    a paper's prose than a sloppy one does. Measured on Cochrane CD005296
    (PMID 36512807) for "single-visit versus multiple-visit root canal
    treatment":

        raw clinician question                      0.680   rank 20   kept
        bag-of-words generated query                0.585   rank 19   kept
        3-group query, best spec compliance         0.546   rank 11   CUT

    The review was rank 11 in the entire library for the query that missed it —
    the KNN found it and the absolute floor removed it. The better the boolean,
    the worse it scored.

    Embedding the clinician's ORIGINAL question alongside every generated term
    and taking the max removes the coupling: the question is prose and matches
    prose, while the booleans keep contributing their own recall for the
    vocabulary the question does not use.
    """
    from rag import search as _search

    queries, seen_q = [], set()
    for q in [question] + list(generated_terms or []):
        q = (q or "").strip()
        if q and q not in seen_q:
            seen_q.add(q)
            queries.append(q)

    best = {}
    for q in queries:
        try:
            for row in _search(q, level_key=None, limit=limit) or []:
                pmid = row.get("pmid")
                if not pmid:
                    continue
                sim = float(row.get("similarity") or 0)
                prev = best.get(pmid)
                if prev is None or sim > float(prev.get("similarity") or 0):
                    row = dict(row)
                    row["similarity"] = sim
                    best[pmid] = row
        except Exception as e:
            # One bad query must not lose the recall of the others.
            print(f"  [rag] query failed, continuing: {e}")

    out = sorted(best.values(), key=lambda r: -float(r.get("similarity") or 0))
    print(f"  [rag] {len(queries)} quer{'y' if len(queries) == 1 else 'ies'} "
          f"-> {len(out)} distinct papers (best similarity kept per PMID)")
    return out[:limit * 2]


def apply_evidence_floor(relevant: list, floor: float = None,
                         min_papers: int = None) -> list:
    """Which of the routing-admitted papers the MODEL actually reads.

    A42. `relevant` decided the ROUTE — is the library's coverage good enough to
    answer from. This decides what goes in the prompt, and it is a different
    question with a different answer: 18% of that set sits between 0.55 and 0.60
    and is cited 1.1% of the time (181 cited paper-instances, ten A38d runs).

    Two rules, and the second exists because the first was generalised from an
    unrepresentative sample. Measured on the five A38d questions the floor looks
    uniformly cheap; measured on all 29 it guts the thin ones —
    `case-opening-sparse` 103 papers to 6, `pregnancy` 100 to 12. A pool of 6
    manufactures the false evidence gap A5 was about, so below `min_papers` the
    floor does not cut at all and the most similar N are kept instead. It can
    only ever ADD papers back, and it selects by similarity because which papers
    survive is a relevance question (standing rule 19).

    A separate function rather than four lines inline so a test can exercise the
    expression production evaluates rather than a restatement of it (rule 14) —
    two mutations survived while this was inline.
    """
    floor = RELEVANCE_GATE["evidence_floor"] if floor is None else floor
    min_papers = (RELEVANCE_GATE["min_evidence_papers"]
                  if min_papers is None else min_papers)
    by_sim = sorted(relevant or [],
                    key=lambda r: -float(r.get("similarity") or 0))
    kept = [r for r in by_sim if float(r.get("similarity") or 0) >= floor]

    if len(kept) < min_papers:
        topped = by_sim[:min_papers]
        # Standing rule 5, and this is the branch a reader most needs to see:
        # the pool it reports is NOT what the floor would have given.
        print(f"  [evidence_floor] only {len(kept)} of {len(by_sim)} clear "
              f"{floor}; keeping the {len(topped)} most similar instead "
              f"(min_evidence_papers={min_papers})")
        return topped
    if len(kept) < len(by_sim):
        print(f"  [evidence_floor] {len(by_sim)} cleared the routing floor; "
              f"{len(kept)} clear the evidence floor {floor}; dropped "
              f"{len(by_sim) - len(kept)} (best dropped "
              f"{float(by_sim[len(kept)].get('similarity') or 0):.3f})")
    return kept


def cap_by_relevance(bucket: list, cap: int, tier: str = "") -> list:
    """Choose which papers in a tier survive the cap, by relevance.

    A5b. WHICH papers survive is a relevance question; the ORDER they are then
    shown in is a quality question. This used to answer both with the score,
    and the score does not know what was asked.

    Measured on "retreatment in one visit versus two visits": 60 level1 papers
    cleared the 0.55 similarity floor, the cap kept 25 of them BY SCORE, and
    Karaoglan 2022 — the single most on-point RCT in the library, similarity
    0.648 — ranked 54th of 60 by score and was cut. Twenty of the twenty-five
    it kept were LESS similar to the question than the one it dropped, led by
    AAE and ESE position statements at score 90.0 and 87.0 and similarities of
    0.62 and 0.60 (which is also A7: guideline rows sitting in level1 at a
    score no trial in the library reaches). The answer then declared that no
    prospective study directly compares the two protocols — A5's false evidence
    gap, produced here.

    Invariant 1 is untouched: tier is still assigned by study design, and the
    caller still orders within the tier by score. This decides membership, not
    rank.

    Ties break on score, so two equally relevant papers are separated by the
    better one — and the order is deterministic, which a bare sort on a float
    that repeats is not.
    """
    if cap is None or cap <= 0 or len(bucket) <= cap:
        return list(bucket)

    def relevance(p):
        return (float(p.get("similarity") or 0), float(p.get("score") or 0))

    ranked  = sorted(bucket, key=relevance, reverse=True)
    kept    = ranked[:cap]
    dropped = ranked[cap:]
    # Standing rule 5: anything that discards candidate content says what it
    # discarded and how close the closest one was. A silent cap is how the
    # retreatment RCT disappeared without a trace in the first place.
    print(f"  [cap] {tier or 'tier'}: {len(bucket)} above the floor, keeping "
          f"the {cap} most relevant; dropped {len(dropped)}, best dropped "
          f"similarity {float(dropped[0].get('similarity') or 0):.3f} "
          f"(PMID {dropped[0].get('pmid')})")
    return kept


# A32 — `ensure_authoritative` was deleted here, deliberately.
#
# It was a backstop against a top-tier paper being lost to query variance, and
# it had never once fired: `usable()` required similarity at or above the
# floor, and the `relevant` list it was handed already contained every such
# candidate, so its re-inclusion set was empty by construction. Its own log
# line has never printed. Its unit tests passed only because they called it
# with `relevant=[]`, a state the production path never produces — and one of
# them, `test_cochrane_below_the_floor_is_reinstated`, used floor=0.50 against
# a similarity of 0.546, so it asserted the opposite of its own name.
#
# RB's decision was that it must NOT be fixed by reaching below the similarity
# floor: on apicoectomy 183 of 200 candidates sit below it, and admitting
# high-tier papers from there is the error just removed from the cap and the
# ORDER BY, wearing a virtuous hat. A Cochrane review about a different
# question is still about a different question.
#
# The narrower redefinition (restore a paper that clears the floor on ANY
# query but is lost in the union merge or a cap) has nothing left to protect:
#
#   * `multi_query_search` keeps the MAX similarity per PMID across every
#     query, so one bad embedding cannot lose a paper another query found.
#     Measured on the retreatment question with a consistent query set: of the
#     papers clearing the floor on at least one of 8 queries, ZERO failed to
#     reach the merged candidate set, which cuts at similarity 0.558 against a
#     floor of 0.55.
#   * the caps below it now order by relevance (A5b, A30b), so a paper cut
#     there is cut because more relevant papers exist in its tier. Restoring it
#     would be authority overriding relevance — the thing this is not allowed
#     to be.
#
# The protection is real and lives in two layers that do work: the union-of-max
# above, and the eval's `must_include_pmid` on 36512807. A guarantee that
# cannot fire is worse than none, because it gets described — and it was, in
# three handover files.


# ── B2/B5 retrieval concurrency ──────────────────────────
# Bounded by the NCBI rate limit, not by CPU: at 9 req/sec with an API key and
# 2 HTTP calls per fetch, ~6 in flight keeps the limiter saturated without
# queueing work that can only wait. Also stays under DB_POOL_MAX (10).
TIER_FETCH_WORKERS = 6

# B5 early stop, Review mode only.
EARLY_STOP_TIERS      = ("level1",)   # cochrane is fetched before this loop
EARLY_STOP_MIN_PAPERS = 15


# ── Library relevance gate ───────────────────────────────
# The floor and the count that interprets it are ONE setting, kept together
# deliberately. See HANDOVER.md "The similarity floor asks a different question
# than you think" for the measurement behind these numbers.
#
# similarity_floor 0.45 -> 0.55 on 2026-08-30. all-MiniLM-L6-v2 scores any two
# endodontic texts around 0.45 on shared domain vocabulary alone, so at 0.45
# "how many hits clear the floor" answers "is this endodontics?" rather than
# "is this the question?". Measured across the 20 eval questions, 0.45 routed
# 18/20 to the library — including "root canal treatment in pregnancy", 63
# hits above the floor and not one on-topic paper.
#
# Raising the floor without re-reading min_relevant would have been the same
# mistake in the other direction: 12 papers above 0.55 is a much stronger
# claim than 12 above 0.45, and the pair has to be tuned as a pair.
RELEVANCE_GATE = {
    "similarity_floor": 0.55,   # cosine; below this a hit is same-specialty noise
    # A42. The ROUTING floor above and the EVIDENCE floor below answer two
    # different questions, and separating them is what makes raising one a
    # single-variable change.
    #
    #   routing   "does the library cover this question well enough to answer
    #             from?" — read together with min_relevant, as the note above
    #             insists, and NOT moved here.
    #   evidence  "is this particular paper worth putting in front of the
    #             model?" — a per-paper relevance decision, standing rule 19.
    #
    # Measured on 181 cited paper-instances across the ten A38d runs, scored on
    # the similarity the floor actually sees (the MAX across the question and
    # every generated term, which is what multi_query_search keeps — a first
    # pass used question-only similarity, reported a lowest cited value of
    # 0.4633, i.e. below a floor that demonstrably works, and would have handed
    # RB a worse trade than the real one):
    #
    #   lowest cited similarity 0.587   p05 0.628   p10 0.637   median 0.713
    #
    #   floor   pool cut   cited cut          per question, cited below 0.60
    #   0.55          0%    0  (0.0%)         mta-vs-biodentine    0 of 29
    #   0.58          9%    0  (0.0%)         regenerative         0 of 45
    #   0.60         18%    2  (1.1%)  <--    single-vs-multiple   0 of 38
    #   0.62         33%    8  (4.4%)         cbct-vs-periapical   1 of 38
    #   0.65         51%   29 (16.0%)         direct-pulp-capping  1 of 31
    #
    # 18% of the pool is carried into every synthesis prompt and cited 1.1% of
    # the time. That is cost, not evidence.
    #
    # IT IS THE COST FIX, and the estimate that said otherwise was wrong.
    # Pricing context tokens alone at Opus input rates put the saving at $0.15
    # and called this free-but-worthless. Measured on A42b's paired design it
    # is $1.08 — cost $2.67 -> $1.59, a 40% fall, while the cited count held at
    # 18.0 -> 19.0. Seven times the estimate, because the estimate missed
    # everything downstream of the pool: fewer papers means fewer
    # claim-citation pairs for the support checker and fewer retries. Standing
    # rule 1 — the instruction was to measure, and estimating instead produced
    # the wrong answer.
    #
    # How much is really the floor: the two biggest savings are on runs that
    # cost $3.98 and $3.74 at 0.55 and look retry-inflated. Excluding them the
    # saving is $0.30. So $0.30-$1.08 is the floor and the rest is variance;
    # n=5, one run per condition, four of five cheaper.
    #
    # And raising the ROUTING floor instead would push 2 of 29 eval questions
    # onto the LIVE route, which costs more, not less — which is exactly why
    # these are two constants.
    "evidence_floor":   0.60,
    # A FLOOR NEEDS A FLOOR. The table above was measured on the five A38d
    # questions, and those turn out to be the ones where the library is DEEP —
    # they lose 0-40% of their pool. Across all 29 eval questions the same
    # floor is wildly uneven, and it guts exactly the questions that can least
    # afford it:
    #
    #   case-opening-sparse                 103 papers ->   6
    #   dens-evaginatus-premolar-diagnostic  56 papers ->   6
    #   pregnancy                           100 papers ->  12
    #   dens-invaginatus                    135 papers ->  27
    #   sonic-vs-ultrasonic                 115 papers ->  28
    #
    # A pool of 6 cannot support an answer; it manufactures the false evidence
    # gap A5 was about. Generalising "the floor is free" from five deep-pool
    # questions to all 29 was the error, and this is the correction rather
    # than a hedge.
    #
    # So the floor only ever removes SURPLUS: below this many papers it does
    # not cut at all, and the most similar N are kept regardless. On a thin
    # pool there was no cost to save in the first place.
    #
    #   min keep   total pool   vs 0.55   smallest   questions under 20
    #   none             2076       63%          6                    3
    #   40               2204       67%         40                    0
    #
    # The guard costs 4 points of context saving (37% -> 33%) and cannot cost
    # a citation, because it only ever ADDS papers back. 40 rather than 24:
    # A33j's arithmetic floor is ~24 (the smallest pool that could yield ~20
    # references at the best citation rate ever observed, 82%), and 40 leaves
    # headroom above a bound that assumed a rate never seen twice.
    "min_evidence_papers": 40,
    "min_relevant":     12,     # hits that must clear the floor to serve locally
    "min_hits":         20,     # raw KNN hits before relevance is even considered
    "max_topic_age_yr":  3,     # newest on-topic paper older than this -> go live
    # DOES NOT mirror the live path, and said it did for long enough that the
    # claim is why nobody checked. MEASURED 2026-09-05
    # (scripts/measure_library_route_floor.py):
    #
    #   library route     flat 25 per tier, and NO per-tier quality floor at
    #                     all -- _apply_quality_threshold, _tier_floor,
    #                     _tier_cap and MODE_TIER_QUOTAS have no caller in
    #                     this module
    #   live/curriculum   MODE_TIER_QUOTAS: level1 18, level4/level5/guideline 4,
    #                     plus TIER_QUALITY_FLOORS per tier
    #
    # Consequence, measured against the live library: 516 of 3,346 rows
    # (15.4%) sit below their own tier's floor and are served anyway,
    # including 42 papers scoring 40.4-49.9 rendered under "Level I -- RCTs
    # and Systematic Reviews". The weak tiers can contribute 25 each against
    # level1's 25.
    #
    # NOT FIXED HERE, and now for a MEASURED reason rather than a cautious one.
    # The before/after was run across all 29 eval questions forced onto this
    # route (scripts/measure_library_floor_29.py, report
    # eval/reports/library_floor_29.md). Aligning the route would take 8 of 29
    # questions below min_evidence_papers 40 -- six of them newly -- against a
    # pre-declared stop threshold of 5.
    #
    # WHY, and this is the part worth carrying: the two guards sit on
    # DIFFERENT AXES and only one has a rescue. apply_evidence_floor runs on
    # SIMILARITY, before banding, and tops a thin pool back up to 40. A quality
    # floor would run on SCORE, after banding, on the pool that rescue just
    # produced -- and nothing tops it up again. On the sparse diagnostic
    # questions the sequence is 40 -> rescued to 40 -> cut to 26, with no guard
    # noticing, because min_evidence_papers was satisfied upstream by a
    # different measurement. It would cut hardest where the corpus is thinnest,
    # which is where A5's false evidence gap gets manufactured.
    #
    # Each half breaches on its own (rule 22): the MODE_TIER_QUOTAS cap alone
    # puts 7 questions under 40, the floor alone puts 7 under. Neither "just
    # the cap" nor "just the floor" is the conservative option.
    #
    # And the trap is now a number: a floor reading the coalesced 0.0 that
    # rag_results_to_scored returns for a NULL score cuts 44 of the 55 served
    # guideline paper-instances across the 29 questions. Any fix needs the
    # NULL-score exemption first.
    "max_per_tier":     25,
    # A1a. Every condition above is a question about the CORPUS — enough hits,
    # enough of them similar, at least one high tier, not stale. All four are
    # satisfiable by the endodontic HALF of a two-part question, which is how
    # "eliquis in patients who needs apicectomy" scored as well covered while
    # not one retrieved paper mentioned anticoagulation.
    #
    # This one is a question about the QUESTION: each discriminating concept in
    # the generated query must be REPRESENTED in the candidate set. Chosen from
    # the measured distribution in `eval/reports/a1_coverage_gate.md`, not from
    # a cost target (A1c).
    "min_concept_papers": 3,
}


def build_evidence_base_with_progress(job_id: str, question: str,
                                      force_route: str = None,
                                      mode: str = "review",
                                      context_block: str = "",
                                      prior_pmids: list = None) -> dict:
    """
    RAG-first evidence pipeline.
    Searches the full library without level_key filter (level_key is empty
    in the current library build). Falls back to PubMed if < MIN_RAG_RESULTS.

    `mode` is "review", "learn" or "case". It gates the B5 early stop only:
    Review stops once the top tiers have supplied enough evidence, because tier
    banding means a case series cannot override a Level I finding anyway. Learn
    mode always sweeps every tier — a teaching curriculum genuinely wants the
    narrative scaffolding that reviews and editorials provide. Case mode sweeps
    too, for a different reason: an unusual presentation may have no Level I
    literature at all, and the case series the early stop would have skipped is
    then the only evidence that exists.

    `context_block` / `prior_pmids` carry a Review thread's earlier exchanges.
    The block reaches the two search-term generators only — an elliptical
    follow-up ("what about in immature teeth?") cannot be turned into a query
    without it. `prior_pmids` seeds the CANDIDATE set after the routing gate has
    already decided; see the seeding block below for why that ordering is the
    whole safety property.

    force_route pins retrieval for evaluation runs: "live" skips the library
    gate entirely, "library" refuses to fall back to PubMed. Production always
    passes None and lets the gate decide.

    This exists because write-back makes the eval set decay silently. The laser
    regression was a failure of live search-term generation; once that run wrote
    196 papers back, the same question began serving from the library, and the
    eval case stopped exercising the generator that broke. A future "3-7 word"
    regression would have passed. Route has to be pinned by the harness, not
    left to whatever the library happens to hold on the day.
    """
    from endo_ai import (
        generate_search_terms, generate_multi_search_terms,
        fetch_cochrane, fetch_papers,
        tier_query_lanes, fetch_untyped_recent, PROVISIONAL_KEY,
        _tier_cap, TIER_FETCH_DEPTH,
        COCHRANE_TERM, LEVEL_1_TERMS, LEVEL_2_TERMS,
        LEVEL_3A_TERMS, LEVEL_3B_TERMS,
        LEVEL_4_TERMS, LEVEL_5_TERMS,
        detect_outliers, apply_currency_tags,
        build_synthesis_order, TIER_LABEL, TIER_ORDER,
        flag_superseded_by_review, _pubmed_audit_log,
        label_and_expand,
    )
    from rag import (search as rag_search, rag_results_to_scored, library_stats,
                     search_by_pmids as rag_search_by_pmids)

    # Cosine KNN always returns its nearest neighbours, relevant or not, so a
    # bare count gate ("did we get 20 hits?") is really asking "does the
    # library exist?" — it passes for every question against a 1,886-paper
    # library and the network is then almost never consulted. These thresholds
    # make the gate ask whether the library actually COVERS the question.
    # All five read from RELEVANCE_GATE at module level — see the note there
    # on why the floor and the count that interprets it must move together.
    MIN_RAG_RESULTS         = RELEVANCE_GATE["min_hits"]
    RAG_SIMILARITY_FLOOR    = RELEVANCE_GATE["similarity_floor"]
    MIN_RAG_RELEVANT        = RELEVANCE_GATE["min_relevant"]
    MAX_RAG_PAPERS_PER_TIER = RELEVANCE_GATE["max_per_tier"]
    RAG_MAX_TOPIC_AGE_YEARS = RELEVANCE_GATE["max_topic_age_yr"]
    MIN_CONCEPT_PAPERS      = RELEVANCE_GATE["min_concept_papers"]
    evidence        = {}
    all_scored      = []

    if force_route not in (None, "live", "library"):
        raise ValueError(f"force_route must be None, 'live' or 'library', got {force_route!r}")

    # Check if library is populated
    try:
        stats      = library_stats()
        library_ok = stats["total"] >= 50
    except Exception:
        library_ok = False

    if force_route == "live":
        library_ok = False
        print("  [rag_gate] force_route=live — library skipped")

    update_job(job_id, message="Generating smart search terms...", progress=8)
    smart_topic = generate_search_terms(question, context_block=context_block)

    # ── Try RAG for full evidence base ────────────────────
    if library_ok:
        update_job(job_id, message="Searching local library...", progress=15)
        # Search without level_key filter — library stores all levels together
        # A4: never let one generated string decide what the library returns.
        # generate_multi_search_terms is only called on the live path below, so
        # produce the extra angles here too — cheap relative to a wrong answer.
        try:
            _terms = generate_multi_search_terms(question, smart_topic,
                                                 context_block=context_block)
        except Exception as _te:
            print(f"  [rag] multi-term generation failed, using primary only: {_te}")
            _terms = [smart_topic]
        rag_results = multi_query_search(question, _terms, limit=100)

        # Coverage test, not just a count: enough genuinely-similar papers, and
        # at least one high-tier design among them. A library that answers with
        # only case reports and narrative reviews should defer to live PubMed
        # even when it returns plenty of near neighbours.
        relevant = [r for r in rag_results
                    if float(r.get("similarity") or 0) >= RAG_SIMILARITY_FLOOR]
        has_high_tier = any((r.get("level_key") or "") in ("cochrane", "level1")
                            for r in relevant)
        newest_year = max((int(r["year"]) for r in relevant
                           if str(r.get("year", "")).isdigit()), default=0)

        # Staleness escape hatch. Write-back plus a library-first gate means a
        # topic gets searched live once and is then served from that single
        # search forever. If the freshest paper the library holds on this topic
        # is older than the cutoff, go live regardless of coverage and let
        # write-back refresh the topic.
        from datetime import datetime as _dt
        topic_age = _dt.now().year - newest_year if newest_year else 99
        topic_is_stale = topic_age > RAG_MAX_TOPIC_AGE_YEARS

        # A1a — does the retrieved set address the QUESTION?
        #
        # `coverage_groups` returns the top-level AND-groups of the primary
        # generated query with the generic ones dropped, i.e. the concepts the
        # query itself treats as hard requirements. If a concept is represented
        # by fewer than MIN_CONCEPT_PAPERS candidates, the library has not
        # covered the question and no similarity count may say otherwise.
        #
        # A question whose query is ALL generic vocabulary has no discriminating
        # concept to test; that is a plain endodontic question and the library
        # is the right route, so the condition abstains rather than blocking.
        _cov_groups = endo_ai_coverage_groups(smart_topic)
        coverage    = endo_ai_question_coverage(_cov_groups, relevant)
        weakest_cov = min([c["hits"] for c in coverage], default=None)
        covers_concepts = (weakest_cov is None) or (weakest_cov >= MIN_CONCEPT_PAPERS)

        library_covers_question = (
            len(rag_results) >= MIN_RAG_RESULTS
            and len(relevant) >= MIN_RAG_RELEVANT
            and has_high_tier
            and not topic_is_stale
            and covers_concepts
        ) if force_route != "library" else True
        # force_route="library" holds the library path even when coverage is
        # thin, so a library-mode eval case measures what the library actually
        # returns instead of quietly becoming a live-path case.
        # A1b. Every condition's own verdict, and the coverage terms with their
        # hit counts. A gate that short-circuits live retrieval silently is
        # standing rule §1.5 — it discards the entire live candidate pool and
        # says nothing about having done so.
        def _v(ok):
            return "PASS" if ok else "FAIL"
        print(f"  [rag_gate] hits={len(rag_results)}>={MIN_RAG_RESULTS} "
              f"{_v(len(rag_results) >= MIN_RAG_RESULTS)} | "
              f"relevant={len(relevant)}>={MIN_RAG_RELEVANT} "
              f"{_v(len(relevant) >= MIN_RAG_RELEVANT)} | "
              f"high_tier={has_high_tier} {_v(has_high_tier)} | "
              f"newest={newest_year} age={topic_age}y<={RAG_MAX_TOPIC_AGE_YEARS} "
              f"{_v(not topic_is_stale)} | "
              f"concepts>={MIN_CONCEPT_PAPERS} {_v(covers_concepts)}")
        for _c in coverage:
            print(f"    [rag_gate:coverage] {_c['hits']:>4} paper(s) mention "
                  f"{_c['terms'][:4]}"
                  + ("" if _c["hits"] >= MIN_CONCEPT_PAPERS else "   <-- NOT COVERED"))
        if not _cov_groups:
            print("    [rag_gate:coverage] no discriminating concept in the query "
                  "— condition abstains")
        print(f"  [rag_gate] -> {'LIBRARY' if library_covers_question else 'LIVE PUBMED'}")

        if library_covers_question:
            # ── Prior-exchange seeding (Review conversation memory) ──
            # The papers the previous answer cited are added as CANDIDATES, with
            # a similarity recomputed against THIS question, and then judged by
            # every gate below exactly as a KNN hit is: the similarity floor
            # cuts the ones the follow-up moved away from, banding puts them in
            # their own tier, and rag.search_by_pmids applies the same
            # retracted/withdrawn/superseded exclusion as rag.search.
            #
            # It runs HERE, after library_covers_question has been computed,
            # and that ordering is load-bearing. Seeding before the gate would
            # let three carried papers push `len(relevant)` over min_relevant
            # and route a thin topic to the library on the strength of the
            # PREVIOUS question's evidence — context substituting for
            # retrieval, which is the one thing this feature must not do.
            if prior_pmids:
                known = {r.get("pmid") for r in rag_results}
                try:
                    seeds = rag_search_by_pmids(question,
                                                [p for p in prior_pmids if p not in known])
                except Exception as _se:
                    print(f"  [context] prior-PMID seeding failed, continuing: {_se}")
                    seeds = []
                if seeds:
                    rag_results = rag_results + seeds
                    kept = [s for s in seeds
                            if float(s.get("similarity") or 0) >= RAG_SIMILARITY_FLOOR]
                    relevant = relevant + kept
                    print(f"  [context] seeded {len(seeds)} paper(s) cited by the "
                          f"earlier answer; {len(kept)} cleared the "
                          f"{RAG_SIMILARITY_FLOOR} floor for this question")

            update_job(job_id, message=f"Found {len(rag_results)} papers in library — building evidence...", progress=40)
            # Build the evidence from the RELEVANT hits only. Feeding all 100
            # nearest neighbours put topically unrelated papers in front of
            # Claude — a question about regenerative endodontics was answered
            # citing papers on apex locators and sealer heat properties, which
            # the claim-support check then correctly flagged. The similarity
            # floor has to filter the evidence, not merely decide the gate.
            #
            # A42 — and it filters at the EVIDENCE floor, which is above the
            # routing floor. `relevant` decided the route; this decides what the
            # model reads. Standing rule 5: say what was dropped.
            all_rag = rag_results_to_scored(apply_evidence_floor(relevant))

            # Band by STUDY DESIGN (level_key), rank by score WITHIN each band.
            #
            # This previously banded by score alone (>=70 -> cochrane/level1,
            # 50-70 -> level2/3, <50 -> level4/5), which inverted the product's
            # central guarantee: a well-cited recent case series scoring 72 was
            # handed to Claude labelled "Level I — RCTs and Systematic Reviews",
            # while a smaller Cochrane review scoring 58 was demoted to Level
            # II/III — and the system prompt instructs Claude to trust the tier
            # label absolutely. Score must rank papers within a tier and never
            # promote one across tiers.
            #
            # The score-banding was a workaround for 37% of the library having
            # no level_key; that has since been backfilled from PubMed
            # publication types. The backfill left 2 rows unlabelled, but this
            # is not a closed problem — live write-back keeps producing them
            # (14 as of 2026-08-29), so the fallback below is load-bearing
            # rather than a leftover. test_end_to_end.py pins it.
            by_tier = {}
            for p in all_rag:
                tier = (p.get("level_key") or "").strip()
                # Retracted rows are excluded by search() already; this is the
                # second lock. Without it the unlabelled-fallback below would
                # re-band a 'retracted' tier to level5 and hand it to Claude.
                if tier == "retracted" or p.get("has_retraction"):
                    continue
                # An unlabelled paper has an UNKNOWN design. Placing it in the
                # weakest tier is the safe direction: it can still inform the
                # answer but can never masquerade as high-tier evidence.
                if tier not in TIER_ORDER:
                    tier = "level5"
                by_tier.setdefault(tier, []).append(p)

            # A5b — WHICH papers survive the cap is a relevance question; the
            # ORDER they are shown in is a quality question. This used to
            # answer both with the score, and the score does not know what was
            # asked.
            #
            # Measured on "retreatment in one visit versus two visits":
            # 60 level1 papers cleared the 0.55 floor, the cap kept 25 by
            # score, and Karaoglan 2022 — the single most on-point RCT in the
            # library, similarity 0.648 — ranked 54th of 60 BY SCORE and was
            # cut. Twenty of the twenty-five it kept were LESS similar to the
            # question than the one it dropped, led by AAE and ESE position
            # statements at score 90.0 and 87.0 with similarities of 0.57 and
            # 0.56 (which is also A7: guideline rows sitting in level1 at a
            # score no trial reaches). The answer then declared that no
            # prospective study directly compares the two protocols — A5's
            # false evidence gap, caused here.
            #
            # So: SELECT the survivors by similarity to this question, then
            # order them by score for display, which is invariant 1 unchanged
            # — tier by study design, score ranking only within a tier.
            # PROVISIONAL_KEY is not in TIER_ORDER and needs no handling HERE,
            # unlike every other loop over this list. The provisional lane is
            # `fetch_untyped_recent`, a live PubMed query for papers MEDLINE
            # has not yet classified, and it never writes back to the library
            # (pinned by test_provisional_lane). So no library row can carry
            # that key and there is nothing on this branch for the lane to
            # drop. If write-back ever admits one, this loop is where it would
            # silently vanish.
            for tier in TIER_ORDER:
                bucket = by_tier.get(tier)
                if not bucket:
                    continue
                bucket = cap_by_relevance(bucket, MAX_RAG_PAPERS_PER_TIER, tier)
                bucket.sort(key=lambda x: x["score"], reverse=True)
                evidence[tier] = {
                    "text":   _scored_to_text(bucket, TIER_LABEL.get(tier, tier.upper())),
                    "ids":    [p["pmid"] for p in bucket],
                    "scored": bucket,
                    "source": "rag",
                }
                all_scored.extend(bucket)

            update_job(job_id, message="Library search complete — asking Claude...", progress=75)
            # Apply outlier detection and currency tags to RAG results
            all_scored = detect_outliers(apply_currency_tags(all_scored))
            flag_superseded_by_review(evidence, question=question)
            avg_score = sum(p["score"] for p in all_scored) / len(all_scored) if all_scored else 0
            evidence["_summary"] = {
                "total_scored":    len(all_scored),
                "avg_score":       round(avg_score, 1),
                "all_scored":      sorted(all_scored, key=lambda x: x["score"], reverse=True),
                "synthesis_order": build_synthesis_order(evidence),
            }
            return evidence

    # ── Full PubMed fallback ──────────────────────────────
    # Generate multiple search terms for broader coverage (Feature 6)
    update_job(job_id, message="Generating multi-angle search terms...", progress=12)
    # The live path resolves an elliptical follow-up the same way the library
    # path does — through the query. It does NOT seed prior PMIDs: every paper
    # here arrives through fetch_papers(), which is where the per-tier quality
    # floors are applied, and injecting library rows past it would be a gate
    # bypass on exactly the path that has the least other protection.
    search_terms = generate_multi_search_terms(question, smart_topic,
                                               context_block=context_block)
    # A33h-i + A33g. Label each AND-group and expand the scenario group, on the
    # LIVE path only. Both are changes to a PubMed BOOLEAN, and the library
    # route does not evaluate one — it embeds the string. A longer, more
    # operator-dense query embeds FURTHER from a paper's prose (see the note at
    # the top of tests/test_retrieval_consistency.py: the better the boolean,
    # the worse the vector search), so expanding the query the library route
    # uses would be an unmeasured change with a known reason to expect harm.
    search_terms = label_and_expand(question, search_terms)
    print(f"  Multi-term search: {search_terms}")

    update_job(job_id, message="Searching Cochrane Reviews...", progress=15)
    cochrane_direct = fetch_cochrane(smart_topic)
    if cochrane_direct:
        evidence["cochrane"] = {"text": cochrane_direct, "ids": [], "scored": [],
                                "source": "pubmed"}
    else:
        text, ids, scored = fetch_papers(smart_topic, COCHRANE_TERM,
                                         "Cochrane Reviews (PubMed)", "cochrane",
                                         question=question)
        evidence["cochrane"] = {"text": text, "ids": ids, "scored": scored,
                                "source": "pubmed"}
        all_scored.extend(scored)

    # A SECOND, HARDCODED COPY OF THE LANE LIST USED TO LIVE HERE, and it had
    # silently fallen three lanes behind `endo_ai.tier_query_lanes()`.
    #
    # This function is the LIVE path for Review and Case. `build_evidence_base`
    # in endo_ai.py is the curriculum path. They are two implementations of the
    # same idea, and the list below was written out longhand in both — so every
    # lane added to the ladder since had reached the curriculum and nothing
    # else:
    #
    #   observational  A31. Added so cross-sectional, morphometric, imaging and
    #                  diagnostic-accuracy designs would be REACHABLE at all.
    #                  Never reached a Review or Case answer.
    #   guideline      A49 item 5, added the previous night. The whole point was
    #                  that a clinical practice guideline had no query that
    #                  could reach it; it still had none here.
    #   provisional    A49 item 4b, added tonight.
    #
    # Found by asking which call sites reach `fetch_untyped_recent`, not by a
    # test — every test of the lanes exercised `tier_query_lanes()` or
    # `endo_ai.build_evidence_base`, and both were correct. This is standing
    # rule 14 exactly: the helper was right and one of its two callers did not
    # call it. `tests/test_live_path_lane_parity.py` now pins the two lists
    # together so a lane cannot be added to one and not the other again.
    #
    # The fourth tuple element is vestigial — see the note in `_run_tiers`,
    # per-tier percentages are "a lie" once fetches run in parallel — so it is
    # filled with a constant rather than reinvented per lane.
    levels = [(lk, terms, label, 50) for lk, terms, label in tier_query_lanes()]

    # ── B2/B5: parallel tier fetches, in two phases ──────────────────────
    # Every (tier, search-term) pair is an independent HTTP round trip; there
    # are ~7 tiers x ~7 terms of them and they used to run one after another.
    # The NCBI limiter (endo_ai.ncbi_get) now paces departures globally, so
    # concurrency here is bounded by the rate limit rather than by the loop.
    #
    # Two properties have to survive, and both are easy to lose:
    #
    #  1. DEDUP ORDER. seen_pmids gave a duplicate paper to whichever tier was
    #     processed first, and tiers ran strongest-first — so a paper found in
    #     both level1 and level4 was presented as Level I. That is correct and
    #     must not become a race. Fetching is therefore parallel but dedup is
    #     applied afterwards, sequentially, in TIER_ORDER.
    #  2. EVIDENCE ORDER. evidence[] is built in tier order regardless of which
    #     fetch finished first, so completion order cannot leak into the answer.
    # SEEDED WITH THE COCHRANE TIER, which it was not. Cochrane is fetched
    # above, before this set exists, so a Cochrane review re-found by level1
    # was rendered in BOTH blocks and counted twice in all_scored and
    # avg_score. That is a hole inside the one invariant this function's own
    # comment (below) calls load-bearing -- and it hits the highest-authority
    # papers in the base, which are exactly the ones a duplicate presentation
    # misleads most.
    seen_pmids: set = {p["pmid"] for p in
                       ((evidence.get("cochrane") or {}).get("scored") or [])
                       if p.get("pmid")}
    _fetch_lock = threading.Lock()

    def _fetch_one(level_key, terms, label, term):
        """One (tier, term) fetch. Returns raw results; no dedup here."""
        if is_aborted(job_id):
            return level_key, None
        try:
            # `mode` and `max_results` were both dropped here, and both are
            # knobs the shared helper reads. Same class as the hardcoded lane
            # list: a setting added to endo_ai that one of its two callers
            # never forwarded.
            #
            #   mode        selects MODE_TIER_QUOTAS. In production B is only
            #               ever review or case, and those two tables are
            #               key-for-key identical, so this is hygiene today --
            #               but eval/run_eval.py can pass mode="learn", where
            #               they are NOT identical.
            #   max_results TIER_FETCH_DEPTH gives `observational` a depth of
            #               100 because A31 measured that the designs it
            #               admits sit deeper in the result list. The live
            #               path was fetching 50, so half that depth never
            #               existed on Review or Case.
            return level_key, fetch_papers(
                term, " OR ".join(terms), label, level_key,
                mode=mode, question=question,
                max_results=TIER_FETCH_DEPTH.get(level_key, 50))
        except Exception as e:
            print(f"  XX {label}: fetch failed ({e})")
            return level_key, None

    def _run_tiers(tier_specs):
        """Fetch every (tier, term) pair concurrently, then fold the results
        into `evidence` in strict tier order."""
        raw = {lk: [] for lk, _t, _l, _p in tier_specs}
        jobs_list = [(lk, terms, label, term)
                     for lk, terms, label, _pct in tier_specs
                     for term in search_terms]
        if not jobs_list:
            return 0

        done = 0
        with ThreadPoolExecutor(max_workers=TIER_FETCH_WORKERS) as pool:
            futures = [pool.submit(_fetch_one, *j) for j in jobs_list]
            for fut in as_completed(futures):
                lk, res = fut.result()
                done += 1
                if res is not None:
                    with _fetch_lock:
                        raw[lk].append(res)
                # Progress is a completion count, not a per-tier percentage:
                # with parallel fetches a monotonic per-tier pct is a lie.
                update_job(job_id,
                           message=f"Searching PubMed — {done}/{len(jobs_list)} queries",
                           progress=min(72, 20 + int(50 * done / len(jobs_list))))

        added = 0
        for level_key, _terms, _label, _pct in tier_specs:
            level_scored, level_ids = [], []
            for _text, ids, scored in raw[level_key]:
                new_scored = [p for p in scored if p["pmid"] not in seen_pmids]
                new_ids    = [i for i in ids    if i not in seen_pmids]
                for p in new_scored:
                    seen_pmids.add(p["pmid"])
                level_scored.extend(new_scored)
                level_ids.extend(new_ids)
            level_scored.sort(key=lambda x: x["score"], reverse=True)

            # THE PROMPT USED TO CARRY ONE SEARCH TERM'S PAPERS AND COUNT ALL
            # SEVEN. This fold ran once per (tier, term) — ~7 fetches per tier
            # — and accumulated `level_scored` over every one of them while
            # keeping only the FIRST non-empty text block:
            #
            #     if text and not level_text:
            #         level_text = text
            #
            # `_build_evidence_context` renders `block["text"]` and nothing
            # else, so Claude read one term's papers per tier while
            # `_summary` counted them all. MEASURED on a live Review
            # retrieval for "sodium hypochlorite concentration":
            #
            #     level1     73 scored,  3 in the prompt   70 never shown
            #     guideline   4 scored,  1 in the prompt    3 never shown
            #     TOTAL      99 scored, 26 in the prompt
            #     -> the model saw 26.3% of the retrieved evidence, under a
            #        header telling it "Total papers: 80 | Avg score: 62.2"
            #
            # That is the A5 false-evidence-gap mechanism itself: the answer
            # can state that no study addresses X while that study sits in
            # `scored`, and the "Top paper per tier" panel can name a paper
            # whose abstract was never in the prompt. Worse, `raw[lk]` is
            # appended in `as_completed` order, so WHICH term's block survived
            # was decided by whichever HTTP round trip finished first — the
            # same question asked twice could answer differently.
            #
            # The text is now rebuilt from the deduped, capped `level_scored`
            # with the renderer the library route and the differential merge
            # already use, so text and scored are one-to-one by construction —
            # the property `endo_ai.build_evidence_base` has always had.
            #
            # THE CAP MOVES WITH IT, and it has to. `fetch_papers` applies the
            # per-tier cap per CALL, so seven calls could accumulate ~7x the
            # intended quota; rendering all of that would have grown the
            # prompt roughly 24x. Capping the deduped list restores the quota
            # to what MODE_TIER_QUOTAS actually says and matches the
            # curriculum path.
            level_scored = cap_by_relevance(
                level_scored, _tier_cap(mode, level_key), level_key)
            level_text = _scored_to_text(
                level_scored, TIER_LABEL.get(level_key, level_key.upper()))

            evidence[level_key] = {"text": level_text, "ids": level_ids,
                                   "scored": level_scored, "source": "pubmed"}
            all_scored.extend(level_scored)
            added += len(level_scored)
        return added

    if not is_aborted(job_id):
        strong = _run_tiers([l for l in levels if l[0] in EARLY_STOP_TIERS])

        # B5 early stop. In Review mode, once the top tiers have supplied
        # enough evidence, the weak tiers cannot change the recommendation —
        # tier banding means a case series never overrides a Level I finding.
        # Learn mode is exempt: a teaching curriculum genuinely wants the
        # narrative scaffolding that reviews and editorials provide.
        n_strong = len((evidence.get("cochrane") or {}).get("scored") or []) + strong
        early = (mode == "review" and n_strong >= EARLY_STOP_MIN_PAPERS)
        if early:
            print(f"  [early_stop] {n_strong} papers from cochrane+level1 "
                  f">= {EARLY_STOP_MIN_PAPERS}; skipping weaker tiers (mode=review)")
            _pubmed_audit_log("early_stop", "level1",
                              f"n_strong={n_strong} threshold={EARLY_STOP_MIN_PAPERS}",
                              [], 200, 0)
            # ...but the guideline lane is NOT a weak tier, and the early
            # stop's reasoning does not cover it.
            #
            # That reasoning is "tier banding means a case series cannot
            # override a Level I finding, so once the top tiers have supplied
            # enough, the weak ones cannot change the recommendation". A
            # guideline is not weaker evidence — it is a specialty's stated
            # POSITION, a different axis entirely, and A49 item 5 exists
            # because nothing could reach one. Skipping it here would have
            # left it unreachable on exactly the well-covered Review questions
            # a clinician is most likely to ask.
            _run_tiers([l for l in levels if l[0] == "guideline"])
        else:
            _run_tiers([l for l in levels if l[0] not in EARLY_STOP_TIERS])

        # A49 item 4b — the provisional lane, and DELIBERATELY OUTSIDE the
        # early stop.
        #
        # The early stop skips the weak tiers once cochrane+level1 have
        # supplied enough, on the reasoning that tier banding means a case
        # series cannot override a Level I finding anyway. That reasoning does
        # not transfer here. This lane exists to reach papers MEDLINE has not
        # classified yet, and a new trial is MOST valuable precisely when there
        # is established evidence for it to contradict — skipping it on
        # well-covered questions would skip it exactly where it matters.
        #
        # It costs one esearch and a batched efetch, ~5s, and it takes no slot
        # from any tier: PROVISIONAL_KEY is not in TIER_ORDER.
        if not is_aborted(job_id):
            try:
                p_text, p_ids, p_papers = fetch_untyped_recent(smart_topic)
                evidence[PROVISIONAL_KEY] = {
                    "text": p_text, "ids": p_ids, "scored": p_papers,
                    "source": "pubmed"}
            except Exception as e:
                # Never let the newest-literature lane take an answer down.
                print(f"  [provisional] lane failed, continuing without it: {e}")

    # Apply outlier detection and currency tags to PubMed results
    all_scored = detect_outliers(apply_currency_tags(all_scored))
    # `question=` is what stops this path nominating by YEAR while the library
    # branch 300 lines up nominates by RELEVANCE — item 2. Without it the two
    # branches of the SAME function pick a different review on 27 of 29
    # questions, and the blind panel preferred the relevance pick 23 times.
    flag_superseded_by_review(evidence, question=question)
    avg_score = sum(p["score"] for p in all_scored) / len(all_scored) if all_scored else 0
    evidence["_summary"] = {
        "total_scored":    len(all_scored),
        "avg_score":       round(avg_score, 1),
        "all_scored":      sorted(all_scored, key=lambda x: x["score"], reverse=True),
        "synthesis_order": build_synthesis_order(evidence),
        # Distinct PMIDs that came back from PubMed across every term and tier,
        # before the per-tier cap. This is the number that actually collapsed in
        # the laser regression (5 across 28 queries), and it moves independently
        # of the final paper count, so the eval can tell a broken query apart
        # from a genuinely thin topic.
        "distinct_pmids_retrieved": len(seen_pmids),
    }
    return evidence


def build_differential_evidence(job_id: str, case_description: str,
                                candidates: list, progress_lo: int = 15,
                                progress_hi: int = 70,
                                prior_pmids: list = None) -> dict:
    """One evidence base covering EVERY candidate cause (`case-v2` Item 3b).

    Runs the existing evidence engine once per candidate — the candidate's own
    search topic joined to the case features — and unions the results, tagging
    each paper with the candidate that retrieved it. Deliberately the same
    engine: per-tier quality floors, the retracted/withdrawn/superseded
    exclusions, the similarity floor and the routing gate all apply to a
    candidate's papers exactly as they apply to any others. A second retrieval
    path would be a second place for those to be missing.

    WHY PER CANDIDATE AND NOT ONE BROADER QUERY. Measured in Item 1: one query
    generation put trauma in 8 runs of 8 and dens invaginatus in 2, the
    palatogingival groove in 0. Widening the single query does not fix that —
    it is the same one call, asked to be luckier. Asking per candidate makes
    coverage a property of the differential rather than of a sample.

    COST. Each candidate is one retrieval, so a 6-candidate differential is up
    to 6. Two things bound it, and neither is a cap on candidates:

      - the LIBRARY GATE is shared. `build_evidence_base_with_progress` decides
        per candidate whether the library covers it, so a candidate the library
        already knows costs embeddings and no PubMed traffic at all.
      - a candidate that returns nothing is NOT retried and NOT broadened. An
        empty result is information — it is the difference between "the
        literature disagrees" and "nobody has studied this in this
        presentation", and the answer has to be able to say which.

    Returns the usual evidence dict, plus `_differential`:
      {candidate: {"n_papers": int, "pmids": [...], **the candidate's fields}}
    and each paper carries `candidates: [names]`.
    """
    from endo_ai import (TIER_ORDER, TIER_LABEL, build_synthesis_order,
                         detect_outliers, apply_currency_tags,
                         flag_superseded_by_review,
                         PROVISIONAL_KEY, PROVISIONAL_MAX_ADMITTED,
                         _provisional_context_line)
    # RELEVANCE_GATE is this module's, and reading it here rather than copying
    # a number is the point: a per-tier cap that drifted between the two paths
    # would mean a differential answer silently sees more or fewer papers per
    # tier than an ordinary case answer, with nothing saying so.
    max_per_tier = RELEVANCE_GATE["max_per_tier"]

    merged: dict = {}
    by_pmid: dict = {}
    per_candidate: dict = {}
    n = max(1, len(candidates))

    for i, cand in enumerate(candidates):
        name = cand.get("candidate") or ""
        topic = cand.get("search_topic") or name
        # The candidate's topic AND the case features, because "dens
        # invaginatus" alone retrieves the anomaly's treatment literature and
        # the question is about it as a CAUSE in this presentation.
        query = f"{topic} — {case_description}"
        pct = progress_lo + int((progress_hi - progress_lo) * i / n)
        update_job(job_id, progress=pct,
                   message=f"Searching literature for: {name} "
                           f"({i + 1} of {len(candidates)})")
        print(f"\n[differential] {i + 1}/{len(candidates)} — {name}")
        try:
            # `prior_pmids` goes to EVERY candidate, not once: a paper
            # carried from the previous turn is judged against each
            # candidate's own query, so it enters through whichever candidate
            # it is actually relevant to. The union dedupes by pmid.
            ev = build_evidence_base_with_progress(
                job_id, query, mode="case", prior_pmids=prior_pmids)
        except Exception as e:
            print(f"  [differential] retrieval failed for {name!r}: {e}")
            per_candidate[name] = {**cand, "n_papers": 0, "pmids": [],
                                   "error": str(e)}
            continue

        found = []
        # TIER_ORDER + the provisional lane. PROVISIONAL_KEY is deliberately
        # NOT in TIER_ORDER -- that absence is what stops it competing for a
        # tier slot -- but the same absence made this merge DROP every
        # provisional paper the retrieval above had just found. They would
        # have been fetched, paid for, and then silently discarded on the
        # differential path, and because the merge also builds the evidence
        # base, citing one would have scored as a FABRICATION.
        for tier in list(TIER_ORDER) + [PROVISIONAL_KEY]:
            for p in ((ev.get(tier) or {}).get("scored") or []):
                pmid = p.get("pmid")
                if not pmid:
                    continue
                found.append(pmid)
                seen = by_pmid.get(pmid)
                if seen is None:
                    paper = dict(p)
                    paper["candidates"] = [name]
                    by_pmid[pmid] = paper
                    merged.setdefault(tier, []).append(paper)
                elif name not in seen["candidates"]:
                    # One paper can bear on two candidates. Recording both is
                    # what lets the answer say so instead of silently
                    # attributing it to whichever search ran first.
                    seen["candidates"].append(name)
        per_candidate[name] = {**cand, "n_papers": len(found),
                               "pmids": found[:40]}
        print(f"  [differential] {name}: {len(found)} paper(s)")

    evidence: dict = {}
    all_scored = []
    for tier in TIER_ORDER:
        bucket = merged.get(tier) or []
        if not bucket:
            continue
        # A30b, at a site A30a's enumeration missed. This capped the merged
        # union BY SCORE — sort by score, keep the first 25 — which is exactly
        # the membership-by-quality error standing rule 19 exists for, and the
        # twin of the per-tier cap A5b found cutting the most on-point RCT at
        # rank 54 of 60. It is worse here than on the ordinary paths, because
        # the union it cuts is the only place a paper retrieved for a WEAK
        # candidate can reach the answer: the differential's whole purpose is
        # to carry evidence for the causes that are not the leading one, and
        # score does not know which candidate a paper was retrieved for.
        #
        # Papers arriving from the library route carry a similarity;
        # cap_by_relevance ties-breaks on score, so ones from the live route
        # (no similarity) fall back to exactly the score order used before.
        bucket = cap_by_relevance(bucket, max_per_tier, tier)
        bucket.sort(key=lambda x: x.get("score") or 0, reverse=True)
        evidence[tier] = {
            "text":   _scored_to_text(bucket, TIER_LABEL.get(tier, tier.upper())),
            "ids":    [p["pmid"] for p in bucket],
            "scored": bucket,
            "source": "differential",
        }
        all_scored.extend(bucket)

    # The provisional lane, rebuilt as its own block after the tier loop. It
    # is NOT added to `all_scored`: those papers carry score=None, and
    # `sum(p["score"] for p in all_scored)` two lines below would raise on the
    # first one. Same reasoning as endo_ai.build_evidence_base.
    prov_bucket = merged.get(PROVISIONAL_KEY) or []
    if prov_bucket:
        prov_bucket = prov_bucket[:PROVISIONAL_MAX_ADMITTED]
        evidence[PROVISIONAL_KEY] = {
            "text": "".join(_provisional_context_line(p) for p in prov_bucket),
            "ids": [p["pmid"] for p in prov_bucket],
            "scored": prov_bucket,
            "source": "differential",
        }

    all_scored = detect_outliers(apply_currency_tags(all_scored))
    # The differential's "question" is the case description: it is what every
    # candidate retrieval was seeded from, so it is the right thing for the
    # nominated review to be relevant TO.
    flag_superseded_by_review(evidence, question=case_description)
    avg = (sum(p["score"] for p in all_scored) / len(all_scored)
           if all_scored else 0)
    evidence["_summary"] = {
        "total_scored":    len(all_scored),
        "avg_score":       round(avg, 1),
        "all_scored":      sorted(all_scored, key=lambda x: x["score"],
                                  reverse=True),
        "synthesis_order": build_synthesis_order(evidence),
    }
    evidence["_differential"] = per_candidate
    covered = sum(1 for v in per_candidate.values() if v["n_papers"])
    print(f"\n[differential] union: {len(all_scored)} paper(s) across "
          f"{len(per_candidate)} candidate(s); {covered} candidate(s) have "
          f"literature, {len(per_candidate) - covered} have none")
    return evidence


def _scored_to_text(scored_papers: list, label: str) -> str:
    """Convert scored paper dicts back to annotated text for Claude context.

    Uses the SAME renderer as the live-PubMed path so provenance badges (COI,
    pre-registration, corrections, indexing) appear identically regardless of
    which retrieval path answered. Previously this built its own line and
    emitted no badges at all, so library-served evidence reached Claude
    stripped of every integrity signal.

    The title and abstract follow that line. They did not, for as long as this
    function has existed: sharing the metadata RENDERER was mistaken for
    sharing the BLOCK, and the parity test in `tests/test_coi_scoping.py`
    compared the two LINES and passed throughout. The result was a prompt that
    named a paper and never said what it found, while instructing the model to
    write a paragraph on what the evidence shows and to put a [[PMID:N]] marker
    on every clinical claim.

    ONE ASYMMETRY REMAINS, deliberately. The live path
    (`endo_ai.build_evidence_base`, annotated-text loop) emits the abstract OR,
    failing that, the title — never both. This emits both, so a library-served
    paper reaches Claude with about 100 more characters of title than a
    live-served one. The 2026-08-31 before/after measurement of the
    citation-support flag rate (39.4% -> 8.5%) was taken with both present, so
    aligning the two paths means re-measuring, and the sensible direction is to
    give the LIVE path the title as well rather than take it away here.

    2026-09-05, item 2 — AND THAT CHANGED THE CASE DIFFERENTIAL, which is the
    third caller of this function and was not considered above. Papers reaching
    `build_differential_evidence` from the LIVE route carried no `title` and no
    `abstract` on their dicts at all, because `fetch_papers` did not put them
    there. So the two loops below found nothing to emit and a live-sourced
    candidate paper reached Claude through this renderer as a metadata line and
    NOTHING ELSE — the exact defect the 2026-08-31 fix above describes, still
    live on one path, on the mode where a paper is retrieved specifically to
    speak for a candidate cause.
    Item 2 added both fields to `fetch_papers` for a different reason (the
    PRISMA nomination needed text to embed) and closed this as a side effect.
    It is a real prompt change on the Case differential path and the v6
    baseline prediction names it.
    """
    from endo_ai import format_paper_context_line
    text = f"\n[{label}]\n"
    for p in scored_papers:
        text += format_paper_context_line(p)
        title = (p.get("title") or "").strip()
        abstract = (p.get("abstract") or "").strip()
        if title:
            text += title + "\n"
        if abstract:
            text += abstract + "\n"
    return text


# ── Deep Learning history archive ─────────────────────────
# Every completed Deep Learning curriculum is persisted as JSON to
# learn_history/. The 7-day re-use window is enforced via the query_cache
# (see rag.get_cached_answer's max_age_days param) — this folder is the
# durable browsable archive (no auto-cleanup).
import re as _learn_re

_LEARN_HISTORY_DIR        = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                          "learn_history")
LEARN_HISTORY_TTL_DAYS    = 7   # cache window for auto re-use
REVIEW_CACHE_TTL_DAYS     = 30  # clinical answers must be re-derived periodically

def save_learn_output(question: str, answer: str, evidence: dict, cost: float) -> str:
    """Persist one completed Deep Learning curriculum to learn_history/.
    Returns the file path written, or '' on failure."""
    try:
        os.makedirs(_LEARN_HISTORY_DIR, exist_ok=True)

        # Filesystem-safe slug from the question
        slug = _learn_re.sub(r"[^a-z0-9]+", "_",
                             (question or "").lower()).strip("_")[:60]
        ts   = datetime.now().strftime("%Y%m%d_%H%M%S")
        path = os.path.join(_LEARN_HISTORY_DIR, f"{ts}_{slug or 'untitled'}.json")

        summary = (evidence or {}).get("_summary", {}) or {}
        payload = {
            "question":          question,
            "timestamp":         datetime.now().isoformat(),
            "answer":            answer,
            "cost_usd":          round(float(cost or 0.0), 4),
            "total_papers":      summary.get("total_scored", 0),
            "avg_paper_score":   summary.get("avg_score", 0),
            "top_pmids":         [p.get("pmid", "") for p in (summary.get("all_scored") or [])[:10]],
            # Paper metadata (no abstracts) so archived reports can render
            # author-style citations instead of bare PMIDs.
            "papers":            _safe_papers(summary.get("all_scored") or []),
        }
        with open(path, "w", encoding="utf-8") as fh:
            _audit_json.dump(payload, fh, ensure_ascii=False, indent=2)
        print(f"  [learn_history] saved -> {os.path.basename(path)}")
        return path
    except Exception as e:
        print(f"  [learn_history] save failed: {e}")
        return ""


@app.route("/learn_history/<filename>")
def get_learn_history_item(filename: str):
    """Return the full archived Deep Learning curriculum (answer + metadata)."""
    # Tight path validation — only allow filenames that exist in the archive
    if "/" in filename or "\\" in filename or ".." in filename or not filename.endswith(".json"):
        return jsonify({"error": "invalid filename"}), 400
    path = os.path.join(_LEARN_HISTORY_DIR, filename)
    if not os.path.isfile(path):
        return jsonify({"error": "not found"}), 404
    try:
        with open(path, encoding="utf-8") as fh:
            rec = _audit_json.load(fh)
        rec["question"] = display_title(rec.get("question", ""))
        # A16b — same reasoning as /history/<cache_id>. Measured on the 22
        # stored curricula: 13 gain the banner's second half and 18 of 22
        # render a bibliography of the retrieval pool rather than the
        # citation set.
        rec["answer"], _lq = finalise_answer_text(rec.get("answer") or "")
        rec["cited_pmids"] = assemble_bibliography(
            rec["answer"], rec.get("papers") or [])["cited_pmids"]
        return jsonify(rec)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/learn_history/<filename>", methods=["DELETE"])
@require_admin_token
def delete_learn_history_item(filename: str):
    """Permanently delete a single archived Deep Learning curriculum.

    Token-gated (WORKLIST 4.1): each curriculum costs ~$1 to regenerate, so an
    unauthenticated DELETE is the most destructive route in the app. The UI's
    sidebar delete button sends the token injected into the page by index();
    with ADMIN_TOKEN unset the route 403s and the UI says so.
    """
    if "/" in filename or "\\" in filename or ".." in filename or not filename.endswith(".json"):
        return jsonify({"error": "invalid filename"}), 400
    path = os.path.join(_LEARN_HISTORY_DIR, filename)
    if not os.path.isfile(path):
        return jsonify({"error": "not found"}), 404
    try:
        os.remove(path)
        return jsonify({"ok": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/learn_history")
def list_learn_history():
    """List archived Deep Learning curricula, newest first."""
    if not os.path.isdir(_LEARN_HISTORY_DIR):
        return jsonify({"items": [], "ttl_days": LEARN_HISTORY_TTL_DAYS})
    items = []
    for fn in sorted(os.listdir(_LEARN_HISTORY_DIR), reverse=True):
        if not fn.endswith(".json"):
            continue
        path = os.path.join(_LEARN_HISTORY_DIR, fn)
        try:
            with open(path, encoding="utf-8") as fh:
                rec = _audit_json.load(fh)
            ts = datetime.fromisoformat(rec.get("timestamp", ""))
            age_days = (datetime.now() - ts).days
            items.append({
                "file":           fn,
                "question":       display_title(rec.get("question", "")),
                "timestamp":      rec.get("timestamp", ""),
                "age_days":       age_days,
                "in_ttl_window":  age_days <= LEARN_HISTORY_TTL_DAYS,
                "total_papers":   rec.get("total_papers", 0),
                "cost_usd":       rec.get("cost_usd", 0),
            })
        except Exception:
            continue
    return jsonify({"items": items, "ttl_days": LEARN_HISTORY_TTL_DAYS})


# ── Citation audit log (FDA CDS verifiability) ────────────
import re as _audit_re
import json as _audit_json
import hashlib as _audit_hash

_AUDIT_DIR = os.path.join(os.path.dirname(__file__), "audit_logs")

def write_citation_audit(question: str, answer: str, mode: str) -> int:
    """Extract every claim → PMID mapping from the answer and persist to disk
    so a clinician (or auditor) can later reconstruct exactly which paper
    supported each statement. Returns the number of claims logged."""
    if not answer:
        return 0

    # Match: <sentence-ish text>(one-or-more [[PMID:N]] markers)
    # Anchor sentence start at \n, sentence boundary punctuation, or doc start.
    pattern = _audit_re.compile(
        r'([^\n.!?]{10,400}?)\s*((?:\[\[PMID:\d+\]\]\s*)+)',
        _audit_re.MULTILINE,
    )

    claims = []
    for m in pattern.finditer(answer):
        sentence = m.group(1).strip()
        pmids    = _audit_re.findall(r'\[\[PMID:(\d+)\]\]', m.group(2))
        if sentence and pmids:
            claims.append({"claim": sentence, "pmids": pmids})

    if not claims:
        return 0

    try:
        os.makedirs(_AUDIT_DIR, exist_ok=True)
        ts   = datetime.now().strftime("%Y%m%d_%H%M%S")
        slug = _audit_hash.sha256(question.encode("utf-8")).hexdigest()[:10]
        path = os.path.join(_AUDIT_DIR, f"audit_{ts}_{slug}.json")
        payload = {
            "timestamp":       datetime.now().isoformat(),
            "mode":            mode,
            "question":        question,
            "answer_sha256":   _audit_hash.sha256(answer.encode("utf-8")).hexdigest(),
            "claim_count":     len(claims),
            "claim_citations": claims,
        }
        with open(path, "w", encoding="utf-8") as fh:
            _audit_json.dump(payload, fh, indent=2, ensure_ascii=False)
        print(f"  [audit] wrote {len(claims)} claim citations -> {os.path.basename(path)}")
    except Exception as e:
        print(f"  [audit] write failed: {e}")

    return len(claims)


# ── Provenance / inline-citation abstract endpoint ────────
# Backs the [[PMID:n]] click-through side panel in the UI. Per FDA CDS
# verifiability guidance, every clinical claim the AI emits should be
# independently inspectable by the clinician.

import requests as _prov_requests
_ABSTRACT_CACHE     = {}      # PMID -> dict ; simple LRU by insertion order
_ABSTRACT_CACHE_MAX = 500

def _trim_abstract_cache():
    while len(_ABSTRACT_CACHE) > _ABSTRACT_CACHE_MAX:
        _ABSTRACT_CACHE.pop(next(iter(_ABSTRACT_CACHE)))

def _eutils_params(extra: dict) -> dict:
    """Merge eutils params with NCBI tool/email + api_key (if set in env).

    Unauthenticated eutils is rate-limited to 3 req/sec on slow servers; an
    API key bumps that to 10 req/sec on prioritised hardware. Set
    NCBI_API_KEY in .env to enable.
    """
    p = dict(extra or {})
    p.setdefault("tool",  "endo-ai-rag")
    p.setdefault("email", os.getenv("NCBI_EMAIL", "endoai@research.local"))
    api_key = (os.getenv("NCBI_API_KEY") or "").strip()
    if api_key:
        p["api_key"] = api_key
    return p


def _eutils_get(url: str, params: dict):
    """GET an eutils endpoint with one retry on transient failures.

    NCBI eutils routinely takes 5-15s under load; an 8s timeout was making
    abstract panel loads fail visibly. We use (connect=5, read=20) and retry
    once on timeout/connection errors before surfacing the failure.
    """
    last_exc = None
    final_params = _eutils_params(params)
    for attempt in (1, 2):
        try:
            return _prov_requests.get(url, params=final_params, timeout=(5, 20))
        except (_prov_requests.exceptions.ReadTimeout,
                _prov_requests.exceptions.ConnectTimeout,
                _prov_requests.exceptions.ConnectionError) as e:
            last_exc = e
            if attempt == 2:
                raise
    raise last_exc  # unreachable, satisfies linters


_ABSTRACT_KEY_RE = _learn_re.compile(r"^(?:\d{1,10}|[A-Za-z][A-Za-z0-9._-]{1,63})$")


@app.route("/api/abstract/<pmid>")
def get_abstract(pmid):
    """Fetch full abstract + metadata for a single PMID.

    Three-tier cache lookup:
      L1 — in-process dict (zero ms)
      L2 — Postgres abstract_cache (one local round-trip; populated for free
            during build_evidence_base())
      L3 — live NCBI eutils (last resort; slow + rate-limited)

    On L3 fetch, we backfill L2 + L1 so the next click is instant.
    """
    # `trust-surface-v1` Q4. The guard used to be `pmid.isdigit()`, which meant
    # a citation pill for a hand-ingested authority document (ESE-QG-2023,
    # AAE-PS-obturation, NBK430685) answered 400 — so making those markers
    # render as pills without widening this would have swapped a raw marker for
    # a dead one. A synthetic key is served from the library ONLY: it has no
    # PubMed record, so the live-eutils tier below is skipped for it and no
    # pubmed.ncbi.nlm.nih.gov URL is ever handed back for one.
    if not _ABSTRACT_KEY_RE.match(pmid or ""):
        return jsonify({"error": "invalid PMID"}), 400
    is_numeric = pmid.isdigit()

    # L1 — in-process
    if pmid in _ABSTRACT_CACHE:
        resp = dict(_ABSTRACT_CACHE[pmid])
        resp["source"] = "memory"
        return jsonify(resp)

    # L2 — Postgres abstract_cache (populated during build_evidence_base)
    cached = get_cached_abstract(pmid)
    if cached and cached.get("abstract") and len(cached.get("abstract") or "") >= 50:
        result = {
            "pmid":     pmid,
            "title":    (cached.get("title") or "Title unavailable").rstrip("."),
            "abstract": cached.get("abstract") or "Abstract unavailable.",
            "journal":  cached.get("journal") or "",
            "year":     cached.get("year")    or "",
            "authors":  cached.get("authors") or "",
            "url":      (f"https://pubmed.ncbi.nlm.nih.gov/{pmid}/"
                         if is_numeric else ""),
            "source":   "postgres",
        }
        _ABSTRACT_CACHE[pmid] = result
        _trim_abstract_cache()
        return jsonify(result)

    # L3 — live eutils (slow path). Numeric ids only: a synthetic key has no
    # PubMed record, and asking eutils for "ESE-QG-2023" spends a rate-limited
    # round trip to be told so. Say plainly that the library copy is missing.
    if not is_numeric:
        return jsonify({
            "error": "This authority document is held in the library and has no "
                     "PubMed record. Its abstract is not cached.",
            "kind":  "local_only",
        }), 404

    try:
        # esummary for structured metadata (title, journal, year, authors)
        s = _eutils_get(
            "https://eutils.ncbi.nlm.nih.gov/entrez/eutils/esummary.fcgi",
            params={"db": "pubmed", "id": pmid, "retmode": "json"},
        )
        meta = (s.json().get("result", {}) or {}).get(pmid, {}) or {}

        # efetch text for the abstract body
        f = _eutils_get(
            "https://eutils.ncbi.nlm.nih.gov/entrez/eutils/efetch.fcgi",
            params={"db": "pubmed", "id": pmid, "rettype": "abstract", "retmode": "text"},
        )
        raw = f.text or ""

        # PubMed efetch returns: citation line, title, authors, affiliation,
        # then abstract, then DOI/PMID footer. Pull the paragraphs and hand
        # them to the shared selector — "longest paragraph" alone picks the
        # AFFILIATION block on a paper with thirty institutional addresses,
        # and the Portuguese translation on a paper that has both.
        paragraphs = []
        current    = []
        for line in raw.split("\n"):
            line = line.rstrip()
            if line.strip():
                current.append(line.strip())
            else:
                if current:
                    paragraphs.append(" ".join(current))
                current = []
        if current:
            paragraphs.append(" ".join(current))
        from endo_ai import _select_abstract_paragraph
        abstract = _select_abstract_paragraph(paragraphs)

        authors_list = meta.get("authors", []) or []
        names = [a.get("name", "") for a in authors_list if a.get("name")]
        if len(names) > 5:
            authors_str = ", ".join(names[:5]) + ", et al."
        else:
            authors_str = ", ".join(names)

        result = {
            "pmid":     pmid,
            "title":    (meta.get("title", "") or "Title unavailable").rstrip("."),
            "abstract": abstract or "Abstract unavailable.",
            "journal":  meta.get("fulljournalname", "") or meta.get("source", ""),
            "year":     (meta.get("pubdate", "") or "")[:4],
            "authors":  authors_str,
            "url":      f"https://pubmed.ncbi.nlm.nih.gov/{pmid}/",
            "source":   "eutils_live",
        }
        _ABSTRACT_CACHE[pmid] = result
        _trim_abstract_cache()
        # Backfill Postgres so the next click on this PMID is instant
        if abstract and len(abstract) >= 50:
            try:
                cache_abstract(
                    pmid     = pmid,
                    title    = result["title"],
                    abstract = result["abstract"],
                    journal  = result["journal"],
                    year     = result["year"],
                    authors  = result["authors"],
                    source   = "eutils_live",
                )
            except Exception as ce:
                print(f"  [abstract_cache] backfill skipped (pmid={pmid}): {ce}")
        return jsonify(result)

    except (_prov_requests.exceptions.ReadTimeout,
            _prov_requests.exceptions.ConnectTimeout) as e:
        return jsonify({
            "error": "PubMed (NCBI eutils) is slow right now — request timed out after retry. "
                     "Try again in a moment, or open the source directly:",
            "url":   f"https://pubmed.ncbi.nlm.nih.gov/{pmid}/",
            "kind":  "timeout",
        }), 504
    except _prov_requests.exceptions.ConnectionError as e:
        return jsonify({
            "error": "Could not reach PubMed (NCBI eutils). Check your internet connection, "
                     "or open the source directly:",
            "url":   f"https://pubmed.ncbi.nlm.nih.gov/{pmid}/",
            "kind":  "connection",
        }), 502
    except Exception as e:
        return jsonify({
            "error": f"Could not load abstract: {e}",
            "url":   f"https://pubmed.ncbi.nlm.nih.gov/{pmid}/",
            "kind":  "other",
        }), 502


# ── Cost telemetry ────────────────────────────────────────
# Reads endo_ai's append-only cost_log.jsonl and returns aggregate stats
# so the operator can measure savings from the model-routing changes.

@app.route("/admin/costs")
@require_admin_token
def admin_costs():
    """Aggregate Claude API cost data over the last N days.

    Optional query params:
      ?days=7       window length (default 7)
      ?source=      which writer's rows to count. DEFAULT `product`, so this
                    endpoint answers "what did the product spend?" — which is
                    the question it is asked. `test`, `script`, or `all` for
                    the others. A row with no `source` reads as `product`,
                    because that is what the rows written before the field
                    existed mostly were.

                    This is here because the suite used to append to this log
                    and put $5.70 of stubbed TTS into it. Those rows are still
                    there and are NOT edited — an append-only audit log is not
                    rewritten after the fact — so the reader filters instead.
                    `?source=all` reproduces the old, contaminated number, and
                    `excluded_calls` below says how many rows the filter
                    dropped so a silent filter cannot look like a clean log.

    Returns:
      total_cost_usd
      total_calls
      window_days
      by_mode          — { mode: { calls, total_cost, avg_cost_per_call } }
      by_model         — { model: { calls, total_cost, in_tokens, out_tokens } }
      by_function      — { fn:    { calls, total_cost, avg_cost_per_call } }
      avg_cost_per_request_by_mode — see note below
    """
    from endo_ai import _COST_LOG_PATH
    from datetime import timedelta
    import json

    try:
        days = int(request.args.get("days", "7"))
    except ValueError:
        days = 7
    cutoff = datetime.now() - timedelta(days=days)

    want_source = (request.args.get("source") or "product").strip().lower()
    if want_source not in ("product", "test", "script", "all"):
        want_source = "product"

    if not os.path.exists(_COST_LOG_PATH):
        return jsonify({
            "total_cost_usd": 0.0, "total_calls": 0, "window_days": days,
            "source": want_source, "excluded_calls": 0, "by_source": {},
            "by_mode": {}, "by_model": {}, "by_function": {},
            "avg_cost_per_request_by_mode": {},
            "note": "cost_log.jsonl does not exist yet — no Claude calls logged.",
        })

    by_mode:     dict = {}
    by_model:    dict = {}
    by_function: dict = {}
    by_source:   dict = {}
    total_cost   = 0.0
    total_calls  = 0
    excluded     = 0

    with open(_COST_LOG_PATH, encoding="utf-8") as fh:
        for line in fh:
            line = line.strip()
            if not line:
                continue
            try:
                rec = json.loads(line)
                ts  = datetime.fromisoformat(rec.get("ts", ""))
            except Exception:
                continue
            if ts < cutoff:
                continue

            cost = float(rec.get("cost_usd", 0.0))
            # A row written before `source` existed reads as product.
            src  = (rec.get("source") or "product").strip().lower()
            b = by_source.setdefault(src, {"calls": 0, "total_cost": 0.0})
            b["calls"] += 1; b["total_cost"] += cost
            if want_source != "all" and src != want_source:
                excluded += 1
                continue

            mode = rec.get("mode", "unknown")
            mdl  = rec.get("model", "unknown")
            fn   = rec.get("function", "unknown")
            in_t = int(rec.get("input_tokens", 0))
            out_t = int(rec.get("output_tokens", 0))

            total_cost  += cost
            total_calls += 1

            m = by_mode.setdefault(mode, {"calls": 0, "total_cost": 0.0})
            m["calls"] += 1; m["total_cost"] += cost

            mm = by_model.setdefault(mdl, {"calls": 0, "total_cost": 0.0,
                                            "in_tokens": 0, "out_tokens": 0})
            mm["calls"] += 1; mm["total_cost"] += cost
            mm["in_tokens"] += in_t; mm["out_tokens"] += out_t

            f = by_function.setdefault(fn, {"calls": 0, "total_cost": 0.0})
            f["calls"] += 1; f["total_cost"] += cost

    # Round + add per-call averages
    for d in (by_mode, by_function):
        for k, v in d.items():
            v["total_cost"] = round(v["total_cost"], 4)
            v["avg_cost_per_call"] = round(v["total_cost"] / max(v["calls"], 1), 6)
    for k, v in by_model.items():
        v["total_cost"] = round(v["total_cost"], 4)
    for k, v in by_source.items():
        v["total_cost"] = round(v["total_cost"], 4)

    # Approximate "cost per user-initiated request" by mode.
    # Each mode has a designated "primary" function whose call count proxies
    # the request count. Multi-step pipelines (Learn) sum all sub-call costs
    # then divide by the count of the terminal stitch call.
    PRIMARY_FN = {
        "review":     "ask_clinical_question",
        "learn":      "stitch_curriculum",       # 1 stitch == 1 user request
        "case":       "ask_case_question",       # per-turn, not per-conv
        "assessment": "generate_referral_letter",
        "export":     "generate_slides_content", # treat any export as a "request"
    }
    avg_per_request = {}
    for mode, primary in PRIMARY_FN.items():
        n_requests = (by_function.get(primary) or {}).get("calls", 0)
        mode_total = (by_mode.get(mode)        or {}).get("total_cost", 0.0)
        if n_requests:
            avg_per_request[mode] = {
                "requests":           n_requests,
                "avg_cost_per_request": round(mode_total / n_requests, 4),
                "primary_fn":         primary,
            }

    return jsonify({
        "window_days":                  days,
        "source":                       want_source,
        "excluded_calls":               excluded,
        "by_source":                    by_source,
        "total_calls":                  total_calls,
        "total_cost_usd":               round(total_cost, 4),
        "by_mode":                      by_mode,
        "by_model":                     by_model,
        "by_function":                  by_function,
        "avg_cost_per_request_by_mode": avg_per_request,
        "note": (
            "avg_cost_per_request_by_mode uses a primary-function heuristic: "
            "for Deep Learning, request count = stitch_curriculum calls "
            "(each Learn request fires 1 syllabus + 4 modules + 1 stitch). "
            "Totals cover source=" + want_source + "; by_source counts EVERY "
            "row in the window, filtered or not, so a filter can never look "
            "like an empty log."
        ),
    })


# ── Evidence-mapping telemetry ────────────────────────────
# Reads endo_ai's append-only evidence_mapping.jsonl and returns aggregate
# pass/fail/retry stats so the operator can monitor whether Claude is
# actually grounding its claims in the retrieved evidence base.

@app.route("/admin/evidence-mapping")
@require_admin_token
def admin_evidence_mapping():
    """Aggregate evidence-mapping validation results over the last N days.

    Optional query params:
      ?days=7          window length (default 7)
      ?failures=true   include the most recent failure records (capped at 50)

    Returns:
      window_days
      total_attempts            (= 1st-attempts + retry-attempts)
      total_first_attempts
      first_attempt_pass_rate   (proportion of attempts=1 that passed)
      retry_count               (count of attempts=2)
      retry_pass_rate           (proportion of attempts=2 that passed)
      by_function               { fn: { attempts, first_pass, retries, retry_pass,
                                        avg_score, fabricated_total, unattributed_total,
                                        gap_total } }
      failure_breakdown         { reason_prefix: count }
      recent_failures           (only if failures=true)
    """
    from endo_ai import _EVMAP_LOG_PATH
    from datetime import timedelta
    import json

    try:
        days = int(request.args.get("days", "7"))
    except ValueError:
        days = 7
    include_failures = request.args.get("failures", "").lower() in ("1", "true", "yes")
    cutoff = datetime.now() - timedelta(days=days)

    if not os.path.exists(_EVMAP_LOG_PATH):
        return jsonify({
            "window_days": days, "total_attempts": 0, "total_first_attempts": 0,
            "first_attempt_pass_rate": None, "retry_count": 0, "retry_pass_rate": None,
            "by_function": {}, "failure_breakdown": {},
            "note": "evidence_mapping.jsonl does not exist yet — no validations logged.",
        })

    by_function: dict = {}
    failure_breakdown: dict = {}
    recent_failures: list = []
    total_attempts = 0
    first_attempts_total = 0
    first_attempts_passed = 0
    retry_count = 0
    retry_passed = 0

    with open(_EVMAP_LOG_PATH, encoding="utf-8") as fh:
        for line in fh:
            line = line.strip()
            if not line:
                continue
            try:
                rec = json.loads(line)
                ts  = datetime.fromisoformat(rec.get("ts", ""))
            except Exception:
                continue
            if ts < cutoff:
                continue

            attempt = int(rec.get("attempt", 1))
            passed  = bool(rec.get("passed"))
            score   = float(rec.get("score") or 0.0)
            fn      = rec.get("function", "unknown")
            reason  = rec.get("failure_reason") or ""

            total_attempts += 1
            if attempt == 1:
                first_attempts_total += 1
                if passed:
                    first_attempts_passed += 1
            elif attempt == 2:
                retry_count += 1
                if passed:
                    retry_passed += 1

            f = by_function.setdefault(fn, {
                "attempts": 0, "first_pass": 0, "retries": 0, "retry_pass": 0,
                "score_sum": 0.0, "fabricated_total": 0,
                "unattributed_total": 0, "gap_total": 0,
            })
            f["attempts"] += 1
            f["score_sum"] += score
            f["fabricated_total"]   += int(rec.get("n_fabricated") or 0)
            f["unattributed_total"] += int(rec.get("n_unattributed") or 0)
            f["gap_total"]          += int(rec.get("n_gap_sections") or 0)
            if attempt == 1 and passed:
                f["first_pass"] += 1
            if attempt == 2:
                f["retries"] += 1
                if passed:
                    f["retry_pass"] += 1

            if reason:
                # Prefix = part before first colon, e.g. "FABRICATED_PMIDS"
                prefix = reason.split(":", 1)[0].strip()
                failure_breakdown[prefix] = failure_breakdown.get(prefix, 0) + 1
                if include_failures and len(recent_failures) < 50:
                    recent_failures.append({
                        "ts":               rec.get("ts"),
                        "function":         fn,
                        "attempt":          attempt,
                        "score":            score,
                        "failure_reason":   reason,
                        "fabricated_pmids": rec.get("fabricated_pmids") or [],
                        "gap_sections":     rec.get("gap_sections") or [],
                        "unattributed_sample": rec.get("unattributed_sample") or [],
                    })

    # Finalise per-function averages
    for k, v in by_function.items():
        v["avg_score"] = round(v["score_sum"] / max(v["attempts"], 1), 1)
        del v["score_sum"]

    response = {
        "window_days":              days,
        "total_attempts":           total_attempts,
        "total_first_attempts":     first_attempts_total,
        "first_attempt_pass_rate":  (round(first_attempts_passed / first_attempts_total, 3)
                                      if first_attempts_total else None),
        "retry_count":              retry_count,
        "retry_pass_rate":          (round(retry_passed / retry_count, 3)
                                      if retry_count else None),
        "by_function":              by_function,
        "failure_breakdown":        failure_breakdown,
    }
    if include_failures:
        response["recent_failures"] = recent_failures
    return jsonify(response)


# ── History ──────────────────────────────────────────────

@app.route("/history")
def history():
    """Return recent cached queries for the history sidebar."""
    from rag import get_conn
    conn = get_conn()
    cur  = conn.cursor()
    try:
        cur.execute("""
            SELECT id, question_text, created_at, hit_count,
                   CASE WHEN jsonb_typeof(papers) = 'array'
                        THEN jsonb_array_length(papers) ELSE 0 END AS paper_count
            FROM query_cache
            ORDER BY created_at DESC
            LIMIT 50;
        """)
        rows = cur.fetchall()
        items = []
        for r in rows:
            qt = r[1] or ""
            # Parse mode prefix from cache key  [review]/[learn]/[case]
            mode_tag = "review"
            question = qt
            for tag in ("learn", "review", "case"):
                prefix = f"[{tag}] "
                if qt.startswith(prefix):
                    mode_tag = tag
                    question = qt[len(prefix):]
                    break
            items.append({
                "id":         r[0],
                # A15f.1 — the stored string keeps the clarification block
                # because it is the semantic cache key; the TITLE does not.
                "question":   display_title(question),
                "mode":       mode_tag,
                "created_at": r[2].isoformat() if r[2] else None,
                "hit_count":  r[3] or 0,
                # A19e — the drawer row shows what tells two entries apart.
                # The cost of a Review answer is not stored on the cache row,
                # so the count goes out alone rather than beside a made-up
                # figure; curricula carry their own measured cost.
                "paper_count": r[4] or 0,
            })
        return jsonify(items)
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        cur.close()
        conn.close()


@app.route("/history/<int:cache_id>")
def history_detail(cache_id: int):
    """Return full answer + papers for a cached history entry."""
    from rag import get_conn
    conn = get_conn()
    cur  = conn.cursor()
    try:
        cur.execute(
            "SELECT question_text, answer, papers FROM query_cache WHERE id = %s",
            (cache_id,)
        )
        row = cur.fetchone()
        if not row:
            return jsonify({"error": "Not found"}), 404
        qt, answer, papers = row
        papers = papers or []
        # A16b. A cache is a time capsule of the rendering that produced it,
        # and this route hands a stored answer straight to the browser. Every
        # Stage 1 fix that lives server-side — the impact-factor strip (Q3),
        # the quarantine block (Q2), the banner's second number (Q1/A3c) and
        # the cited-set bibliography (Q5) — reached /status and stopped here.
        #
        # Measured on the real rows: 6 of 10 carry an impact factor, 7 gain the
        # banner's second half, 1 gains a quarantine block, and 10 of 10 render
        # a bibliography of the whole retrieval pool instead of the citation
        # set. Re-rendering at read time is the option A16b prefers, because
        # every one of those fixes is presentational — nothing about the
        # stored answer is wrong, only how it was being shown.
        answer, _hq = finalise_answer_text(answer or "")
        cited = assemble_bibliography(answer, papers)["cited_pmids"]
        mode_tag = "review"
        question = qt or ""
        for tag in ("learn", "review", "case"):
            prefix = f"[{tag}] "
            if question.startswith(prefix):
                mode_tag = tag
                question = question[len(prefix):]
                break
        return jsonify({
            "question": display_title(question),
            "mode":     mode_tag,
            "answer":   answer or "",
            "cited_pmids": cited,
            # THE WHITELIST BELONGS HERE TOO. This was the one route that
            # serialised paper dicts without it, and it was harmless only by
            # accident: nothing in `query_cache.papers` had ever carried an
            # abstract, because the live path keeps abstracts in a side map
            # and never in the scored dict. The moment the library path
            # started carrying them, every library answer cached from then on
            # would have handed full abstract text to any client opening it
            # from history. "Abstracts stay server-side" has to be enforced at
            # every exit, not at the ones that happened to be tested.
            "papers":   _safe_papers(papers or []),
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        cur.close()
        conn.close()


# ── Case Discussion ───────────────────────────────────────

@app.route("/case_chat", methods=["POST"])
def case_chat():
    data         = request.json or {}
    messages     = data.get("messages", [])
    conv_id      = data.get("conv_id") or str(uuid.uuid4())
    skip_clarify = data.get("skip_clarify", False)
    question     = (messages[0]["content"] if messages else "").strip()
    if not question:
        return jsonify({"error": "No message provided"}), 400

    # Clarify gate — first message only, not a follow-up in an ongoing chat.
    # Deliberately NOT generate_clarifying_questions: that one is shared with
    # Review and asks 2-3 questions on principle, which here produced an
    # interrogation that re-asked facts the clinician had already written.
    # generate_case_followups re-reads the description first and returns [] when
    # it is already sufficient, so a complete case goes straight to an answer.
    if not skip_clarify and len(messages) == 1:
        try:
            from endo_ai import generate_case_followups
            questions = generate_case_followups(question)
            if questions:
                return jsonify({"needs_clarification": True, "questions": questions})
        except Exception:
            pass

    job_id = create_job(question, mode="case")
    thread = threading.Thread(
        target=run_case_chat,
        args=(job_id, messages, conv_id),
        daemon=True,
    )
    thread.start()
    return jsonify({"job_id": job_id, "conv_id": conv_id})


def run_case_chat(job_id: str, messages: list, conv_id: str):
    """`conv_id` is the CLIENT's conversation identity and is deliberately
    unused here.

    There used to be a `case_convs[conv_id] = {"evidence": evidence}` store
    below. Nothing ever read it, and nothing usefully could:

    - the whole conversation arrives in `messages` on every turn, and the
      client re-sends it in full (`templates/index.html`, `caseMessages`), so
      the server has no conversation state to remember;
    - the evidence base is rebuilt per turn from a query that combines the case
      with the LATEST follow-up, so a cached turn-1 evidence base would answer
      turn 3 from the wrong literature — the value was not merely unread, it
      would have been wrong to read;
    - the one plausible consumer, a sources panel, is already served by
      `update_job(..., papers=...)` below, which the case poller already
      fetches, exactly as Review mode does it.

    Meanwhile it retained a full evidence base — annotated abstracts included,
    ~277 KB — per CLIENT-SUPPLIED `conv_id`, with no cap and no eviction,
    unlike `review_threads` (`REVIEW_THREADS_MAX`). Deleted rather than wired
    up. The parameter and the API field stay: the client owns that identifier.
    """
    try:
        original_q = messages[0]["content"] if messages else ""
        # Latest user message drives THIS turn's literature search
        latest_user = next((m["content"] for m in reversed(messages)
                            if (m or {}).get("role") == "user"), original_q)
        is_followup = len(messages) > 1

        # Search query: combine original case context with the latest follow-up
        # so vague follow-ups ("what about MTA?") still hit relevant literature.
        if is_followup and latest_user.strip().lower() != original_q.strip().lower():
            search_q = f"{original_q} -- {latest_user}"
        else:
            search_q = original_q

        # ── Is this turn asking WHY, or asking WHAT TO DO? ──
        # "What could the cause be?" and "how should I manage this?" are
        # different questions and were being answered by one pipeline.
        # Fails open to treatment, which is the path that shipped and is
        # measured; see `classify_case_intent`.
        from endo_ai import (classify_case_intent, generate_case_differential,
                             CASE_INTENT_DIAGNOSTIC)
        intent = classify_case_intent(original_q, latest_user)
        print(f"\n[case] intent: {intent}")
        update_job(job_id, case_intent=intent)

        differential, diff_cost = [], 0.0
        if intent == CASE_INTENT_DIAGNOSTIC:
            update_job(job_id, message="Working out the differential...",
                       progress=8)
            differential, diff_cost = generate_case_differential(original_q,
                                                                 latest_user)
            if not differential:
                # LOUD, and retried once. An empty differential on a turn the
                # router called DIAGNOSTIC silently produces a treatment-shaped
                # answer — which is the failure this whole path exists to fix,
                # arriving through the back door. It happened for real: a
                # `max_tokens` truncation returned [], and the fixture case was
                # answered "Proceed with non-surgical root canal treatment".
                print("  [differential] EMPTY on a diagnostic turn — retrying "
                      "once before falling back to the treatment path")
                retry, retry_cost = generate_case_differential(original_q,
                                                               latest_user)
                diff_cost += retry_cost
                differential = retry
                if not differential:
                    print("  [differential] STILL EMPTY — this turn will be "
                          "answered on the treatment path, which is not what "
                          "the clinician asked for")
            for c in differential:
                print(f"  [differential] candidate: {c['candidate']}")
            # Published so the UI can show what is being searched for, and so a
            # trace can read the differential without re-deriving it. The
            # CANDIDATES only — never the evidence base. A full evidence base on
            # the job dict is how `case_convs` came to retain ~277 KB of
            # annotated abstracts per client-supplied id, and invariant 13 says
            # abstract text never reaches a browser.
            update_job(job_id, differential=[
                {k: c.get(k) for k in
                 ("candidate", "supports", "against", "discriminator")}
                for c in differential])

        # ── Carry the previous turn's papers as CANDIDATES ──
        # `case-v3` Item E. A follow-up rebuilt its evidence base from
        # scratch, so the papers the clinician was just reading about were
        # re-found or not depending on how the combined query embedded, and
        # the continuity between turns lived only in the prose. These seed
        # retrieval and never bypass it: `build_evidence_base_with_progress`
        # adds them after the routing gate has decided, and every gate then
        # applies — the similarity floor recomputed against THIS turn's
        # question, tier banding, and the retracted/superseded exclusions.
        #
        # NOT a cache. The evidence base is still rebuilt every turn and no
        # answer is ever reused; only the candidate set carries over. Caching
        # a turn-1 evidence base and serving it at turn 3 is the mistake
        # `case_convs` made, and it would answer the follow-up from the wrong
        # literature.
        from endo_ai import case_prior_pmids
        prior_pmids = case_prior_pmids(messages) if is_followup else []
        if prior_pmids:
            print(f"  [case] carrying {len(prior_pmids)} paper(s) cited by the "
                  f"earlier turn(s) as candidates")

        if differential:
            # One retrieval per candidate, unioned. The candidates share the
            # library gate, so a candidate the library already covers costs no
            # PubMed traffic.
            evidence = build_differential_evidence(job_id, original_q,
                                                   differential,
                                                   prior_pmids=prior_pmids)
        else:
            update_job(job_id,
                       message=("Searching literature for this question..."
                                if is_followup else
                                "Searching evidence base for this case..."),
                       progress=10)
            # mode="case", not the default "review". The review-mode early stop
            # skips level2-level5 and invitro once cochrane+level1 clear 15
            # papers, and those are exactly the tiers a case discussion needs —
            # a case series is often the only literature on an unusual
            # presentation. Measured before this: case answers cited a median
            # of 2 papers from a median-100 evidence base.
            evidence = build_evidence_base_with_progress(
                job_id, search_q, mode="case", prior_pmids=prior_pmids)

        if is_aborted(job_id):
            update_job(job_id, status="aborted", progress=100, message="Cancelled")
            return

        # ── Stream, then check (`case-v3` Item E) ──
        # A case turn's wall time is dominated by synthesis and the
        # post-checks, not retrieval — retrieval is seconds and pennies. The
        # clinician used to watch a spinner through all of it. Now the answer
        # arrives as it is written and the chips say "checking…" for exactly
        # the window between the model stopping and the guardrails finishing.
        #
        # The papers are published BEFORE synthesis so the [[PMID:N]] pills
        # rendered mid-stream resolve to author names rather than bare
        # numbers. `_safe_papers` still strips the abstracts on the way out.
        update_job(job_id,
                   message   = "Asking Claude…",
                   progress  = 80,
                   papers    = evidence.get("_summary", {}).get("all_scored", []),
                   streaming = True,
                   partial_answer = "",
                   checks_status  = "pending")

        def _on_partial(text: str):
            # Raw, UNCHECKED model text. It goes to `partial_answer` — never
            # to `answer` — so nothing downstream can mistake it for a
            # validated result.
            update_job(job_id,
                       partial_answer = text,
                       streaming      = True,
                       checks_status  = "pending",
                       message        = "Writing the answer…",
                       progress       = min(95, 80 + len(text) // 400))

        def _on_phase(_label: str):
            update_job(job_id, streaming=False, checks_status="pending",
                       progress=97,
                       message="Checking citations against the abstracts…")

        from endo_ai import ask_case_question
        try:
            answer, cost = ask_case_question(
                messages, evidence,
                differential = differential,
                stream_cb    = _on_partial,
                abort_cb     = lambda: is_aborted(job_id),
                phase_cb     = _on_phase)
        except StreamAborted:
            update_job(job_id, status="aborted", progress=100,
                       message="Cancelled", streaming=False, partial_answer="")
            return
        cost += diff_cost

        papers = evidence.get("_summary", {}).get("all_scored", [])
        update_job(
            job_id,
            status   = "complete",
            progress = 100,
            message  = "Done",
            answer   = answer,
            papers   = papers,
            images   = [],
            cost_usd = round(cost, 4),
            # The guardrails have finished; the chips can stop saying
            # "checking…". `partial_answer` is cleared so nothing downstream
            # can read UNCHECKED text once a checked answer exists.
            streaming      = False,
            partial_answer = "",
            checks_status  = "complete",
        )

        # The case discussion joins the History sidebar. The read path was
        # built for this — `/history`, `/history/<cache_id>` and the browser's
        # `loadHistoryItem` all parse the `[case] ` prefix — and nothing ever
        # wrote it, so Case was the one mode whose answers vanished when the
        # tab was closed.
        #
        # HISTORY, NOT CACHE: `get_cached_answer` excludes every `[case] ` row,
        # so this can never be served back to another patient's question
        # (invariant 21). The write is deliberately AFTER `update_job` — a
        # history failure must not cost the clinician the answer already on
        # screen, so it cannot take the job down with it.
        try:
            from rag import save_case_history
            save_case_history(conv_id, original_q, answer, papers)
        except Exception as e:
            print(f"  [case] history save failed (answer unaffected): {e}")
    except Exception as e:
        update_job(job_id, status="error", progress=100, error=str(e), message=str(e))


# ── Audio Export ──────────────────────────────────────────

# Caps on client-supplied export text. The export endpoints accept an answer
# from the browser when no live job matches (history-loaded answers have no
# server-side job), and that text is fed to a paid TTS API — an unbounded field
# there is a cheap way to run up a bill or tie up the single worker process.
MAX_EXPORT_ANSWER_CHARS   = 200_000
MAX_EXPORT_QUESTION_CHARS = 2_000


class ExportSourceTooLarge(ValueError):
    """Raised by _resolve_export_source; endpoints turn it into a 413."""


def _resolve_export_source(data: dict):
    """Return (question, answer) for an export request.

    Exports used to require a live job in THIS process's memory, which broke
    the moment an answer was loaded from the history sidebar or the server
    restarted: the audio/slides/video buttons 404'd ("Job not found") or,
    worse, exported a PREVIOUS question's answer through a stale job id. The
    client already holds the rendered answer, so it now sends question+answer
    with the request and a job id is just the preferred source, not the only
    one.
    """
    job_id = data.get("job_id", "")
    with jobs_lock:
        job = jobs.get(job_id)
    if job and job.get("answer"):
        return job.get("question", ""), job["answer"]

    # Client-supplied fallback. This text is untrusted input that flows into a
    # paid TTS pipeline, so it is capped: real answers measured across the
    # library run 7.7k-11.5k chars, and the cap is ~17x the largest observed.
    # Truncating instead of rejecting would silently narrate half an answer,
    # so an oversized body is refused outright.
    answer = (data.get("answer") or "").strip()
    if len(answer) > MAX_EXPORT_ANSWER_CHARS:
        raise ExportSourceTooLarge(
            f"answer is {len(answer):,} characters; the limit is "
            f"{MAX_EXPORT_ANSWER_CHARS:,}")
    question = (data.get("question") or "").strip()[:MAX_EXPORT_QUESTION_CHARS]
    if answer:
        return question, answer
    return None, None


@app.route("/generate_audio", methods=["POST"])
def generate_audio_endpoint():
    if not TTS_AVAILABLE:
        return jsonify({"error": "No TTS backend available"}), 503

    data           = request.json or {}
    job_id         = data.get("job_id", "")
    length_minutes = int(data.get("length_minutes", 10))
    length_minutes = max(5, min(60, length_minutes))
    voice          = (data.get("voice") or "").strip()   # "" -> narration.resolve_voice()
    style          = data.get("style", "lecture")   # "lecture" | "conversation"

    try:
        src_question, src_answer = _resolve_export_source(data)
    except ExportSourceTooLarge as e:
        return jsonify({"error": f"Answer too large to export: {e}"}), 413
    if not src_answer:
        return jsonify({"error": "Job not found or no answer to convert"}), 404

    audio_id = str(uuid.uuid4())
    with audio_jobs_lock:
        import time as _t_audio
        audio_jobs[audio_id] = {
            "status":         "running",
            "phase":          "script",     # "script" → "audio" → "complete"
            "file_path":      None,
            "script":         None,         # list of {host,text} for conversation
            "style":          style,
            "question":       src_question,
            "length_minutes": length_minutes,
            "error":          None,
            "started_at":     _t_audio.time(),
            "slides_done":    0, "slides_total": 0,
        }

    thread = threading.Thread(
        target=run_generate_audio,
        args=(audio_id, src_answer, src_question, length_minutes, voice, style),
        daemon=True,
    )
    thread.start()
    return jsonify({"audio_id": audio_id})


@app.route("/audio_status/<audio_id>")
def audio_status(audio_id: str):
    import time as _t_status
    with audio_jobs_lock:
        job = audio_jobs.get(audio_id)
        if job is not None:
            job_copy = dict(job)
        else:
            job_copy = None
    if not job_copy:
        return jsonify({"error": "Not found"}), 404

    elapsed = None
    eta     = None
    started = job_copy.get("started_at")
    done    = int(job_copy.get("slides_done", 0) or 0)
    total   = int(job_copy.get("slides_total", 0) or 0)
    if started:
        elapsed = max(0, int(_t_status.time() - started))
        if job_copy.get("status") in ("complete", "error", "cancelled"):
            eta = 0
        elif done > 0 and total > 0 and done < total:
            per_slide = elapsed / done
            eta = max(0, int(per_slide * (total - done)))
        elif total > 0:
            # No slide done yet -- rough heuristic: 8s per slide for content+TTS+encode
            eta = max(0, int(8 * total) - elapsed)

    return jsonify({
        "status":       job_copy["status"],
        "phase":        job_copy.get("phase", "script"),
        "style":        job_copy.get("style", "lecture"),
        "script":       job_copy.get("script"),
        "error":        job_copy.get("error"),
        "slides_done":  done,
        "slides_total": total,
        "elapsed_seconds": elapsed,
        "eta_seconds":     eta,
    })


@app.route("/audio_download/<audio_id>")
def audio_download(audio_id: str):
    with audio_jobs_lock:
        job = audio_jobs.get(audio_id)
    if not job or job["status"] != "complete" or not job.get("file_path"):
        return jsonify({"error": "Audio not ready"}), 404
    style = job.get("style", "lecture")
    fname = f"endo_ai_podcast_{audio_id[:8]}.mp3" if style == "conversation" else f"endo_ai_lecture_{audio_id[:8]}.mp3"
    return send_file(
        job["file_path"],
        as_attachment=True,
        download_name=fname,
        mimetype="audio/mpeg",
    )


def run_generate_audio(audio_id: str, answer: str, question: str,
                       length_minutes: int, voice: str = "",
                       style: str = "lecture"):
    try:
        # ── CONVERSATION style (two-host podcast) ───────────────
        if style == "conversation":
            with audio_jobs_lock:
                audio_jobs[audio_id]["phase"]  = "script"
                audio_jobs[audio_id]["status"] = "generating_script"

            print(f"  Generating {length_minutes}-min conversation script (DR. CHEN / ALEX)...")
            lines = generate_podcast_script(answer, question, length_minutes)
            print(f"  Script: {len(lines)} lines")

            with audio_jobs_lock:
                audio_jobs[audio_id]["script"] = lines   # expose immediately for transcript
                audio_jobs[audio_id]["phase"]  = "audio"
                audio_jobs[audio_id]["status"] = "converting_to_audio"

            tmp = tempfile.NamedTemporaryFile(suffix=".mp3", delete=False)
            tmp.close()

            if OPENAI_TTS_AVAILABLE:
                # Two voices are deliberate here, so this cannot use
                # resolve_voice() — but everything else the other exports get
                # applies: the pronunciation dictionary, tts-1-hd, sentence
                # splitting instead of a silent notes[:4096] truncation, and a
                # single cost row. Without it the podcast was the last export
                # still saying "er cr ysgg".
                HOST1_VOICE = "onyx"   # DR. CHEN
                HOST2_VOICE = "nova"   # ALEX
                audio_bytes = b""
                spoken_chars = 0
                used_model = narration.resolve_model(None)
                for i, line in enumerate(lines):
                    host  = line.get("host", "")
                    text  = (line.get("text") or "").strip()
                    if not text:
                        continue
                    v = HOST1_VOICE if "CHEN" in host.upper() else HOST2_VOICE
                    try:
                        seg = narration.synthesize_segment(
                            text, voice=v, model=used_model,
                            label=f"{host} {i+1}/{len(lines)}",
                            allow_gtts=False)
                        audio_bytes += seg["audio"]
                        spoken_chars += seg["characters"]
                        print(f"    [{i+1}/{len(lines)}] {host} OK ({v})")
                    except Exception as tts_err:
                        print(f"    [{i+1}/{len(lines)}] TTS error: {tts_err}")
                with open(tmp.name, "wb") as f:
                    f.write(audio_bytes)
                if spoken_chars:
                    narration.log_narration_cost(
                        "run_generate_audio", used_model, spoken_chars,
                        request_id=audio_id,
                        voice=f"{HOST1_VOICE}+{HOST2_VOICE}")
            elif GTTS_AVAILABLE:
                # Fallback: concatenate all text as one block
                full_text = " ".join(l.get("text", "") for l in lines)
                tts = gTTS(text=full_text[:5000], lang="en", slow=False)
                tts.save(tmp.name)
            else:
                raise RuntimeError("No TTS backend available")

            with audio_jobs_lock:
                audio_jobs[audio_id]["status"]    = "complete"
                audio_jobs[audio_id]["file_path"] = tmp.name
                q   = audio_jobs[audio_id].get("question", "")
                dur = audio_jobs[audio_id].get("length_minutes", 10)
            _persist_media(tmp.name, audio_id, "mp3", q, "conversation", "audio", dur)
            print(f"  Podcast audio saved: {tmp.name}")
            return

        # ── LECTURE style (single-voice, default) ───────────────
        with audio_jobs_lock:
            audio_jobs[audio_id]["phase"]  = "script"
            audio_jobs[audio_id]["status"] = "generating_script"

        print(f"  Generating {length_minutes}-min lecture script...")
        script = generate_audio_script(answer, question, length_minutes)
        words  = len(script.split())
        print(f"  Script: {words} words")

        with audio_jobs_lock:
            audio_jobs[audio_id]["phase"]  = "audio"
            audio_jobs[audio_id]["status"] = "converting_to_audio"

        tmp = tempfile.NamedTemporaryFile(suffix=".mp3", delete=False)
        tmp.close()

        # One call covers WORKLIST 4.1/4.2/4.3/4.5: pronunciation dictionary
        # and marker stripping, OpenAI-primary with a LOUD gTTS fallback, the
        # sidecar timestamp map the web deck reads, and the per-character cost
        # row. The sidecar lands beside the mp3 under the same stem, so
        # _persist_media needs no change.
        narration.synthesize_lecture(
            script, tmp.name, audio_id=audio_id, voice=voice,
            style="lecture", mode="export",
            function_name="run_generate_audio")

        with audio_jobs_lock:
            audio_jobs[audio_id]["status"]    = "complete"
            audio_jobs[audio_id]["file_path"] = tmp.name
            q   = audio_jobs[audio_id].get("question", "")
            dur = audio_jobs[audio_id].get("length_minutes", 10)
        _persist_media(tmp.name, audio_id, "mp3", q, "lecture", "audio", dur)
        print(f"  Audio saved: {tmp.name}")

    except Exception as e:
        print(f"  Audio generation error: {e}")
        with audio_jobs_lock:
            audio_jobs[audio_id]["status"] = "error"
            audio_jobs[audio_id]["error"]  = str(e)


# ── PPTX Slides ──────────────────────────────────────────

@app.route("/generate_slides", methods=["POST"])
def generate_slides_endpoint():
    if not PPTX_AVAILABLE:
        return jsonify({"error": "python-pptx not installed"}), 503

    data           = request.json or {}
    job_id         = data.get("job_id", "")
    length_minutes = int(data.get("length_minutes", 10))
    length_minutes = max(5, min(60, length_minutes))
    voice          = (data.get("voice") or "").strip()   # "" -> narration.resolve_voice()

    try:
        src_question, src_answer = _resolve_export_source(data)
    except ExportSourceTooLarge as e:
        return jsonify({"error": f"Answer too large to export: {e}"}), 413
    if not src_answer:
        return jsonify({"error": "Job not found or no answer"}), 404

    audio_id = str(uuid.uuid4())
    with audio_jobs_lock:
        import time as _t_slides
        audio_jobs[audio_id] = {
            "status": "running", "file_path": None, "error": None, "type": "pptx",
            "question": src_question, "length_minutes": length_minutes,
            "started_at": _t_slides.time(),
            "slides_done": 0, "slides_total": 0,
        }

    thread = threading.Thread(
        target=run_generate_slides,
        args=(audio_id, src_answer, src_question, length_minutes, voice,
              data.get("papers") or []),
        daemon=True,
    )
    thread.start()
    return jsonify({"audio_id": audio_id})


@app.route("/slides_download/<audio_id>")
def slides_download(audio_id: str):
    with audio_jobs_lock:
        job = audio_jobs.get(audio_id)
    if not job or job["status"] != "complete" or not job.get("file_path"):
        return jsonify({"error": "Slides not ready"}), 404
    return send_file(
        job["file_path"],
        as_attachment=True,
        download_name=f"endo_ai_slides_{audio_id[:8]}.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# ── Web deck export (PRESENTATION_WORKLIST §3) ────────────
# Same shape as the three exports above, and deliberately the same gating:
# NONE. The other export routes are ungated on purpose — exports are a user
# feature, an unset ADMIN_TOKEN would 403 them by design (HANDOVER.md bug
# class (d)), and /ask, the expensive route, is ungated too. What this route
# DOES inherit is the size cap, via _resolve_export_source.
MAX_EXPORT_PAPERS = 400


def _sanitize_export_papers(raw) -> list:
    """Client-supplied paper metadata for a history-loaded answer.

    Only the fields the deck actually renders survive, values are capped, and
    the list is bounded — this is untrusted input that ends up inside a file
    served from the app's origin.
    """
    allowed = {"pmid", "title", "authors", "year", "journal", "journal_abbrev",
               "volume", "issue", "pages", "sample_size", "level_key", "score",
               "has_retraction"}
    out = []
    for p in (raw or [])[:MAX_EXPORT_PAPERS]:
        if not isinstance(p, dict) or not str(p.get("pmid") or "").isdigit():
            continue
        row = {}
        for k in allowed:
            v = p.get(k)
            if isinstance(v, str):
                row[k] = v[:400]
            elif isinstance(v, (int, float, bool)) or v is None:
                row[k] = v
        out.append(row)
    return out


def _embed_abstracts(pmids) -> dict:
    """Abstracts baked into the deck so it works with the server stopped (§3.2).

    L1 (in-process) + L2 (Postgres abstract_cache) only. The live eutils path
    that /api/abstract falls back to is deliberately NOT used here: 40 serial
    rate-limited fetches would dominate the export, and a PMID missing at build
    time is exactly the case the deck's own server-first lookup covers.
    """
    out = {}
    for pmid in pmids:
        rec = _ABSTRACT_CACHE.get(pmid)
        if not rec:
            try:
                cached = get_cached_abstract(pmid)
            except Exception as e:
                print(f"  [webdeck] abstract lookup failed for {pmid}: {e}")
                continue
            if not (cached and (cached.get("abstract") or "")):
                continue
            rec = {"pmid": pmid,
                   "title": (cached.get("title") or "").rstrip("."),
                   "abstract": cached.get("abstract") or "",
                   "journal": cached.get("journal") or "",
                   "year": cached.get("year") or "",
                   "authors": cached.get("authors") or ""}
        if (rec.get("abstract") or "").strip():
            out[str(pmid)] = {k: rec.get(k, "") for k in
                              ("pmid", "title", "abstract", "journal",
                               "year", "authors")}
    return out


def _find_narration_audio_id(question: str, length_minutes: int) -> str:
    """Newest audio render of the same question at the same length.

    The timestamp sidecar (narration.build_timestamp_map) records the audio_id,
    voice and duration but nothing that identifies WHICH answer was narrated,
    so this is the only available link and it is deliberately narrow: a
    different length is a different script, and a different question is a
    different deck. No match simply means a deck without audio.
    """
    q = (question or "").strip()
    if not q:
        return ""
    for item in _load_media_index():          # newest first
        if item.get("type") != "audio":
            continue
        if (item.get("question") or "").strip() != q:
            continue
        if int(item.get("length_minutes") or 0) != int(length_minutes):
            continue
        if os.path.exists(os.path.join(MEDIA_DIR,
                                       f"{item['id']}.timestamps.json")):
            return item["id"]
    return ""


def build_deck_narration_sections(slides: list) -> list:
    """One narration section per SPEC slide, in spec order.

    `webdeck.narration.load_narration` arms auto-advance only when the
    sidecar's segment count equals the spec slide count, so the 1:1 is the
    whole point and a dropped slide silently disarms it.
    `narration.synthesize_lecture` discards a section whose text is empty,
    which is why a slide with no speaker notes falls back to its title and
    then to a bare placeholder: a slide the deck SHOWS must be spoken for, or
    the audio and the slides drift by one from that point on.
    """
    out = []
    for i, slide in enumerate(slides or []):
        title = (slide.get("title") or "").replace("\n", " ").strip()
        text  = (slide.get("speaker_notes") or "").strip()
        if not text:
            text = title or f"Slide {i + 1}."
        out.append({"title": title or f"Slide {i + 1}", "text": text})
    return out


def _build_synced_narration(audio_id: str, slides: list, question: str,
                            length_minutes: int, voice: str = None) -> str:
    """Record narration cut per spec slide. Returns the audio_id, or "".

    Written under the DECK's own id rather than the audio export's: this
    soundtrack is cut to this spec and is meaningless against any other, and
    `find_sidecar` refuses to guess which render belongs to which answer.

    Every failure returns "" and the deck falls back to whatever unsynced
    render exists — §3.3's "graceful without audio". A deck that builds
    without a soundtrack beats an export that dies.
    """
    sections = build_deck_narration_sections(slides)
    if not sections:
        return ""
    out_path = os.path.join(MEDIA_DIR, f"{audio_id}.mp3")
    try:
        os.makedirs(MEDIA_DIR, exist_ok=True)
        summary = narration.synthesize_lecture(
            "", out_path, audio_id=audio_id,
            voice=voice, sections=sections, style="deck",
            function_name="run_generate_webdeck",
            media_dir=MEDIA_DIR, per_section=True)
    except Exception as e:
        print(f"  [webdeck] per-slide narration failed ({e}); "
              f"falling back to any existing render")
        return ""

    tmap = summary.get("timestamp_map") or {}
    got  = len(tmap.get("slides") or [])
    if got != len(sections):
        # Not a warning: a map that does not describe these slides is exactly
        # what auto-advance must not be armed on, and using it would advance
        # the deck to the wrong slide for the sentence being spoken.
        print(f"  [webdeck] per-slide map has {got} segment(s) for "
              f"{len(sections)} slide(s) — not using it")
        return ""
    print(f"  [webdeck] narration cut per slide: {got} segment(s), "
          f"{summary.get('duration_seconds', 0):.1f}s, "
          f"${summary.get('cost_usd', 0):.4f}")
    return audio_id


@app.route("/generate_webdeck", methods=["POST"])
def generate_webdeck_endpoint():
    data           = request.json or {}
    job_id         = data.get("job_id", "")
    length_minutes = int(data.get("length_minutes", 10))
    length_minutes = max(5, min(60, length_minutes))
    # "auto"      — reuse a synced render; otherwise record one per slide.
    # "reuse"     — the pre-2026-09-01 behaviour: attach whatever render exists
    #               for this question and leave auto-advance off if it does not
    #               match. Costs nothing.
    # "per_slide" — always record, even over a render that already matches.
    # "off"       — no audio at all.
    narrate = str(data.get("narrate") or "auto").lower()
    if narrate not in ("auto", "reuse", "per_slide", "off"):
        narrate = "auto"

    try:
        src_question, src_answer = _resolve_export_source(data)
    except ExportSourceTooLarge as e:
        return jsonify({"error": f"Answer too large to export: {e}"}), 413
    if not src_answer:
        return jsonify({"error": "Job not found or no answer"}), 404

    # Paper metadata drives the mandatory evidence-shape card and the
    # references slides. Prefer the live job; fall back to what the browser
    # holds, the same way the answer text itself does.
    with jobs_lock:
        job = jobs.get(job_id)
    papers = list((job or {}).get("papers") or [])
    if not papers:
        papers = _sanitize_export_papers(data.get("papers"))

    audio_id = str(uuid.uuid4())
    import time as _t_deck
    with audio_jobs_lock:
        audio_jobs[audio_id] = {
            "status": "running", "file_path": None, "error": None,
            "file_ext": "html", "type": "webdeck",
            "question": src_question, "length_minutes": length_minutes,
            "started_at": _t_deck.time(),
            "slides_done": 0, "slides_total": 0,
        }

    thread = threading.Thread(
        target=run_generate_webdeck,
        args=(audio_id, src_answer, src_question, length_minutes, papers,
              data.get("audio_id") or "", narrate, data.get("voice") or ""),
        daemon=True,
    )
    thread.start()
    return jsonify({"audio_id": audio_id})


def run_generate_webdeck(audio_id: str, answer: str, question: str,
                         length_minutes: int, papers: list,
                         narration_audio_id: str = "",
                         narrate: str = "auto", voice: str = ""):
    try:
        import slide_spec_cache
        from webdeck import build_web_deck, load_narration
        from webdeck.citations import extract_pmids

        with audio_jobs_lock:
            audio_jobs[audio_id]["status"] = "generating_content"

        # The canonical text object (§0 prime rule, §5.1): the PPTX export and
        # this one read the SAME spec, so their content hashes can be compared.
        spec, spec_hash, from_cache = slide_spec_cache.get_or_build(
            answer, question, length_minutes)
        slides = (spec or {}).get("slides") or []
        if not slides:
            raise ValueError("Slide generator returned 0 slides after retry.")
        print(f"  [webdeck] {len(slides)} slide specs "
              f"({'cached' if from_cache else 'generated'}), hash {spec_hash[:12]}")

        with audio_jobs_lock:
            audio_jobs[audio_id]["status"] = "building_slides"
            audio_jobs[audio_id]["slides_total"] = len(slides)

        cited = list(dict.fromkeys(extract_pmids(answer) + extract_pmids(slides)))
        abstracts = _embed_abstracts(cited)
        print(f"  [webdeck] embedded {len(abstracts)}/{len(cited)} abstracts")

        # §3.3 — attach an existing narration render for the SAME answer. The
        # sidecar carries no answer identity of its own, so the link is the
        # media index: same question, same length, most recent audio render.
        narr_id = "" if narrate == "off" else (
            narration_audio_id or _find_narration_audio_id(
                question, length_minutes))

        # The 13-vs-34 mismatch, closed. A lecture render is cut on the
        # SCRIPT's own structure — 13 sections for a 10-minute laser answer —
        # while this spec has 25 slides that the §1.3 body budget renders as
        # 34 sections. Those 13 boundaries describe nothing on screen, so
        # auto-advance stayed off and the deck said so. Recording the
        # narration against THIS spec, one segment per slide, is the only
        # thing that can arm it; nothing derivable from the lecture sidecar
        # can, because `char_start` indexes the spoken script and not the
        # answer.
        n_spec = len(slides)
        if narrate in ("auto", "per_slide") and narration.openai_available():
            existing = load_narration(MEDIA_DIR, narr_id,
                                      spec_slide_count=n_spec) if narr_id else None
            already_synced = bool((existing or {}).get("synced"))
            if narrate == "per_slide" or not already_synced:
                with audio_jobs_lock:
                    audio_jobs[audio_id]["status"] = "generating_audio"
                synced_id = _build_synced_narration(
                    audio_id, slides, question, length_minutes, voice or None)
                if synced_id:
                    narr_id = synced_id
        elif narrate in ("auto", "per_slide"):
            print("  [webdeck] no OpenAI TTS backend — keeping the existing "
                  "render; auto-advance stays off if it does not match")

        def _load_narration(spec_slide_count, spec_to_section):
            if not narr_id:
                return None
            return load_narration(MEDIA_DIR, narr_id,
                                  spec_slide_count=spec_slide_count,
                                  spec_to_section=spec_to_section)

        html_out = build_web_deck(spec, question, answer, papers_list=papers,
                                  abstracts=abstracts, spec_hash=spec_hash,
                                  narration_loader=_load_narration)

        tmp = tempfile.NamedTemporaryFile(suffix=".html", delete=False,
                                          mode="w", encoding="utf-8")
        tmp.write(html_out)
        tmp.close()

        with audio_jobs_lock:
            audio_jobs[audio_id]["status"] = "complete"
            audio_jobs[audio_id]["file_path"] = tmp.name
            audio_jobs[audio_id]["spec_hash"] = spec_hash
            q   = audio_jobs[audio_id].get("question", "")
            dur = audio_jobs[audio_id].get("length_minutes", 10)
        _persist_media(tmp.name, audio_id, "html", q, "webdeck", "webdeck", dur)
        print(f"  [webdeck] OK {len(html_out):,} bytes -> {tmp.name}")

    except Exception as e:
        import traceback; traceback.print_exc()
        with audio_jobs_lock:
            audio_jobs[audio_id]["status"] = "error"
            audio_jobs[audio_id]["error"]  = str(e)


def _webdeck_path(audio_id: str):
    with audio_jobs_lock:
        job = audio_jobs.get(audio_id)
    if job and job.get("status") == "complete" and job.get("file_path") \
            and os.path.exists(job["file_path"]):
        return job["file_path"]
    # Media-tab copy: survives the in-memory job (single worker, restarts).
    item = next((i for i in _load_media_index()
                 if i.get("id") == audio_id and i.get("ext") == "html"), None)
    if item:
        p = os.path.join(MEDIA_DIR, item["filename"])
        if os.path.exists(p):
            return p
    return None


@app.route("/webdeck_view/<audio_id>")
def webdeck_view(audio_id: str):
    """Served inline so the deck's citation pills can reach /api/abstract on
    the same origin. Every string the deck renders is escaped at build time —
    see webdeck/citations.py and the config-JSON escaping in builder.py."""
    path = _webdeck_path(audio_id)
    if not path:
        return jsonify({"error": "Web deck not ready"}), 404
    resp = send_file(path, mimetype="text/html")
    resp.headers["X-Content-Type-Options"] = "nosniff"
    return resp


@app.route("/webdeck_download/<audio_id>")
def webdeck_download(audio_id: str):
    path = _webdeck_path(audio_id)
    if not path:
        return jsonify({"error": "Web deck not ready"}), 404
    return send_file(path, as_attachment=True,
                     download_name=f"curo_deck_{audio_id[:8]}.html",
                     mimetype="text/html")


@app.route("/generate_video", methods=["POST"])
def generate_video_endpoint():
    if not MOVIEPY_AVAILABLE:
        return jsonify({"error": "moviepy not installed — run: pip install moviepy"}), 503

    data           = request.json or {}
    job_id         = data.get("job_id", "")
    length_minutes = int(data.get("length_minutes", 10))
    length_minutes = max(5, min(60, length_minutes))
    voice          = (data.get("voice") or "").strip()   # "" -> narration.resolve_voice()

    try:
        src_question, src_answer = _resolve_export_source(data)
    except ExportSourceTooLarge as e:
        return jsonify({"error": f"Answer too large to export: {e}"}), 413
    if not src_answer:
        return jsonify({"error": "Job not found or no answer"}), 404

    audio_id = str(uuid.uuid4())
    import time as _t_init
    with audio_jobs_lock:
        audio_jobs[audio_id] = {
            "status": "running", "file_path": None,
            "error": None, "file_ext": "mp4", "type": "video",
            "question": src_question, "length_minutes": length_minutes,
            "started_at": _t_init.time(),
            "slides_done": 0, "slides_total": 0,
        }

    thread = threading.Thread(
        target=run_generate_video,
        args=(audio_id, src_answer, src_question, length_minutes, voice),
        daemon=True,
    )
    thread.start()
    return jsonify({"audio_id": audio_id})


@app.route("/video_download/<audio_id>")
def video_download(audio_id: str):
    with audio_jobs_lock:
        job = audio_jobs.get(audio_id)
    if not job or job["status"] != "complete" or not job.get("file_path"):
        return jsonify({"error": "Video not ready"}), 404
    return send_file(
        job["file_path"],
        as_attachment=True,
        download_name=f"endo_ai_lecture_{audio_id[:8]}.mp4",
        mimetype="video/mp4",
    )


import base64 as _b64, zipfile as _zipfile, re as _re

# Tiny 1×1 transparent PNG — used as the audio-shape icon in each slide
_AUDIO_ICON_PNG = _b64.b64decode(
    'iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk'
    '+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=='
)
_NS_P    = 'http://schemas.openxmlformats.org/presentationml/2006/main'
_NS_A    = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_NS_R    = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
_NS_REL  = 'http://schemas.openxmlformats.org/package/2006/relationships'
_NS_CT   = 'http://schemas.openxmlformats.org/package/2006/content-types'
_AUDIO_REL = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio'
_MEDIA_REL = "http://schemas.microsoft.com/office/2007/relationships/media"
_IMAGE_REL = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image'


def _patch_slide_xml_zip(data: bytes, slide_num: int) -> bytes:
    """Inject a hidden autoplay audio shape into one slide's XML.

    Every element below is copied from what PowerPoint itself writes when you
    insert an audio file and tick "Start automatically". That ground truth was
    obtained by driving PowerPoint over COM and reading the saved XML, because
    three earlier hand-authored attempts all produced files PowerPoint refused
    to open with "corrupted and unreadable" — and the COM render path used for
    verification opened them anyway, so nothing caught it.

    The three things the hand-written version got wrong:

      * `<p:audioFile>` — no such element. Audio lives in the DrawingML
        namespace as `<a:audioFile>`. This was the fatal one.
      * `<p:audio>` was nested inside the click-sequence's `<p:par>` chain. It
        is a SIBLING of `<p:seq>`, both children of the root node's childTnLst.
      * `<p:cond evt="onPrevClick"><p:tn/></p:cond>` — the event names are
        `onPrev`/`onNext` and the target is `<p:tgtEl><p:sldTgt/></p:tgtEl>`;
        a bare `<p:tn/>` has no `val` and is invalid.

    PowerPoint also declares the mp3 TWICE: once as an `audio` relationship for
    `a:audioFile`, and once as a Microsoft `media` relationship referenced from
    a `p14:media` extension. Both are required for the media to play in modern
    PowerPoint, which is why `_patch_slide_rels_zip` writes two rels per slide.
    """
    xml = data.decode('utf-8')
    shape_id  = 900 + slide_num
    rId_audio = f'rIdAudio{slide_num}'      # .../relationships/audio
    rId_media = f'rIdMedia{slide_num}'      # microsoft .../relationships/media
    rId_icon  = f'rIdAIcon{slide_num}'

    pic = (
        f'<p:pic>'
        f'<p:nvPicPr>'
        f'<p:cNvPr id="{shape_id}" name="Narration {slide_num}">'
        f'<a:hlinkClick r:id="" action="ppaction://media"/>'
        f'</p:cNvPr>'
        f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
        f'<p:nvPr>'
        f'<a:audioFile r:link="{rId_audio}"/>'
        f'<p:extLst>'
        f'<p:ext uri="{{DAA4B4D4-6D71-4841-9C94-3DE7FCFB9230}}">'
        f'<p14:media xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" '
        f'r:embed="{rId_media}"/>'
        f'</p:ext>'
        f'</p:extLst>'
        f'</p:nvPr>'
        f'</p:nvPicPr>'
        f'<p:blipFill>'
        f'<a:blip r:embed="{rId_icon}"/>'
        f'<a:stretch><a:fillRect/></a:stretch>'
        f'</p:blipFill>'
        f'<p:spPr>'
        f'<a:xfrm><a:off x="457200" y="457200"/><a:ext cx="457200" cy="457200"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'</p:spPr>'
        f'</p:pic>'
    )
    if '</p:spTree>' in xml:
        xml = xml.replace('</p:spTree>', pic + '</p:spTree>', 1)

    xml = _re.sub(r'<p:timing\b[^>]*>[\s\S]*?</p:timing>', '', xml)

    # Structure mirrors PowerPoint's own: a mainSeq holding one clickEffect
    # that issues playFrom(0.0), and a SIBLING p:audio media node.
    timing = (
        f'<p:timing><p:tnLst><p:par>'
        f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>'
        f'<p:seq concurrent="1" nextAc="seek">'
        f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
        f'<p:par><p:cTn id="3" fill="hold">'
        f'<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst><p:childTnLst>'
        f'<p:par><p:cTn id="4" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>'
        f'<p:par><p:cTn id="5" presetID="1" presetClass="mediacall" '
        f'presetSubtype="0" fill="hold" nodeType="clickEffect">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>'
        f'<p:cmd type="call" cmd="playFrom(0.0)"><p:cBhvr>'
        f'<p:cTn id="6" dur="indefinite" fill="hold"/>'
        f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
        f'</p:cBhvr></p:cmd>'
        f'</p:childTnLst></p:cTn></p:par>'
        f'</p:childTnLst></p:cTn></p:par>'
        f'</p:childTnLst></p:cTn></p:par>'
        f'</p:childTnLst></p:cTn>'
        f'<p:prevCondLst><p:cond evt="onPrev" delay="0">'
        f'<p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        f'<p:nextCondLst><p:cond evt="onNext" delay="0">'
        f'<p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        f'</p:seq>'
        f'<p:audio><p:cMediaNode vol="80000" showWhenStopped="0">'
        f'<p:cTn id="7" fill="hold" display="0">'
        f'<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>'
        f'<p:endCondLst><p:cond evt="onStopAudio" delay="0">'
        f'<p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:endCondLst>'
        f'</p:cTn>'
        f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
        f'</p:cMediaNode></p:audio>'
        f'</p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'
    )
    if '</p:sld>' in xml:
        xml = xml.replace('</p:sld>', timing + '</p:sld>', 1)
    return xml.encode('utf-8')


def _patch_slide_rels_zip(data: bytes, slide_num: int) -> bytes:
    """String-based injection of audio + icon relationships into slide rels XML."""
    xml = data.decode('utf-8')
    rId_audio = f'rIdAudio{slide_num}'
    rId_media = f'rIdMedia{slide_num}'
    rId_icon  = f'rIdAIcon{slide_num}'
    inserts = ''
    if rId_audio not in xml:
        inserts += (f'<Relationship Id="{rId_audio}" Type="{_AUDIO_REL}"'
                    f' Target="../media/narration_s{slide_num}.mp3"/>')
    # PowerPoint declares the same mp3 twice: the standard `audio` rel above
    # for a:audioFile, and this Microsoft `media` rel for the p14:media
    # extension. Modern PowerPoint needs both to play embedded media.
    if rId_media not in xml:
        inserts += (f'<Relationship Id="{rId_media}" Type="{_MEDIA_REL}"'
                    f' Target="../media/narration_s{slide_num}.mp3"/>')
    if rId_icon not in xml:
        inserts += (f'<Relationship Id="{rId_icon}" Type="{_IMAGE_REL}"'
                    f' Target="../media/audio_icon.png"/>')
    if inserts and '</Relationships>' in xml:
        xml = xml.replace('</Relationships>', inserts + '</Relationships>', 1)
    return xml.encode('utf-8')


def _inject_audio_into_pptx(pptx_path: str, slide_audios: dict) -> str:
    """
    Rewrite a saved PPTX, injecting per-slide MP3 narration via direct ZIP surgery.
    slide_audios: {slide_num_1based: mp3_bytes}
    Returns path to the new narrated PPTX (original is untouched).
    """
    import io as _sysio

    in_data = open(pptx_path, 'rb').read()
    out_buf = _sysio.BytesIO()

    with _zipfile.ZipFile(_sysio.BytesIO(in_data), 'r') as zin:
        patched_rels: set = set()

        with _zipfile.ZipFile(out_buf, 'w', _zipfile.ZIP_DEFLATED) as zout:

            # ── Patch [Content_Types].xml (string-based, no lxml) ─
            ct_xml = zin.read('[Content_Types].xml').decode('utf-8')
            for ext, ctype in [('mp3', 'audio/mpeg'), ('png', 'image/png')]:
                if f'Extension="{ext}"' not in ct_xml:
                    ct_xml = ct_xml.replace(
                        '</Types>',
                        f'<Default Extension="{ext}" ContentType="{ctype}"/></Types>', 1)
            zout.writestr('[Content_Types].xml', ct_xml.encode('utf-8'))

            # ── Process all existing ZIP members ───────────────
            for item in zin.infolist():
                fname = item.filename
                if fname == '[Content_Types].xml':
                    continue
                data = zin.read(fname)

                m = _re.match(r'ppt/slides/slide(\d+)\.xml$', fname)
                if m:
                    snum = int(m.group(1))
                    if snum in slide_audios:
                        try:
                            data = _patch_slide_xml_zip(data, snum)
                        except Exception as ex:
                            print(f"  WARN XML patch failed slide {snum}: {ex}")

                m = _re.match(r'ppt/slides/_rels/slide(\d+)\.xml\.rels$', fname)
                if m:
                    snum = int(m.group(1))
                    if snum in slide_audios:
                        try:
                            data = _patch_slide_rels_zip(data, snum)
                            patched_rels.add(snum)
                        except Exception as ex:
                            print(f"  WARN Rels patch failed slide {snum}: {ex}")

                zout.writestr(item, data)

            # ── Create rels files for slides that had none ─────
            for snum in slide_audios:
                if snum not in patched_rels:
                    zout.writestr(
                        f'ppt/slides/_rels/slide{snum}.xml.rels',
                        (f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                         f'<Relationships xmlns="{_NS_REL}">'
                         f'<Relationship Id="rIdAudio{snum}" Type="{_AUDIO_REL}"'
                         f' Target="../media/narration_s{snum}.mp3"/>'
                         f'<Relationship Id="rIdAIcon{snum}" Type="{_IMAGE_REL}"'
                         f' Target="../media/audio_icon.png"/>'
                         f'</Relationships>').encode('utf-8')
                    )

            # ── Add media files ────────────────────────────────
            zout.writestr('ppt/media/audio_icon.png', _AUDIO_ICON_PNG)
            for snum, mp3_bytes in slide_audios.items():
                zout.writestr(f'ppt/media/narration_s{snum}.mp3', mp3_bytes)

    out_path = pptx_path[:-5] + '_narrated.pptx'
    with open(out_path, 'wb') as f:
        f.write(out_buf.getvalue())
    return out_path


# ── PLACEHOLDER (replaced by ZIP injection below) ─────────
def _embed_slide_audio(slide, prs, mp3_bytes: bytes, slide_num: int) -> bool:
    """
    Embed MP3 narration into a slide with autoplay on slide entry.
    Uses a 1×1 invisible sp shape with sndAc + OOXML timing animation.
    Returns True on success.
    """
    try:
        from pptx.opc.part import Part
        from pptx.opc.packuri import PackURI
        from lxml import etree

        AUDIO_REL = ('http://schemas.openxmlformats.org/'
                     'officeDocument/2006/relationships/audio')
        NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

        shape_id = 900 + slide_num  # high id, no collision with content shapes

        # 1. Register audio part with the slide
        audio_part = Part(
            PackURI(f'/ppt/media/narration_s{slide_num}.mp3'),
            'audio/mpeg',
            mp3_bytes,
            prs.part.package,
        )
        rId = slide.part.relate_to(audio_part, AUDIO_REL)

        # 2. Invisible 1×1 EMU shape carrying the audio reference
        sp_xml = (
            f'<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}" xmlns:r="{NS_R}">'
            f'<p:nvSpPr>'
            f'<p:cNvPr id="{shape_id}" name="Narration {slide_num}"/>'
            f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
            f'<p:nvPr>'
            f'<p:sndAc><p:stSnd>'
            f'<p:snd r:embed="{rId}" name="narration{slide_num}.mp3"/>'
            f'</p:stSnd></p:sndAc>'
            f'</p:nvPr>'
            f'</p:nvSpPr>'
            f'<p:spPr>'
            f'<a:xfrm><a:off x="0" y="0"/><a:ext cx="1" cy="1"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            f'<a:noFill/>'
            f'</p:spPr>'
            f'<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>'
            f'</p:sp>'
        )
        slide.shapes._spTree.append(etree.fromstring(sp_xml))

        # 3. Autoplay timing — triggers audio at delay=0 (slide entry, no click needed)
        timing_xml = (
            f'<p:timing xmlns:p="{NS_P}" xmlns:a="{NS_A}">'
            f'<p:tnLst><p:par>'
            f'<p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">'
            f'<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
            f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
            f'<p:childTnLst><p:par>'
            f'<p:cTn id="3" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst><p:par>'
            f'<p:cTn id="4" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst><p:audio>'
            f'<p:cMediaNode vol="80000" mute="0" numSld="0" showWhenStopped="0">'
            f'<p:cTn id="5" fill="hold" display="0">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'</p:cTn>'
            f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
            f'</p:cMediaNode>'
            f'</p:audio></p:childTnLst>'
            f'</p:cTn>'
            f'</p:par></p:childTnLst>'
            f'</p:cTn>'
            f'</p:par></p:childTnLst>'
            f'</p:cTn>'
            f'<p:nav>'
            f'<p:prevCondLst><p:cond evt="onPrevClick" delay="0"><p:tn/></p:cond></p:prevCondLst>'
            f'<p:nextCondLst><p:cond evt="onNextClick" delay="0"><p:tn/></p:cond></p:nextCondLst>'
            f'</p:nav>'
            f'</p:seq></p:childTnLst>'
            f'</p:cTn>'
            f'</p:par></p:tnLst>'
            f'<p:bldLst>'
            f'<p:bldP spid="{shape_id}" grpId="0" uiExpand="1" build="p"/>'
            f'</p:bldLst>'
            f'</p:timing>'
        )
        slide._element.append(etree.fromstring(timing_xml))
        return True

    except Exception as exc:
        print(f"    WARN  Audio embed failed (slide {slide_num}): {exc}")
        return False


# ── Slide image renderer (Pillow) ────────────────────────

def _draw_wrapped_text(draw, text: str, x: int, y: int, max_w: int,
                       font, color, align: str = 'left') -> int:
    """Word-wrap text into max_w pixels. Returns new y after last line."""
    if not text:
        return y
    words = text.split()
    lines, cur = [], ''
    for w in words:
        test = (cur + ' ' + w).strip()
        try:
            tw = draw.textlength(test, font=font)
        except Exception:
            tw = len(test) * 12  # fallback
        if tw > max_w and cur:
            lines.append(cur); cur = w
        else:
            cur = test
    if cur:
        lines.append(cur)
    try:
        lh = draw.textbbox((0, 0), 'Ag', font=font)[3] + 8
    except Exception:
        lh = 28
    for line in lines:
        if align == 'center':
            try:
                lw = draw.textlength(line, font=font)
            except Exception:
                lw = len(line) * 12
            draw.text((x - int(lw) // 2, y), line, fill=color, font=font)
        else:
            draw.text((x, y), line, fill=color, font=font)
        y += lh
    return y


# ── Three palettes; one is picked at random per generation ───
# Keys are stored as (R, G, B) tuples; the picker copies them into both
# _BRAND_RGB (Pillow video) and _BRAND (PPTX RGBColor) at the start of each job.
_PALETTES = {
    "warm_clinical": {
        "name":       "Warm Clinical",
        "bg":         (0x0F, 0x2A, 0x3F),  # deep navy
        "bg_dark":    (0x0A, 0x1F, 0x2F),  # deeper navy
        "card":       (0xF4, 0xF1, 0xEC),  # warm cream
        "card_text":  (0x1F, 0x29, 0x37),  # charcoal
        "card_alt":   (0xE5, 0xE0, 0xD7),
        "accent":     (0xC9, 0x95, 0x6B),  # muted copper
        "accent2":    (0xE1, 0xC6, 0x99),  # soft gold
        "teal":       (0x4A, 0x6F, 0xA5),  # slate blue (secondary)
        "white":      (0xFF, 0xFF, 0xFF),
        "muted":      (0xB8, 0xC4, 0xD0),  # pale steel
        "subtle":     (0x6B, 0x72, 0x80),
        "header_alt": (0xE5, 0xE0, 0xD7),
    },
    "cool_corporate": {
        "name":       "Cool Corporate",
        "bg":         (0x1F, 0x29, 0x37),  # charcoal slate
        "bg_dark":    (0x11, 0x18, 0x27),
        "card":       (0xF2, 0xEF, 0xE9),  # cool cream
        "card_text":  (0x1F, 0x29, 0x37),
        "card_alt":   (0xE3, 0xDF, 0xD7),
        "accent":     (0x5C, 0x7A, 0xEA),  # muted indigo
        "accent2":    (0xE0, 0xB8, 0x72),  # warm amber
        "teal":       (0x4F, 0x86, 0xC6),  # steel blue
        "white":      (0xFF, 0xFF, 0xFF),
        "muted":      (0xC7, 0xD2, 0xE0),
        "subtle":     (0x6B, 0x72, 0x80),
        "header_alt": (0xE3, 0xDF, 0xD7),
    },
    "teal_rose": {
        "name":       "Deep Teal & Dusty Rose",
        "bg":         (0x0E, 0x3B, 0x43),  # deep teal
        "bg_dark":    (0x08, 0x2B, 0x30),
        "card":       (0xF7, 0xF4, 0xED),  # paper
        "card_text":  (0x1F, 0x29, 0x37),
        "card_alt":   (0xE6, 0xE0, 0xD3),
        "accent":     (0xC2, 0x82, 0x85),  # dusty rose
        "accent2":    (0xD4, 0xB8, 0x96),  # champagne
        "teal":       (0x2A, 0x6F, 0x77),  # mid teal
        "white":      (0xFF, 0xFF, 0xFF),
        "muted":      (0xC0, 0xD6, 0xD8),
        "subtle":     (0x6B, 0x72, 0x80),
        "header_alt": (0xE6, 0xE0, 0xD3),
    },
}

# Default palette (overridden per-job by _apply_random_palette)
_BRAND_RGB = dict(_PALETTES["warm_clinical"])
_BADGE_RGB = {
    "green": (0x10, 0xB9, 0x81), "amber": (0xF5, 0x9E, 0x0B),
    "red":   (0xEF, 0x44, 0x44), "teal":  (0x0E, 0x8C, 0x8B),
    "coral": (0xE7, 0x6F, 0x51), "gold":  (0xE9, 0xC4, 0x6A),
}


def _load_fonts():
    """Returns dict of font sizes -> ImageFont."""
    from PIL import ImageFont
    import os as _os
    sizes = {"xxl": 96, "xl": 64, "lg": 44, "md": 28, "sm": 20, "xs": 16, "xxs": 13}
    families = {"reg": "segoeui.ttf", "bold": "segoeuib.ttf"}
    fonts = {}
    for kind, fam in families.items():
        path = f"C:/Windows/Fonts/{fam}"
        if not _os.path.exists(path):
            path = "C:/Windows/Fonts/arial.ttf"
        for tag, sz in sizes.items():
            try:
                fonts[f"{kind}_{tag}"] = ImageFont.truetype(path, sz)
            except Exception:
                fonts[f"{kind}_{tag}"] = ImageFont.load_default()
    return fonts


def _wrap_lines(draw, text, max_w, font):
    if not text:
        return []
    words = str(text).split()
    lines, cur = [], ""
    for w in words:
        test = (cur + " " + w).strip()
        try:    tw = draw.textlength(test, font=font)
        except: tw = len(test) * 12
        if tw > max_w and cur:
            lines.append(cur); cur = w
        else:
            cur = test
    if cur: lines.append(cur)
    return lines


def _draw_text(draw, text, box, font, color, *, align="left", anchor="top",
               line_gap=6, max_lines=None):
    """Draw word-wrapped text inside box=(x,y,w,h). Returns y after last line."""
    x, y, w, h = box
    lines = _wrap_lines(draw, text, w, font)
    if max_lines:
        lines = lines[:max_lines]
    try:
        lh = draw.textbbox((0, 0), "Ag", font=font)[3] + line_gap
    except Exception:
        lh = 28
    total_h = lh * len(lines)
    if anchor == "middle":
        cy = y + (h - total_h) // 2
    elif anchor == "bottom":
        cy = y + h - total_h
    else:
        cy = y
    for line in lines:
        try:
            lw = draw.textlength(line, font=font)
        except Exception:
            lw = len(line) * 12
        if align == "center":
            lx = x + (w - int(lw)) // 2
        elif align == "right":
            lx = x + w - int(lw)
        else:
            lx = x
        draw.text((lx, cy), line, fill=color, font=font)
        cy += lh
    return cy


def _rrect(draw, box, radius, fill):
    """Rounded rectangle. box=(x0,y0,x1,y1)"""
    try:
        draw.rounded_rectangle(box, radius=radius, fill=fill)
    except Exception:
        draw.rectangle(box, fill=fill)


def _draw_chrome(draw, F, eyebrow, slide_num, total, footer, W, H):
    # Top accent bar
    draw.rectangle([0, 0, W, 12], fill=_BRAND_RGB["accent"])
    if eyebrow:
        _draw_text(draw, eyebrow, (60, 36, W - 120, 32),
                   F["bold_xs"], _BRAND_RGB["accent2"])
    # Footer band
    if footer:
        _draw_text(draw, footer[:80], (60, H - 38, int(W * 0.7), 28),
                   F["bold_xxs"], _BRAND_RGB["muted"])
    if slide_num and total:
        _draw_text(draw, f"{slide_num} / {total}",
                   (W - 220, H - 38, 160, 28),
                   F["bold_xxs"], _BRAND_RGB["muted"], align="right")


def _render_title_img(draw, F, data, deck_meta, length_minutes, W, H):
    # Coral top band
    draw.rectangle([0, 0, W, 56], fill=_BRAND_RGB["accent"])
    _draw_text(draw, data.get("eyebrow", "CLINICAL EDUCATION"),
               (80, 130, W - 160, 36), F["bold_xs"], _BRAND_RGB["muted"])
    title = data.get("title") or deck_meta.get("title", "")
    _draw_text(draw, title, (80, 200, W - 160, 320),
               F["bold_xxl"], _BRAND_RGB["white"], line_gap=10)
    subtitle = data.get("subtitle") or deck_meta.get(
        "subtitle", f"{length_minutes}-Minute Clinical Lecture")
    _draw_text(draw, subtitle, (80, 540, W - 160, 80),
               F["reg_md"], _BRAND_RGB["muted"])
    stats = (data.get("stats") or [])[:3]
    if stats:
        n = len(stats)
        gap = 24
        card_w = (W - 160 - gap * (n - 1)) // n
        x = 80; y = 720; h = 220
        for s in stats:
            _rrect(draw, (x, y, x + card_w, y + h), 18, _BRAND_RGB["bg_dark"])
            _draw_text(draw, s.get("value", ""),
                       (x + 24, y + 30, card_w - 48, 80),
                       F["bold_xl"], _BRAND_RGB["accent2"], align="center")
            _draw_text(draw, s.get("label", ""),
                       (x + 24, y + 130, card_w - 48, 80),
                       F["reg_xs"], _BRAND_RGB["muted"], align="center")
            x += card_w + gap


def _render_stat_cards_img(draw, F, data, W, H):
    _draw_text(draw, data.get("title", ""), (80, 130, W - 160, 80),
               F["bold_lg"], _BRAND_RGB["white"])
    cards = (data.get("cards") or [])[:3]
    if not cards: return
    n = len(cards); gap = 28
    card_w = (W - 160 - gap * (n - 1)) // n
    y = 280; h = 580
    x = 80
    for c in cards:
        _rrect(draw, (x, y, x + card_w, y + h), 24, _BRAND_RGB["card"])
        # Top accent strip
        _rrect(draw, (x, y, x + card_w, y + 28), 0, _BRAND_RGB["accent"])
        _draw_text(draw, c.get("value", ""),
                   (x + 24, y + 90, card_w - 48, 200),
                   F["bold_xxl"], _BRAND_RGB["accent"], align="center")
        _draw_text(draw, c.get("label", ""),
                   (x + 32, y + 320, card_w - 64, h - 360),
                   F["reg_md"], _BRAND_RGB["card_text"], align="center")
        x += card_w + gap


def _render_type_cards_img(draw, F, data, W, H):
    _draw_text(draw, data.get("title", ""), (80, 130, W - 160, 80),
               F["bold_lg"], _BRAND_RGB["white"])
    cards = (data.get("cards") or [])[:4]
    if not cards: return
    n = len(cards); gap = 22
    card_w = (W - 160 - gap * (n - 1)) // n
    y = 240; h = 700
    x = 80
    for c in cards:
        _rrect(draw, (x, y, x + card_w, y + h), 22, _BRAND_RGB["card"])
        # Label strip
        _rrect(draw, (x, y, x + card_w, y + 80), 0, _BRAND_RGB["teal"])
        _draw_text(draw, c.get("label", ""),
                   (x, y + 22, card_w, 40), F["bold_md"],
                   _BRAND_RGB["white"], align="center")
        _draw_text(draw, c.get("heading", ""),
                   (x + 28, y + 110, card_w - 56, 110),
                   F["bold_md"], _BRAND_RGB["card_text"])
        _draw_text(draw, c.get("body", ""),
                   (x + 28, y + 230, card_w - 56, h - 340),
                   F["reg_xs"], _BRAND_RGB["card_text"], line_gap=4)
        badge = c.get("badge")
        if badge:
            bcol = _BADGE_RGB.get((c.get("badge_color") or "teal").lower(),
                                  _BRAND_RGB["teal"])
            bx0, by0 = x + 28, y + h - 70
            bx1, by1 = x + card_w - 28, y + h - 24
            _rrect(draw, (bx0, by0, bx1, by1), 14, bcol)
            _draw_text(draw, badge, (bx0, by0 + 8, bx1 - bx0, 32),
                       F["bold_xxs"], _BRAND_RGB["white"], align="center")
        x += card_w + gap


def _render_numbered_grid_img(draw, F, data, W, H):
    _draw_text(draw, data.get("title", ""), (80, 130, W - 160, 80),
               F["bold_lg"], _BRAND_RGB["white"])
    items = data.get("items") or []
    if not items: return
    cols = 3 if len(items) >= 5 else 2
    rows = (len(items) + cols - 1) // cols
    gap = 22
    card_w = (W - 160 - gap * (cols - 1)) // cols
    card_h = (760 - gap * (rows - 1)) // rows
    x0 = 80; y0 = 240
    for idx, it in enumerate(items[:cols * rows]):
        r, c = divmod(idx, cols)
        x = x0 + (card_w + gap) * c
        y = y0 + (card_h + gap) * r
        _rrect(draw, (x, y, x + card_w, y + card_h), 22, _BRAND_RGB["card"])
        # Number circle
        cx = x + 56; cy = y + 56; rad = 38
        draw.ellipse([cx - rad, cy - rad, cx + rad, cy + rad],
                     fill=_BRAND_RGB["accent"])
        _draw_text(draw, str(it.get("n", idx + 1)),
                   (cx - rad, cy - 22, rad * 2, 44),
                   F["bold_md"], _BRAND_RGB["white"], align="center")
        _draw_text(draw, it.get("heading", ""),
                   (x + 120, y + 26, card_w - 140, 70),
                   F["bold_md"], _BRAND_RGB["card_text"])
        _draw_text(draw, it.get("body", ""),
                   (x + 28, y + 130, card_w - 56, card_h - 150),
                   F["reg_xs"], _BRAND_RGB["subtle"], line_gap=4)


def _render_chart_bar_img(draw, F, data, W, H):
    _draw_text(draw, data.get("title", ""), (80, 130, W - 160, 80),
               F["bold_lg"], _BRAND_RGB["white"])
    cats = data.get("categories") or []
    vals = data.get("values") or []
    unit = data.get("unit", "")
    if not cats or not vals: return
    n = min(len(cats), len(vals))
    cats = cats[:n]; vals = [float(v) for v in vals[:n]]
    vmax = max(vals) if vals else 1
    vmax = max(vmax, 1)
    # Plot area
    px, py, pw, ph = 320, 260, W - 380, 660
    # Frame line
    draw.line([(px, py), (px, py + ph)], fill=_BRAND_RGB["muted"], width=2)
    draw.line([(px, py + ph), (px + pw, py + ph)], fill=_BRAND_RGB["muted"], width=2)
    # Bars
    row_h = ph // n
    bar_h = int(row_h * 0.55)
    bar_pad = (row_h - bar_h) // 2
    for i, (cat, v) in enumerate(zip(cats, vals)):
        y = py + i * row_h + bar_pad
        bw = int((v / vmax) * (pw - 80))
        _rrect(draw, (px + 2, y, px + 2 + bw, y + bar_h), 8, _BRAND_RGB["accent"])
        # Category label (left of axis)
        _draw_text(draw, str(cat), (60, y + bar_h // 2 - 16, 240, 40),
                   F["reg_xs"], _BRAND_RGB["white"], align="right")
        # Value at bar end
        _draw_text(draw, f"{v:g}{unit}",
                   (px + 2 + bw + 14, y + bar_h // 2 - 16, 200, 40),
                   F["bold_xs"], _BRAND_RGB["accent2"])


def _render_comparison_table_img(draw, F, data, W, H):
    _draw_text(draw, data.get("title", ""), (80, 130, W - 160, 80),
               F["bold_lg"], _BRAND_RGB["white"])
    headers = data.get("headers") or []
    rows = data.get("rows") or []
    if not headers or not rows: return
    ncols = len(headers)
    table_x = 80; table_y = 260
    table_w = W - 160; table_h = 680
    nrows = len(rows) + 1
    col_w = table_w // ncols
    row_h = table_h // nrows
    # Header row
    _rrect(draw, (table_x, table_y, table_x + table_w, table_y + row_h),
           14, _BRAND_RGB["accent"])
    for c, h in enumerate(headers):
        cx = table_x + c * col_w
        _draw_text(draw, str(h),
                   (cx + 16, table_y + 18, col_w - 32, row_h - 36),
                   F["bold_sm"], _BRAND_RGB["white"], align="center")
    # Body rows
    for r, row in enumerate(rows, start=1):
        y = table_y + r * row_h
        bg = _BRAND_RGB["card"] if r % 2 == 1 else _BRAND_RGB["card_alt"]
        rad = 14 if r == nrows - 1 else 0
        # Use plain rectangle for middle rows; rounded only for last row's bottom
        if r == nrows - 1:
            _rrect(draw, (table_x, y, table_x + table_w, y + row_h), 14, bg)
            # Repaint the top so it's flat (only round bottom corners)
            draw.rectangle([table_x, y, table_x + table_w, y + 14], fill=bg)
        else:
            draw.rectangle([table_x, y, table_x + table_w, y + row_h], fill=bg)
        for c in range(ncols):
            cx = table_x + c * col_w
            cell = str(row[c]) if c < len(row) else ""
            font = F["bold_sm"] if c == 0 else F["reg_sm"]
            _draw_text(draw, cell, (cx + 16, y + 14, col_w - 32, row_h - 28),
                       font, _BRAND_RGB["card_text"])


def _render_bullets_img(draw, F, data, W, H):
    _draw_text(draw, data.get("title", ""), (80, 130, W - 160, 80),
               F["bold_lg"], _BRAND_RGB["white"])
    bullets = data.get("bullets") or []
    if not bullets: return
    y = 260
    for b in bullets[:7]:
        # Coral square bullet
        draw.rectangle([80, y + 14, 104, y + 38], fill=_BRAND_RGB["accent"])
        y_after = _draw_text(draw, b, (130, y, W - 220, 100),
                             F["reg_md"], _BRAND_RGB["white"])
        y = max(y_after, y + 64)
        if y > H - 120:
            break


def _render_summary_img(draw, F, data, W, H):
    # Inset card
    _rrect(draw, (60, 110, W - 60, H - 70), 28, _BRAND_RGB["bg_dark"])
    _draw_text(draw, data.get("eyebrow", "KEY TAKEAWAYS"),
               (110, 150, W - 220, 32), F["bold_xs"], _BRAND_RGB["accent2"])
    _draw_text(draw, data.get("title", "Clinical Pearls"),
               (110, 200, W - 220, 100), F["bold_lg"], _BRAND_RGB["white"])
    bullets = data.get("bullets") or []
    y = 340
    for i, b in enumerate(bullets[:6], start=1):
        # Coral chip
        cx = 140; cy = y + 26
        draw.ellipse([cx - 26, cy - 26, cx + 26, cy + 26],
                     fill=_BRAND_RGB["accent"])
        _draw_text(draw, str(i), (cx - 26, cy - 22, 52, 44),
                   F["bold_md"], _BRAND_RGB["white"], align="center")
        _draw_text(draw, b, (200, y, W - 280, 80),
                   F["reg_md"], _BRAND_RGB["white"])
        y += 90
        if y > H - 110: break


def _render_references_img(draw, F, data, W, H):
    _draw_text(draw, data.get("title", "References"),
               (80, 130, W - 160, 80), F["bold_lg"], _BRAND_RGB["white"])
    items = data.get("items") or data.get("bullets") or []
    if not items: return
    y = 260
    for i, item in enumerate(items[:10], start=1):
        _draw_text(draw, f"{i}.", (90, y, 60, 32), F["bold_xs"],
                   _BRAND_RGB["accent2"])
        y_after = _draw_text(draw, str(item), (150, y, W - 240, 90),
                             F["reg_xs"], _BRAND_RGB["muted"], line_gap=2)
        y = max(y_after, y + 50)
        if y > H - 110: break


_IMG_RENDERERS = {
    "stat_cards":       _render_stat_cards_img,
    "type_cards":       _render_type_cards_img,
    "numbered_grid":    _render_numbered_grid_img,
    "chart_bar":        _render_chart_bar_img,
    "comparison_table": _render_comparison_table_img,
    "bullets":          _render_bullets_img,
    "summary":          _render_summary_img,
    "references":       _render_references_img,
}


def _render_slide_image(slide_data: dict, slide_num: int, total_slides: int,
                         deck_title: str, length_minutes: int,
                         deck_meta: dict = None) -> bytes:
    """Render one slide as a 1920x1080 PNG using Pillow with the brand layouts."""
    import io as _io2
    from PIL import Image, ImageDraw

    W, H = 1920, 1080
    img = Image.new("RGB", (W, H), color=_BRAND_RGB["bg"])
    draw = ImageDraw.Draw(img)
    F = _load_fonts()

    deck_meta = deck_meta or {"title": deck_title}
    stype = (slide_data.get("type") or slide_data.get("slide_type")
             or ("title" if slide_num == 1 else "bullets")).lower()
    footer = (deck_meta.get("footer") or deck_title or "").upper()[:80]

    try:
        if stype == "title" or slide_num == 1:
            _render_title_img(draw, F, slide_data, deck_meta, length_minutes, W, H)
        else:
            renderer = _IMG_RENDERERS.get(stype, _render_bullets_img)
            renderer(draw, F, slide_data, W, H)
            _draw_chrome(draw, F, slide_data.get("eyebrow", ""),
                         slide_num, total_slides, footer, W, H)
    except Exception as e:
        # Fallback: title-only slide so the video isn't blank
        import traceback; traceback.print_exc()
        print(f"  [video] slide {slide_num} render error ({stype}): {e}")
        _draw_text(draw, slide_data.get("title", ""),
                   (80, 80, W - 160, 200), F["bold_lg"], _BRAND_RGB["white"])

    buf = _io2.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ── Narrated video generation ─────────────────────────────

def run_generate_video(audio_id: str, answer: str, question: str,
                       length_minutes: int, voice: str = ""):
    """Generate a narrated MP4 video: Pillow slide images + OpenAI TTS audio."""
    import os as _os3, tempfile as _tmp3
    try:
        from endo_ai import generate_slides_content

        _apply_random_palette()
        with audio_jobs_lock:
            audio_jobs[audio_id]["status"] = "generating_content"

        print(f"  [video] Generating {length_minutes}-min slide content...")
        deck = generate_slides_content(answer, question, length_minutes)
        slides     = deck.get("slides", [])
        deck_title = deck.get("title", question)
        total      = len(slides)

        if not slides:
            # Both attempts in generate_slides_content failed. The detailed
            # parse-failure reason was already printed there (look for
            # "[slides] first/retry attempt produced 0 slides"). Surface a
            # user-facing message that points to the most likely cause.
            raise ValueError(
                "Slide generator returned 0 slides after retry. "
                "Most common cause: the JSON exceeded max_tokens and got truncated — "
                "try a shorter length (e.g. 5 min instead of 15 min) or simplify the topic. "
                "Check server log for the parse-failure reason."
            )

        with audio_jobs_lock:
            audio_jobs[audio_id]["status"] = "building_slides"
            audio_jobs[audio_id]["slides_total"] = total
            audio_jobs[audio_id]["slides_done"]  = 0

        tmpdir = _tmp3.mkdtemp(prefix="endo_vid_")
        img_paths   = [None] * total
        audio_paths = [None] * total

        # ── 1. Render slide images sequentially (Pillow is fast) ──
        for i, slide_data in enumerate(slides):
            slide_num = i + 1
            with audio_jobs_lock:
                if audio_jobs[audio_id].get("cancelled"):
                    print("  [video] Cancelled"); return
            img_bytes = _render_slide_image(slide_data, slide_num, total,
                                            deck_title, length_minutes, deck)
            ip = _os3.path.join(tmpdir, f"slide_{slide_num:03d}.png")
            with open(ip, 'wb') as f: f.write(img_bytes)
            img_paths[i] = ip

        # ── 2. Generate TTS in parallel (OpenAI handles concurrent reqs fine) ──
        with audio_jobs_lock:
            audio_jobs[audio_id]["status"] = "generating_audio"

        from concurrent.futures import ThreadPoolExecutor, as_completed
        import threading as _th

        done_lock = _th.Lock()
        done_count = [0]

        # Resolved ONCE, outside the pool: every slide of one video must be
        # spoken by the same voice, and resolve_voice() prints a warning on an
        # unknown name that would otherwise repeat per slide.
        tts_voice = narration.resolve_voice(voice)
        tts_model = narration.resolve_model()
        billed_chars = [0]

        def _tts_one(idx, slide_data):
            # Per-slide, NOT synthesize_lecture: this clip is paired with this
            # slide's image below (ffmpeg -shortest), so its audio boundary has
            # to be the slide's own. synthesize_segment gives the same
            # pronunciation dictionary, voice and model as the audio export
            # without merging the slides into one file.
            slide_num = idx + 1
            notes = (slide_data.get("speaker_notes") or "").strip()
            if not notes:
                return idx, None, 0
            seg = narration.synthesize_segment(
                notes, voice=tts_voice, model=tts_model,
                label=f"video slide {slide_num}",
                allow_gtts=GTTS_AVAILABLE)
            if not seg["audio"]:
                return idx, None, 0
            ap = _os3.path.join(tmpdir, f"audio_{slide_num:03d}.mp3")
            with open(ap, 'wb') as f:
                f.write(seg["audio"])
            return idx, ap, seg["characters"]

        # OpenAI TTS comfortably handles 6 parallel reqs on default tier.
        # gTTS uses Google's free endpoint -- keep concurrency lower if it kicks in.
        max_workers = 6 if OPENAI_TTS_AVAILABLE else 3
        print(f"  [video] narration: voice={tts_voice} model={tts_model}")
        with ThreadPoolExecutor(max_workers=max_workers) as pool:
            futures = [pool.submit(_tts_one, i, sd) for i, sd in enumerate(slides)]
            for fut in as_completed(futures):
                with audio_jobs_lock:
                    if audio_jobs[audio_id].get("cancelled"):
                        print("  [video] Cancelled mid-TTS"); return
                try:
                    idx, ap, chars = fut.result()
                    audio_paths[idx] = ap
                    with done_lock:
                        billed_chars[0] += chars
                    if ap:
                        print(f"    [video] slide {idx+1}/{total} TTS OK")
                except Exception as e:
                    print(f"    [video] TTS task failed: {e}")
                with done_lock:
                    done_count[0] += 1
                    with audio_jobs_lock:
                        audio_jobs[audio_id]["slides_done"] = done_count[0]

        # ONE cost row for the whole export, carrying the job id — the job made
        # one API call per slide but it is one billable video.
        narration.log_narration_cost(
            "run_generate_video", tts_model, billed_chars[0],
            request_id=audio_id, voice=tts_voice)

        # ── Assemble video with moviepy ───────────────────
        with audio_jobs_lock:
            audio_jobs[audio_id]["status"] = "injecting_audio"

        print("  [video] Assembling video clips with native ffmpeg...")
        out_path = _os3.path.join(tmpdir, f"narrated_{audio_id[:8]}.mp4")

        if _native_ffmpeg:
            # ── Direct ffmpeg path: bypass moviepy entirely ──
            # For each (image, audio) pair, encode a per-slide MP4 (~1-3s/slide
            # for stillimages). Then concat losslessly with -c copy (instant).
            import subprocess as _sp, time as _t
            t0 = _t.time()
            slide_clips = []
            for i, (ip, ap) in enumerate(zip(img_paths, audio_paths)):
                clip_out = _os3.path.join(tmpdir, f"clip_{i+1:03d}.mp4")
                if ap and _os3.path.exists(ap):
                    cmd = [
                        _native_ffmpeg, "-y", "-loglevel", "error",
                        "-loop", "1", "-i", ip,
                        "-i", ap,
                        "-c:v", "libx264", "-tune", "stillimage",
                        "-preset", "ultrafast", "-crf", "23",
                        "-pix_fmt", "yuv420p",
                        "-c:a", "aac", "-b:a", "128k",
                        "-shortest",
                        "-r", "12",
                        clip_out,
                    ]
                else:
                    # Silent 4-second slide
                    cmd = [
                        _native_ffmpeg, "-y", "-loglevel", "error",
                        "-loop", "1", "-i", ip,
                        "-f", "lavfi", "-i", "anullsrc=channel_layout=stereo:sample_rate=44100",
                        "-t", "4",
                        "-c:v", "libx264", "-tune", "stillimage",
                        "-preset", "ultrafast", "-crf", "23",
                        "-pix_fmt", "yuv420p",
                        "-c:a", "aac", "-b:a", "128k",
                        "-r", "12",
                        clip_out,
                    ]
                _sp.run(cmd, check=True, capture_output=True)
                slide_clips.append(clip_out)

            # Concat all slide clips losslessly
            list_path = _os3.path.join(tmpdir, "concat_list.txt")
            with open(list_path, "w", encoding="utf-8") as f:
                for c in slide_clips:
                    f.write(f"file '{c.replace(chr(92), '/')}'\n")
            concat_cmd = [
                _native_ffmpeg, "-y", "-loglevel", "error",
                "-f", "concat", "-safe", "0",
                "-i", list_path,
                "-c", "copy",
                out_path,
            ]
            _sp.run(concat_cmd, check=True, capture_output=True)
            print(f"  [video] OK MP4 saved in {_t.time()-t0:.1f}s: {out_path}")

        else:
            # ── Moviepy fallback (slow on Windows; only if ffmpeg missing) ──
            print("  [video] Native ffmpeg not found, falling back to moviepy...")
            clips = []
            for ip, ap in zip(img_paths, audio_paths):
                if ap and _os3.path.exists(ap):
                    ac   = AudioFileClip(ap)
                    dur  = ac.duration
                    clip = ImageClip(ip, duration=dur).with_audio(ac)
                else:
                    clip = ImageClip(ip, duration=4)
                clips.append(clip)
            final = concatenate_videoclips(clips, method="compose")
            final.write_videofile(out_path, fps=12, codec='libx264',
                                  audio_codec='aac', logger=None,
                                  ffmpeg_params=['-preset', 'ultrafast',
                                                 '-tune', 'stillimage',
                                                 '-crf', '23', '-pix_fmt', 'yuv420p'],
                                  threads=0)
            for c in clips: c.close()
            final.close()
            print(f"  [video] OK MP4 saved: {out_path}")

        with audio_jobs_lock:
            audio_jobs[audio_id]["status"]    = "complete"
            audio_jobs[audio_id]["file_path"] = out_path
            audio_jobs[audio_id]["file_ext"]  = "mp4"
            q   = audio_jobs[audio_id].get("question", "")
            dur = audio_jobs[audio_id].get("length_minutes", 10)
        _persist_media(out_path, audio_id, "mp4", q, "video", "video", dur)

    except Exception as e:
        import traceback; traceback.print_exc()
        with audio_jobs_lock:
            audio_jobs[audio_id]["status"] = "error"
            audio_jobs[audio_id]["error"]  = str(e)


# ── PPTX brand palette (mirrors _BRAND_RGB; overwritten per job) ─
def _palette_to_pptx(p):
    """Convert a palette tuple-dict into RGBColor objects for PPTX."""
    keys = ("bg","bg_dark","card","card_text","accent","accent2",
            "teal","white","muted","subtle")
    return {k: RGBColor(*p[k]) for k in keys}

_BRAND = _palette_to_pptx(_PALETTES["warm_clinical"])


def _apply_random_palette():
    """Pick one of the three palettes at random and update both
    _BRAND_RGB (Pillow video) and _BRAND (PPTX) in place. Renderers
    read these module-level dicts so they pick up the new colors."""
    import random as _random_pal
    name = _random_pal.choice(list(_PALETTES.keys()))
    p    = _PALETTES[name]
    _BRAND_RGB.clear(); _BRAND_RGB.update(p)
    new_pptx = _palette_to_pptx(p)
    _BRAND.clear(); _BRAND.update(new_pptx)
    print(f"  [palette] {p['name']}")
    return name

_BADGE_COLORS = {
    "green":   RGBColor(0x10, 0xB9, 0x81),
    "amber":   RGBColor(0xF5, 0x9E, 0x0B),
    "red":     RGBColor(0xEF, 0x44, 0x44),
    "teal":    RGBColor(0x0E, 0x8C, 0x8B),
    "coral":   RGBColor(0xE7, 0x6F, 0x51),
    "gold":    RGBColor(0xE9, 0xC4, 0x6A),
}


def _set_solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _add_rect(slide, x, y, w, h, rgb, rounded=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    if rounded:
        try:
            shp.adjustments[0] = 0.10
        except Exception:
            pass
    _set_solid(shp, rgb)
    return shp


def _add_text(slide, x, y, w, h, text, *, size=14, bold=False,
              color=None, align="left", anchor="top"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM
    }.get(anchor, MSO_ANCHOR.TOP)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = str(text or "")
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    return tb


def _slide_bg(slide, rgb):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb


def _add_chrome(slide, eyebrow, slide_num, total, footer):
    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.08))
    _set_solid(bar, _BRAND["accent"])
    # Eyebrow
    if eyebrow:
        _add_text(slide, Inches(0.5), Inches(0.25), Inches(10), Inches(0.4),
                  eyebrow, size=11, bold=True, color=_BRAND["accent2"])
    # Footer
    if footer:
        _add_text(slide, Inches(0.5), Inches(7.05), Inches(10), Inches(0.4),
                  footer, size=9, bold=True, color=_BRAND["muted"])
    # Page number
    if slide_num and total:
        _add_text(slide, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.4),
                  f"{slide_num} / {total}", size=10, bold=True,
                  color=_BRAND["muted"], align="right")


def _render_title(slide, data, deck, length_minutes):
    _slide_bg(slide, _BRAND["bg"])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.5))
    _set_solid(bar, _BRAND["accent"])
    _add_text(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.45),
              data.get("eyebrow", "CLINICAL EDUCATION"),
              size=12, bold=True, color=_BRAND["muted"])
    _add_text(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(2.2),
              data.get("title") or deck.get("title", ""),
              size=54, bold=True, color=_BRAND["white"])
    _add_text(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(1.2),
              data.get("subtitle") or deck.get("subtitle", f"{length_minutes}-Minute Clinical Lecture"),
              size=22, color=_BRAND["muted"])
    # Stat strip across bottom
    stats = data.get("stats") or []
    if stats:
        x = Inches(0.6); top = Inches(5.2); card_w = Inches(4.0); card_h = Inches(1.4)
        gap = Inches(0.15)
        for s in stats[:3]:
            _add_rect(slide, x, top, card_w, card_h, _BRAND["bg_dark"])
            _add_text(slide, x + Inches(0.2), top + Inches(0.15),
                      card_w - Inches(0.4), Inches(0.6),
                      s.get("value", ""), size=28, bold=True, color=_BRAND["accent2"])
            _add_text(slide, x + Inches(0.2), top + Inches(0.75),
                      card_w - Inches(0.4), Inches(0.55),
                      s.get("label", ""), size=11, color=_BRAND["muted"])
            x += card_w + gap


def _render_stat_cards(slide, data):
    _slide_bg(slide, _BRAND["bg"])
    _add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.7),
              data.get("title", ""), size=30, bold=True, color=_BRAND["white"])
    cards = (data.get("cards") or [])[:3]
    if not cards: return
    n = len(cards)
    total_w = Inches(12.3); gap = Inches(0.25)
    card_w = (total_w - gap * (n - 1)) / n
    x = Inches(0.5); y = Inches(2.5); h = Inches(3.8)
    for c in cards:
        _add_rect(slide, x, y, card_w, h, _BRAND["card"])
        # Accent strip on top of card
        _add_rect(slide, x, y, card_w, Inches(0.18), _BRAND["accent"], rounded=False)
        _add_text(slide, x + Inches(0.3), y + Inches(0.5),
                  card_w - Inches(0.6), Inches(1.4),
                  c.get("value", ""), size=54, bold=True, color=_BRAND["accent"], align="center")
        _add_text(slide, x + Inches(0.3), y + Inches(2.0),
                  card_w - Inches(0.6), h - Inches(2.2),
                  c.get("label", ""), size=14, color=_BRAND["card_text"], align="center")
        x += card_w + gap


def _render_type_cards(slide, data):
    _slide_bg(slide, _BRAND["bg"])
    _add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.7),
              data.get("title", ""), size=28, bold=True, color=_BRAND["white"])
    cards = (data.get("cards") or [])[:4]
    if not cards: return
    n = len(cards)
    total_w = Inches(12.3); gap = Inches(0.2)
    card_w = (total_w - gap * (n - 1)) / n
    x = Inches(0.5); y = Inches(2.2); h = Inches(4.4)
    for c in cards:
        _add_rect(slide, x, y, card_w, h, _BRAND["card"])
        # Label strip
        _add_rect(slide, x, y, card_w, Inches(0.6), _BRAND["teal"], rounded=False)
        _add_text(slide, x, y + Inches(0.05), card_w, Inches(0.5),
                  c.get("label", ""), size=14, bold=True,
                  color=_BRAND["white"], align="center")
        _add_text(slide, x + Inches(0.25), y + Inches(0.85),
                  card_w - Inches(0.5), Inches(0.7),
                  c.get("heading", ""), size=15, bold=True, color=_BRAND["card_text"])
        _add_text(slide, x + Inches(0.25), y + Inches(1.6),
                  card_w - Inches(0.5), h - Inches(2.5),
                  c.get("body", ""), size=11, color=_BRAND["card_text"])
        # Badge
        badge = c.get("badge")
        if badge:
            bcol = _BADGE_COLORS.get((c.get("badge_color") or "teal").lower(),
                                     _BRAND["teal"])
            bw = card_w - Inches(0.6); bh = Inches(0.4)
            _add_rect(slide, x + Inches(0.3), y + h - Inches(0.65),
                      bw, bh, bcol)
            _add_text(slide, x + Inches(0.3), y + h - Inches(0.62),
                      bw, bh, badge, size=11, bold=True,
                      color=_BRAND["white"], align="center")
        x += card_w + gap


def _render_numbered_grid(slide, data):
    _slide_bg(slide, _BRAND["bg"])
    _add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.7),
              data.get("title", ""), size=28, bold=True, color=_BRAND["white"])
    items = data.get("items") or []
    if not items: return
    cols = 3 if len(items) >= 5 else 2
    rows = (len(items) + cols - 1) // cols
    total_w = Inches(12.3); total_h = Inches(4.6)
    gap = Inches(0.2)
    card_w = (total_w - gap * (cols - 1)) / cols
    card_h = (total_h - gap * (rows - 1)) / rows
    x0 = Inches(0.5); y0 = Inches(2.2)
    for idx, it in enumerate(items[:cols*rows]):
        r, c = divmod(idx, cols)
        x = x0 + (card_w + gap) * c
        y = y0 + (card_h + gap) * r
        _add_rect(slide, x, y, card_w, card_h, _BRAND["card"])
        # Number circle
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       x + Inches(0.2), y + Inches(0.2),
                                       Inches(0.65), Inches(0.65))
        _set_solid(circ, _BRAND["accent"])
        _add_text(slide, x + Inches(0.2), y + Inches(0.22),
                  Inches(0.65), Inches(0.65),
                  str(it.get("n", idx + 1)), size=18, bold=True,
                  color=_BRAND["white"], align="center")
        _add_text(slide, x + Inches(1.0), y + Inches(0.2),
                  card_w - Inches(1.2), Inches(0.6),
                  it.get("heading", ""), size=14, bold=True, color=_BRAND["card_text"])
        _add_text(slide, x + Inches(0.3), y + Inches(0.95),
                  card_w - Inches(0.6), card_h - Inches(1.1),
                  it.get("body", ""), size=10, color=_BRAND["subtle"])


def _render_chart_bar(slide, data):
    _slide_bg(slide, _BRAND["bg"])
    _add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.7),
              data.get("title", ""), size=28, bold=True, color=_BRAND["white"])
    cats = data.get("categories") or []
    vals = data.get("values") or []
    if not cats or not vals: return
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series(data.get("unit", "Value"), vals)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.6),
        cd)
    chart = chart_shape.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(12); dl.font.bold = True
    try: dl.font.color.rgb = _BRAND["white"]
    except Exception: pass
    # Color the bars coral
    try:
        for series in chart.series:
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = _BRAND["accent"]
    except Exception:
        pass
    # Axis text white
    try:
        for ax in (chart.category_axis, chart.value_axis):
            ax.tick_labels.font.size = Pt(11)
            ax.tick_labels.font.color.rgb = _BRAND["white"]
    except Exception:
        pass


def _render_comparison_table(slide, data):
    _slide_bg(slide, _BRAND["bg"])
    _add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.7),
              data.get("title", ""), size=28, bold=True, color=_BRAND["white"])
    headers = data.get("headers") or []
    rows    = data.get("rows") or []
    if not headers or not rows: return
    nrows = len(rows) + 1
    ncols = len(headers)
    table_shape = slide.shapes.add_table(
        nrows, ncols,
        Inches(0.6), Inches(2.1), Inches(12.1), Inches(4.4))
    table = table_shape.table
    # Header row
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = ""
        cell.fill.solid(); cell.fill.fore_color.rgb = _BRAND["accent"]
        para = cell.text_frame.paragraphs[0]
        run = para.add_run(); run.text = str(h)
        run.font.size = Pt(13); run.font.bold = True
        run.font.color.rgb = _BRAND["white"]
        para.alignment = PP_ALIGN.CENTER
    # Body rows
    for r, row in enumerate(rows, start=1):
        for c in range(ncols):
            cell = table.cell(r, c)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = (_BRAND["card"] if r % 2 == 1
                                         else RGBColor(0xE5, 0xE0, 0xD7))
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = str(row[c]) if c < len(row) else ""
            run.font.size = Pt(12)
            run.font.color.rgb = _BRAND["card_text"]
            run.font.bold = (c == 0)
            para.alignment = PP_ALIGN.LEFT


def _render_bullets(slide, data):
    _slide_bg(slide, _BRAND["bg"])
    _add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.7),
              data.get("title", ""), size=28, bold=True, color=_BRAND["white"])
    bullets = data.get("bullets") or []
    if not bullets: return
    y = Inches(2.2)
    for b in bullets[:7]:
        # Coral bullet square
        sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.7), y + Inches(0.18),
                                     Inches(0.18), Inches(0.18))
        _set_solid(sq, _BRAND["accent"])
        _add_text(slide, Inches(1.05), y, Inches(11.5), Inches(0.7),
                  b, size=16, color=_BRAND["white"])
        y += Inches(0.65)


def _render_summary(slide, data):
    _slide_bg(slide, _BRAND["bg"])
    _add_rect(slide, Inches(0.5), Inches(0.9), Inches(12.3), Inches(5.8),
              _BRAND["bg_dark"])
    _add_text(slide, Inches(0.8), Inches(1.1), Inches(12), Inches(0.5),
              data.get("eyebrow", "KEY TAKEAWAYS"), size=12, bold=True,
              color=_BRAND["accent2"])
    _add_text(slide, Inches(0.8), Inches(1.55), Inches(12), Inches(0.9),
              data.get("title", "Clinical Pearls"),
              size=32, bold=True, color=_BRAND["white"])
    bullets = data.get("bullets") or []
    y = Inches(2.7)
    for i, b in enumerate(bullets[:6], start=1):
        # Numbered chip
        chip = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(0.9), y + Inches(0.05),
                                       Inches(0.45), Inches(0.45))
        _set_solid(chip, _BRAND["accent"])
        _add_text(slide, Inches(0.9), y + Inches(0.05),
                  Inches(0.45), Inches(0.45),
                  str(i), size=14, bold=True, color=_BRAND["white"], align="center")
        _add_text(slide, Inches(1.6), y, Inches(11), Inches(0.6),
                  b, size=15, color=_BRAND["white"])
        y += Inches(0.6)


def _render_references(slide, data):
    _slide_bg(slide, _BRAND["bg"])
    _add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.7),
              data.get("title", "References"), size=28, bold=True, color=_BRAND["white"])
    items = data.get("items") or data.get("bullets") or []
    if not items: return
    y = Inches(2.0)
    for i, item in enumerate(items[:10], start=1):
        _add_text(slide, Inches(0.7), y, Inches(0.5), Inches(0.4),
                  f"{i}.", size=11, bold=True, color=_BRAND["accent2"])
        _add_text(slide, Inches(1.1), y, Inches(11.5), Inches(0.5),
                  item, size=11, color=_BRAND["muted"])
        y += Inches(0.42)


_RENDERERS = {
    "stat_cards":       _render_stat_cards,
    "type_cards":       _render_type_cards,
    "numbered_grid":    _render_numbered_grid,
    "chart_bar":        _render_chart_bar,
    "comparison_table": _render_comparison_table,
    "bullets":          _render_bullets,
    "summary":          _render_summary,
    "references":       _render_references,
}


def run_generate_slides(audio_id: str, answer: str, question: str,
                        length_minutes: int, voice: str = "",
                        papers: list | None = None):
    try:
        from presentations.build_deck import build_deck_from_specs
        from presentations.chart_data import tier_counts_from_papers
        import slide_spec_cache

        with audio_jobs_lock:
            audio_jobs[audio_id]["status"] = "generating_content"

        print(f"  Generating {length_minutes}-min pattern-based slide specs...")
        # §5.1: ONE canonical text object per answer. The web deck already
        # reads through this cache; the pptx path used to call the generator
        # directly, so the two decks were laying out two DIFFERENT LLM
        # generations of the same answer and their content hashes could never
        # have matched.
        deck, spec_hash, from_cache = slide_spec_cache.get_or_build(
            answer, question, length_minutes)
        print(f"  slide spec {spec_hash[:12]} "
              f"({'cached' if from_cache else 'generated'})")

        slides_list = deck.get("slides", []) or []
        if not slides_list:
            raise ValueError(
                "Slide generator returned 0 slides after retry. "
                "Try a shorter length or check server log for parse-failure reason."
            )

        with audio_jobs_lock:
            audio_jobs[audio_id]["status"] = "building_slides"
            audio_jobs[audio_id]["slides_total"] = len(slides_list)
            audio_jobs[audio_id]["slides_done"]  = 0

        print(f"  Building {len(slides_list)}-slide PPTX with design-token patterns...")
        # source_text is what the chart gate verifies plotted values against.
        # Without it every chart is correctly suppressed — §1.5 forbids
        # plotting a number that is not verbatim in the cited source, and the
        # deck's own speaker notes cannot serve as that corpus because the
        # same model wrote both.
        prs, slides_queue = build_deck_from_specs(
            deck, source_text=answer,
            tier_counts=tier_counts_from_papers(papers or []))

        # ── Save base PPTX (no audio yet) ─────────────────────
        tmp_base = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        prs.save(tmp_base.name)
        tmp_base.close()
        print(f"  Base PPTX saved: {len(slides_queue)} slides")

        # ── Generate TTS per slide and inject via ZIP surgery ──
        slide_audios: dict = {}   # {slide_num_1based: mp3_bytes}

        if OPENAI_TTS_AVAILABLE and slides_queue:
            with audio_jobs_lock:
                audio_jobs[audio_id]["status"] = "generating_audio"
                audio_jobs[audio_id]["slides_done"] = 0
            total = len(slides_queue)

            # Resolved once for the whole deck; see run_generate_video.
            tts_voice_p = narration.resolve_voice(voice)
            tts_model_p = narration.resolve_model()
            billed_chars_p = [0]
            print(f"  Recording narration for {total} slides "
                  f"(voice={tts_voice_p} model={tts_model_p}) in parallel...")

            from concurrent.futures import ThreadPoolExecutor, as_completed
            import threading as _th_pptx
            done_lock_p = _th_pptx.Lock(); done_count_p = [0]

            def _tts_pptx(slide_num, notes_text):
                # Per-slide, NOT synthesize_lecture: these bytes are injected
                # into ONE pptx slide part, so the audio boundary must be that
                # slide's. `notes_text`, not a local named `narration` — the
                # module import is what carries the dictionary.
                seg = narration.synthesize_segment(
                    notes_text, voice=tts_voice_p, model=tts_model_p,
                    label=f"pptx slide {slide_num}",
                    allow_gtts=False)
                if not seg["audio"]:
                    return slide_num, None, 0
                return slide_num, seg["audio"], seg["characters"]

            with ThreadPoolExecutor(max_workers=6) as pool:
                futures = [pool.submit(_tts_pptx, sn, n)
                           for _, n, sn in slides_queue]
                for fut in as_completed(futures):
                    with audio_jobs_lock:
                        if audio_jobs[audio_id].get("cancelled"):
                            print("  Job cancelled by user"); return
                    sn, content, chars = fut.result()
                    if content is not None:
                        slide_audios[sn] = content
                        print(f"    slide {sn}/{total} OK")
                    with done_lock_p:
                        billed_chars_p[0] += chars
                        done_count_p[0] += 1
                        with audio_jobs_lock:
                            audio_jobs[audio_id]["slides_done"] = done_count_p[0]

            # ONE cost row for the whole deck, carrying the job id.
            narration.log_narration_cost(
                "run_generate_slides", tts_model_p, billed_chars_p[0],
                request_id=audio_id, voice=tts_voice_p)

        elif GTTS_AVAILABLE and slides_queue:
            with audio_jobs_lock:
                audio_jobs[audio_id]["status"] = "generating_audio"
            import io as _sysio2
            all_notes = " ".join(n for _, n, _ in slides_queue if (n or "").strip())
            if all_notes:
                try:
                    buf = _sysio2.BytesIO()
                    gTTS(text=all_notes[:5000], lang="en", slow=False).write_to_fp(buf)
                    slide_audios[1] = buf.getvalue()
                    print("  gTTS narration recorded for slide 1")
                except Exception as gtts_err:
                    print(f"  gTTS error: {gtts_err}")
        else:
            print("  No TTS backend -- slides will include text notes only")

        # ── Inject audio into PPTX via ZIP manipulation ────────
        final_path = tmp_base.name
        if slide_audios:
            with audio_jobs_lock:
                audio_jobs[audio_id]["status"] = "injecting_audio"
            print(f"  Injecting audio into {len(slide_audios)} slides...")
            try:
                final_path = _inject_audio_into_pptx(tmp_base.name, slide_audios)
                print(f"  OK Narrated PPTX: {final_path}")
                import os as _os; _os.unlink(tmp_base.name)
            except Exception as inj_err:
                import traceback; traceback.print_exc()
                print(f"  Audio injection failed: {inj_err} -- delivering base PPTX")
                final_path = tmp_base.name

        with audio_jobs_lock:
            audio_jobs[audio_id]["status"]    = "complete"
            audio_jobs[audio_id]["file_path"] = final_path
            q   = audio_jobs[audio_id].get("question", "")
            dur = audio_jobs[audio_id].get("length_minutes", 10)
        _persist_media(final_path, audio_id, "pptx", q, "slides", "slides", dur)
        # len(prs.slides), not len(slides_list): the body-budget rule splits
        # overfull slides onto continuation pages, so the spec count under-
        # reports what was actually rendered.
        print(f"  OK Done ({len(prs.slides)} slides, "
              f"{len(slide_audios)} with audio)")

    except Exception as e:
        import traceback; traceback.print_exc()
        print(f"  Slides generation error: {e}")
        with audio_jobs_lock:
            audio_jobs[audio_id]["status"] = "error"
            audio_jobs[audio_id]["error"]  = str(e)


# ── Cancel audio/slides job ───────────────────────────────

@app.route("/cancel_audio/<audio_id>", methods=["POST"])
def cancel_audio_job(audio_id: str):
    with audio_jobs_lock:
        job = audio_jobs.get(audio_id)
        if job and job.get("status") not in ("complete", "error", "cancelled"):
            job["status"]    = "cancelled"
            job["cancelled"] = True
    return jsonify({"ok": True})


# ── Media Library (persistent generated files) ───────────

import json as _json_mod, shutil as _shutil

MEDIA_DIR        = os.path.join(os.path.dirname(__file__), "generated_media")
_MEDIA_INDEX_PATH = os.path.join(os.path.dirname(__file__), "media_index.json")
_media_index_lock = threading.Lock()
os.makedirs(MEDIA_DIR, exist_ok=True)

def _load_media_index() -> list:
    if os.path.exists(_MEDIA_INDEX_PATH):
        try:
            with open(_MEDIA_INDEX_PATH, "r", encoding="utf-8") as f:
                return _json_mod.load(f)
        except Exception:
            pass
    return []

def _append_media_item(item: dict):
    """Thread-safe: prepend item to media index JSON (newest first)."""
    with _media_index_lock:
        items = _load_media_index()
        items.insert(0, item)
        with open(_MEDIA_INDEX_PATH, "w", encoding="utf-8") as f:
            _json_mod.dump(items, f, indent=2)

def _remove_media_item(media_id: str) -> bool:
    with _media_index_lock:
        items = _load_media_index()
        target = next((i for i in items if i.get("id") == media_id), None)
        if not target:
            return False
        fpath = os.path.join(MEDIA_DIR, target.get("filename", ""))
        try:
            if os.path.exists(fpath):
                os.unlink(fpath)
        except Exception:
            pass
        items = [i for i in items if i.get("id") != media_id]
        with open(_MEDIA_INDEX_PATH, "w", encoding="utf-8") as f:
            _json_mod.dump(items, f, indent=2)
        return True

def _persist_media(src_path: str, media_id: str, ext: str,
                   question: str, style: str, media_type: str,
                   length_minutes: int):
    """Copy temp file to MEDIA_DIR and record in index."""
    try:
        filename = f"{media_id}.{ext}"
        dst = os.path.join(MEDIA_DIR, filename)
        _shutil.copy2(src_path, dst)
        _append_media_item({
            "id":             media_id,
            "type":           media_type,   # "audio" | "video" | "slides"
            "style":          style,        # "conversation" | "lecture" | "video" | "slides"
            "ext":            ext,
            "filename":       filename,
            "question":       question,
            "length_minutes": length_minutes,
            "created_at":     datetime.now().isoformat(),
        })
        print(f"  [media] Saved {filename}")
    except Exception as e:
        print(f"  [media] Failed to persist: {e}")


@app.route("/api/media")
def list_media():
    """Return all saved media items, newest first."""
    return jsonify(_load_media_index())


@app.route("/api/media/<media_id>/download")
def download_media(media_id: str):
    items = _load_media_index()
    item  = next((i for i in items if i.get("id") == media_id), None)
    if not item:
        return jsonify({"error": "Not found"}), 404
    fpath = os.path.join(MEDIA_DIR, item["filename"])
    if not os.path.exists(fpath):
        return jsonify({"error": "File missing"}), 404
    ext_map = {"mp3": "audio/mpeg", "mp4": "video/mp4",
               "html": "text/html",
               "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation"}
    mime = ext_map.get(item.get("ext", "mp3"), "application/octet-stream")
    return send_file(fpath, as_attachment=True,
                     download_name=item["filename"], mimetype=mime)


@app.route("/api/media/<media_id>/stream")
def stream_media(media_id: str):
    """Stream audio/video for in-browser playback (no attachment header)."""
    items = _load_media_index()
    item  = next((i for i in items if i.get("id") == media_id), None)
    if not item:
        return jsonify({"error": "Not found"}), 404
    fpath = os.path.join(MEDIA_DIR, item["filename"])
    if not os.path.exists(fpath):
        return jsonify({"error": "File missing"}), 404
    ext_map = {"mp3": "audio/mpeg", "mp4": "video/mp4",
               "html": "text/html",
               "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation"}
    mime = ext_map.get(item.get("ext", "mp3"), "application/octet-stream")
    return send_file(fpath, mimetype=mime)


@app.route("/api/media/<media_id>", methods=["DELETE"])
def delete_media(media_id: str):
    ok = _remove_media_item(media_id)
    return jsonify({"success": ok})


# ── Case Difficulty Assessment ────────────────────────────

_PROFILE_PATH = os.path.join(os.path.dirname(__file__), "provider_profile.json")

def _load_profile() -> dict:
    if os.path.exists(_PROFILE_PATH):
        try:
            with open(_PROFILE_PATH, "r", encoding="utf-8") as f:
                return _json_mod.load(f)
        except Exception:
            pass
    return {}

def _save_profile(data: dict):
    with open(_PROFILE_PATH, "w", encoding="utf-8") as f:
        _json_mod.dump(data, f, indent=2)


@app.route("/api/profile", methods=["GET"])
def get_profile():
    return jsonify(_load_profile())


@app.route("/api/profile", methods=["POST"])
def save_profile():
    data = request.json or {}
    _save_profile(data)
    return jsonify({"success": True})


@app.route("/api/settings", methods=["GET"])
def get_settings():
    """Return AI-behaviour settings (explanation level, content balance, theme)."""
    profile = _load_profile()
    return jsonify({
        "explanation_level": profile.get("explanation_level", "clinician"),
        "content_balance":   profile.get("content_balance",   "balanced"),
        "theme":             profile.get("theme",             "dark"),
    })


@app.route("/api/settings", methods=["POST"])
def save_settings():
    """Persist AI-behaviour settings into the profile JSON."""
    data    = request.json or {}
    profile = _load_profile()
    for key in ("explanation_level", "content_balance", "theme"):
        if key in data:
            profile[key] = data[key]
    _save_profile(profile)
    return jsonify({"success": True})


# ── X-ray gating (PHI) ───────────────────────────────────
# Patient radiographs are PHI. The vision path ships DISABLED and is enabled
# per-deployment via ENABLE_XRAY=true — a decision recorded in WORKLIST §5:
# enabling it in production requires a BAA with the vision provider
# (Gemini / OpenAI). See HANDOVER.md.

def _xray_enabled() -> bool:
    """Read ENABLE_XRAY at request time so tests/deploys can toggle it."""
    return (os.getenv("ENABLE_XRAY") or "").strip().lower() in ("1", "true", "yes", "on")


def _sanitize_tooth_hint(raw: str) -> str:
    """Reduce the tooth hint to a bare tooth designation (e.g. '14', '#30',
    'B', '4.6'). Anything longer or containing free text is dropped entirely,
    so patient case narrative can never ride along with the image to the
    vision provider."""
    hint = (raw or "").strip()
    if _audit_re.fullmatch(r"#?[A-Ta-t0-9][0-9.]{0,3}", hint):
        return hint.lstrip("#")
    return ""


def _strip_image_metadata(image_bytes: bytes, ext: str) -> bytes:
    """Re-encode the uploaded image with Pillow, dropping ALL ancillary
    metadata: EXIF/GPS/IPTC/XMP on JPEG, tEXt/iTXt/zTXt chunks on PNG.
    Radiograph exports routinely embed patient name/DOB in these fields.
    Raises on any decode failure — the caller fails CLOSED (rejects the
    upload) rather than forwarding un-stripped bytes."""
    from PIL import Image
    import io as _io
    img = Image.open(_io.BytesIO(image_bytes))
    out = _io.BytesIO()
    if ext in ("jpg", "jpeg"):
        img.convert("RGB").save(out, format="JPEG", quality=95)
    else:
        # Pillow drops text chunks unless an explicit pnginfo is passed.
        img.save(out, format="PNG")
    return out.getvalue()


@app.route("/api/analyze-xray", methods=["POST"])
def analyze_xray():
    """Upload a PA radiograph; a vision model pre-fills the assessment form.

    Disabled by default (ENABLE_XRAY unset/false -> 403): radiographs are PHI
    and sending them to a third-party vision API requires a BAA. When enabled,
    the image is re-encoded to strip EXIF/PNG metadata, and only a sanitized
    tooth number — never case text — accompanies it."""
    if not _xray_enabled():
        return jsonify({
            "error": "X-ray analysis is disabled on this deployment. "
                     "Patient radiographs are PHI; enabling the vision path "
                     "requires a BAA with the vision provider. Set "
                     "ENABLE_XRAY=true only once that is in place.",
            "feature": "xray", "enabled": False,
        }), 403
    if "image" not in request.files:
        return jsonify({"error": "No image provided"}), 400
    file = request.files["image"]
    ext  = (file.filename or "").rsplit(".", 1)[-1].lower()
    if ext not in ("png", "jpg", "jpeg"):
        return jsonify({"error": "Use PNG or JPG"}), 400
    media_type = "image/jpeg" if ext in ("jpg", "jpeg") else "image/png"
    image_bytes = file.read()
    if len(image_bytes) > 12 * 1024 * 1024:
        return jsonify({"error": "Image too large (max 12 MB)"}), 400
    # Strip embedded metadata (EXIF, GPS, PNG text chunks) before the bytes
    # leave this server. Fail closed: an image Pillow cannot decode is
    # rejected, never forwarded raw.
    try:
        image_bytes = _strip_image_metadata(image_bytes, ext)
    except Exception:
        return jsonify({"error": "Could not process image (metadata "
                                 "stripping failed) — upload not sent."}), 400
    tooth_hint = _sanitize_tooth_hint(request.form.get("tooth_hint"))
    provider   = (request.form.get("provider")   or "auto").strip()
    if provider not in ("gemini", "openai"):
        provider = "auto"
    try:
        raw      = analyze_radiograph(image_bytes, media_type, tooth_hint=tooth_hint, provider=provider)
        prefill  = _analysis_to_prefill(raw)
        meta     = raw.get("_meta", {"provider": "unknown", "fallback_reason": None})
        return jsonify({"success": True, "prefill": prefill, "raw": raw,
                        "provider":        meta.get("provider"),
                        "fallback_reason": meta.get("fallback_reason"),
                        "cost_usd":        meta.get("cost_usd", 0.0),
                        "input_tokens":    meta.get("input_tokens", 0),
                        "output_tokens":   meta.get("output_tokens", 0)})
    except Exception as e:
        return jsonify({"success": False, "error": str(e), "fallback": "manual"}), 200


@app.route("/api/assess", methods=["POST"])
def assess_case():
    """Score a case and compare against the provider's profile."""
    data         = request.json or {}
    answers      = data.get("answers", {})
    profile      = data.get("profile") or _load_profile()

    # Build AAE factor scores from questionnaire answers
    scores = _answers_to_aae_scores(answers)
    aae    = calculate_case_difficulty(scores)

    # Personalised match (only if profile exists)
    if profile:
        match = match_case_to_profile(answers, profile)
    else:
        match = {
            "recommendation":    aae["level"],
            "summary":           f"AAE difficulty: {aae['level']}",
            "exceeds_comfort":   aae["high_factors"],
            "within_comfort":    [],
            "equipment_warnings":[]
        }

    # Echo back clinical diagnoses for the result panel
    pulpal_dx     = answers.get("pulpal_diagnosis", "")
    periapical_dx = answers.get("periapical_diagnosis", "")

    # Override recommendation for unusual/conflicting diagnosis combinations
    conflict_msgs = answers.get("_diagnosisConflictMessages", [])
    if answers.get("_diagnosisConflict") and conflict_msgs:
        match["recommendation"] = "REFER"
        match["exceeds_comfort"] = conflict_msgs + match.get("exceeds_comfort", [])
        match["summary"] = (
            "Unusual diagnosis combination — endodontist evaluation required "
            "before proceeding with treatment"
        )

    return jsonify({
        "success": True,
        "aae": aae,
        "match": match,
        "pulpal_diagnosis": pulpal_dx,
        "periapical_diagnosis": periapical_dx,
        "diagnosis_conflict": bool(conflict_msgs),
    })


@app.route("/api/referral-letter", methods=["POST"])
def referral_letter():
    data    = request.json or {}
    profile = data.get("profile") or _load_profile()
    reasons = data.get("reasons", [])
    case    = data.get("case", {})
    try:
        letter = generate_referral_letter(case, profile, reasons)
        return jsonify({"success": True, "letter": letter})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


def _answers_to_aae_scores(a: dict) -> dict:
    """Convert questionnaire answers to 1-3 AAE factor scores."""

    def curv_score(v):
        return {"straight": 1, "mild": 1, "moderate": 2, "severe": 3}.get(v, 1)

    def calc_score(v):
        return {"open": 1, "mild": 1, "moderate": 2, "severe": 3}.get(v, 1)

    def root_score(v):
        return {"normal": 1, "extra_roots": 2, "dilacerated": 3, "short": 2, "long": 2}.get(v, 1)

    def peri_score(v):
        return {"normal": 1, "small_lesion": 2, "large_lesion": 3}.get(v, 1)

    def retx_score(v):
        return {"none": 1, "simple": 2, "carrier": 2, "posts": 3, "separated": 3}.get(v, 1)

    complications = a.get("complications", ["none"])
    patient_factors = a.get("patient_factors", ["none"])

    def patient_score():
        if any(f in patient_factors for f in ["medical_complex", "bisphosphonates", "anticoagulation"]):
            return 3
        if any(f in patient_factors for f in ["hot_tooth", "limited_opening", "gag_reflex", "anxiety"]):
            return 2
        return 1

    # Clinical findings — percussion/palpation/cold test inform diagnosis clarity
    cold_test  = a.get("cold_test", "not_tested")
    percussion = a.get("percussion", "not_tested")
    palpation  = a.get("palpation", "not_tested")

    # Diagnosis clarity: base on periapical status, boosted if clinical findings confirm
    diag_clarity = peri_score(a.get("periapical_status", "normal"))
    if cold_test == "no_response" and a.get("periapical_status", "normal") == "normal":
        diag_clarity = max(diag_clarity, 2)  # necrosis without radiographic lesion = less clear
    if cold_test == "not_tested" and a.get("periapical_status", "normal") != "normal":
        diag_clarity = max(diag_clarity, 2)  # lesion without pulp testing = incomplete workup

    # Symptomatic irreversible pulpitis with percussion → harder anaesthesia
    is_symptomatic = (cold_test == "exaggerated_lingering" or
                      percussion == "yes" or palpation == "yes")
    anesthesia_score = 3 if "hot_tooth" in patient_factors else (2 if is_symptomatic else 1)

    return {
        "medical_history":        patient_score(),
        "anesthesia_history":     anesthesia_score,
        "patient_disposition":    3 if "anxiety" in patient_factors else 1,
        "mouth_opening":          3 if "limited_opening" in patient_factors else 1,
        "gag_reflex":             3 if "gag_reflex" in patient_factors else 1,
        "diagnosis_clarity":      diag_clarity,
        "radiographic_difficulty":2 if a.get("calcification") in ["moderate", "severe"] else 1,
        "tooth_position":         _tooth_position_score(a.get("tooth_number", 0)),
        "isolation_difficulty":   {"intact": 1, "filling": 1, "full_coverage": 2, "post_crown": 3}.get(a.get("crown_type", "intact"), 1),
        "crown_morphology":       {"easy": 1, "moderate": 2, "difficult": 3}.get(a.get("crown_access", "easy"), 1),
        "canal_morphology":       max(curv_score(a.get("canal_curvature", "mild")), calc_score(a.get("calcification", "open"))),
        "root_morphology":        root_score(a.get("root_anatomy", "normal")),
        "resorption":             3 if "resorption" in complications else 1,
        "trauma_history":         {"uncomplicated_fracture": 2, "complicated_fracture": 2, "luxation": 3, "avulsion": 3}.get(a.get("trauma_type", ""), 1) if "trauma" in complications else 1,
        "previous_endo":          retx_score(a.get("retreatment_complexity", "none") if a.get("is_retreatment") == "yes" else "none"),
        "perio_endo":             2 if "perio" in complications else 1,
    }


def _tooth_position_score(tooth_num: int) -> int:
    if tooth_num in [1, 16, 17, 32]:    return 3  # 3rd molars
    if tooth_num in [2, 15, 18, 31]:    return 3  # 2nd molars
    if tooth_num in [3, 14, 19, 30]:    return 2  # 1st molars
    if tooth_num in [4, 5, 12, 13, 20, 21, 28, 29]: return 1  # premolars
    return 1  # anteriors


# ── Main ─────────────────────────────────────────────────

if __name__ == "__main__":
    app.run(debug=True, host="127.0.0.1", port=5000)
